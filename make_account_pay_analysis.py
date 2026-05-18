# -*- coding: utf-8 -*-
from __future__ import annotations

import json
import math
import subprocess
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "account_pay_list_payment_trend_revdate_2025_2026_ytd_fixed.pptx"
PREVIEW_PREFIX = ROOT / "account_pay_list_payment_trend_revdate_fixed_slide_"

KFONT = "Malgun Gothic"
NAVY = RGBColor(19, 17, 118)
INK = RGBColor(29, 37, 52)
MUTED = RGBColor(96, 107, 121)
BLUE = RGBColor(31, 111, 213)
TEAL = RGBColor(21, 154, 168)
GREEN = RGBColor(46, 155, 87)
RED = RGBColor(216, 68, 68)
ORANGE = RGBColor(231, 134, 32)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(217, 225, 234)


def run_query() -> dict:
    php = r"""
<?php
include 'include/dbconn.php';
@mysql_query("SET NAMES utf8", $dbConn);

$start = '2025-01-01';
$end = '2026-05-01';
$baseWhere = "a.p_code = b.p_code
  AND a.parent = 'MAIN'
  AND a.rev_status in ('DONE','READY','PPAY','OPAY')
  AND a.revDate >= '$start' AND a.revDate < '$end'
  AND MONTH(a.revDate) BETWEEN 1 AND 4";

$out = array(
  'monthly' => array(),
  'pay_status' => array(),
  'rev_status' => array(),
  'ptype' => array(),
  'new_products' => array(),
  'top_balance' => array()
);

$sql = "SELECT YEAR(a.revDate) AS yy, MONTH(a.revDate) AS mm,
        COUNT(DISTINCT a.reserveCode) AS bookings,
        SUM(a.p_cnt) AS pax,
        SUM(a.last_total) AS total,
        SUM(a.last_total - a.last_bal) AS paid,
        SUM(a.last_bal) AS balance
        FROM reserve_info a, product_master b
        WHERE $baseWhere
        GROUP BY YEAR(a.revDate), MONTH(a.revDate)
        ORDER BY yy, mm";
$res = mysql_query($sql, $dbConn);
while ($row = mysql_fetch_assoc($res)) { $out['monthly'][] = $row; }

$sql = "SELECT YEAR(a.revDate) AS yy,
        IFNULL(NULLIF(a.payment_st,''),'UNKNOWN') AS payment_st,
        COUNT(DISTINCT a.reserveCode) AS bookings,
        SUM(a.last_total) AS total,
        SUM(a.last_total - a.last_bal) AS paid,
        SUM(a.last_bal) AS balance
        FROM reserve_info a, product_master b
        WHERE $baseWhere
        GROUP BY YEAR(a.revDate), IFNULL(NULLIF(a.payment_st,''),'UNKNOWN')
        ORDER BY yy, total DESC";
$res = mysql_query($sql, $dbConn);
while ($row = mysql_fetch_assoc($res)) { $out['pay_status'][] = $row; }

$sql = "SELECT YEAR(a.revDate) AS yy,
        IFNULL(NULLIF(a.rev_status,''),'UNKNOWN') AS rev_status,
        COUNT(DISTINCT a.reserveCode) AS bookings,
        SUM(a.last_total) AS total,
        SUM(a.last_bal) AS balance
        FROM reserve_info a, product_master b
        WHERE $baseWhere
        GROUP BY YEAR(a.revDate), IFNULL(NULLIF(a.rev_status,''),'UNKNOWN')
        ORDER BY yy, total DESC";
$res = mysql_query($sql, $dbConn);
while ($row = mysql_fetch_assoc($res)) { $out['rev_status'][] = $row; }

$sql = "SELECT YEAR(a.revDate) AS yy,
        IFNULL(NULLIF(b.p_type,''),'UNKNOWN') AS p_type,
        COUNT(DISTINCT a.reserveCode) AS bookings,
        SUM(a.p_cnt) AS pax,
        SUM(a.last_total) AS total,
        SUM(a.last_bal) AS balance
        FROM reserve_info a, product_master b
        WHERE $baseWhere
        GROUP BY YEAR(a.revDate), IFNULL(NULLIF(b.p_type,''),'UNKNOWN')
        ORDER BY yy, total DESC";
$res = mysql_query($sql, $dbConn);
while ($row = mysql_fetch_assoc($res)) { $out['ptype'][] = $row; }

$sql = "SELECT b.p_code, b.p_name, b.p_type, b.p_day, b.wdate,
        COUNT(DISTINCT a.reserveCode) AS bookings,
        SUM(a.p_cnt) AS pax,
        SUM(a.last_total) AS total,
        SUM(a.last_total - a.last_bal) AS paid,
        SUM(a.last_bal) AS balance
        FROM reserve_info a, product_master b
        WHERE $baseWhere
          AND YEAR(a.revDate) = 2026
          AND b.wdate >= '2026-01-01' AND b.wdate < '2026-05-01'
        GROUP BY b.p_code, b.p_name, b.p_type, b.p_day, b.wdate
        ORDER BY total DESC";
$res = mysql_query($sql, $dbConn);
while ($row = mysql_fetch_assoc($res)) {
  $row['p_name'] = strip_tags($row['p_name']);
  $out['new_products'][] = $row;
}

$sql = "SELECT a.reserveCode, a.p_code, b.p_name, a.revDate,
        a.payment_st, a.rev_status, a.p_cnt, a.last_total,
        (a.last_total - a.last_bal) AS paid, a.last_bal AS balance
        FROM reserve_info a, product_master b
        WHERE $baseWhere
          AND YEAR(a.revDate) = 2026
          AND a.last_bal > 0
        ORDER BY a.last_bal DESC
        LIMIT 10";
$res = mysql_query($sql, $dbConn);
while ($row = mysql_fetch_assoc($res)) {
  $row['p_name'] = strip_tags($row['p_name']);
  $out['top_balance'][] = $row;
}

echo json_encode($out, JSON_UNESCAPED_UNICODE);
?>
"""
    proc = subprocess.run(
        ["php"],
        input=php,
        cwd=str(ROOT),
        text=True,
        encoding="utf-8",
        capture_output=True,
        check=True,
    )
    stdout = proc.stdout.strip()
    start = stdout.find("{")
    return json.loads(stdout[start:] if start >= 0 else stdout)


def val(x) -> float:
    try:
        return float(x or 0)
    except (TypeError, ValueError):
        return 0.0


def money(v: float, digits: int = 1) -> str:
    return f"${v / 1_000_000:.{digits}f}M"


def money_k(v: float) -> str:
    return f"${v / 1_000:.0f}K"


def num(v: float) -> str:
    return f"{v:,.0f}"


def pct(v: float | None) -> str:
    if v is None or math.isnan(v):
        return "n/a"
    return f"{'+' if v > 0 else ''}{v:.1f}%"


def trim(s: str, n: int = 30) -> str:
    s = " ".join(str(s or "").split())
    return s if len(s) <= n else s[: n - 1] + "…"


PAY = {"READY": "미납", "PPAY": "부분완납", "DONE": "완납", "OPAY": "환불", "UNKNOWN": "미지정"}
REV = {"READY": "예약접수", "DONE": "예약확정", "CANCEL": "예약취소", "UNKNOWN": "미지정"}
PTYPE = {"1": "로컬상품", "2": "인바운드", "4": "인센티브", "5": "아웃바운드", "6": "기타", "UNKNOWN": "미지정"}


def summarize(data: dict) -> dict:
    months = [1, 2, 3, 4]
    by_month = {(int(r["yy"]), int(r["mm"])): r for r in data["monthly"]}

    def series(year: int, key: str) -> list[float]:
        return [val(by_month.get((year, m), {}).get(key)) for m in months]

    out = {
        "month_names": ["1월", "2월", "3월", "4월"],
        "total": {2025: series(2025, "total"), 2026: series(2026, "total")},
        "paid": {2025: series(2025, "paid"), 2026: series(2026, "paid")},
        "balance": {2025: series(2025, "balance"), 2026: series(2026, "balance")},
        "pax": {2025: series(2025, "pax"), 2026: series(2026, "pax")},
        "bookings": {2025: series(2025, "bookings"), 2026: series(2026, "bookings")},
    }
    for year in (2025, 2026):
        for key in ("total", "paid", "balance", "pax", "bookings"):
            out[f"{key}{year}"] = sum(out[key][year])
        out[f"collection{year}"] = out[f"paid{year}"] / out[f"total{year}"] * 100 if out[f"total{year}"] else 0
        out[f"balance_rate{year}"] = out[f"balance{year}"] / out[f"total{year}"] * 100 if out[f"total{year}"] else 0

    for key in ("total", "paid", "balance", "pax", "bookings"):
        out[f"{key}_delta"] = out[f"{key}2026"] - out[f"{key}2025"]
        out[f"{key}_pct"] = out[f"{key}_delta"] / out[f"{key}2025"] * 100 if out[f"{key}2025"] else None
    out["collection_delta"] = out["collection2026"] - out["collection2025"]

    out["pay"] = []
    for r in data["pay_status"]:
        out["pay"].append({"label": PAY.get(str(r["payment_st"]), str(r["payment_st"])), "year": int(r["yy"]), "total": val(r["total"]), "balance": val(r["balance"]), "bookings": val(r["bookings"])})
    out["pay2026"] = sorted([r for r in out["pay"] if r["year"] == 2026], key=lambda x: x["total"], reverse=True)

    ptype = {}
    for r in data["ptype"]:
        k = str(r["p_type"])
        item = ptype.setdefault(k, {"label": PTYPE.get(k, k), "t2025": 0.0, "t2026": 0.0})
        item[f"t{int(r['yy'])}"] += val(r["total"])
    out["ptype"] = sorted([{**v, "delta": v["t2026"] - v["t2025"]} for v in ptype.values()], key=lambda x: x["delta"])

    out["rev"] = []
    for r in data["rev_status"]:
        out["rev"].append({"label": REV.get(str(r["rev_status"]), str(r["rev_status"])), "year": int(r["yy"]), "total": val(r["total"]), "bookings": val(r["bookings"])})
    out["rev2026"] = sorted([r for r in out["rev"] if r["year"] == 2026], key=lambda x: x["total"], reverse=True)

    out["new_products"] = [
        {
            "p_code": r["p_code"],
            "p_name": trim(r["p_name"]),
            "wdate": str(r["wdate"])[:10],
            "bookings": val(r["bookings"]),
            "pax": val(r["pax"]),
            "total": val(r["total"]),
            "paid": val(r["paid"]),
            "balance": val(r["balance"]),
        }
        for r in data["new_products"]
    ]
    out["new_total"] = sum(r["total"] for r in out["new_products"])
    out["new_paid"] = sum(r["paid"] for r in out["new_products"])
    out["new_balance"] = sum(r["balance"] for r in out["new_products"])
    out["new_bookings"] = sum(r["bookings"] for r in out["new_products"])
    out["new_share"] = out["new_total"] / out["total2026"] * 100 if out["total2026"] else 0
    out["existing_total2026"] = out["total2026"] - out["new_total"]
    out["existing_total_pct"] = (out["existing_total2026"] - out["total2025"]) / out["total2025"] * 100 if out["total2025"] else None

    out["top_balance"] = [
        {"reserveCode": r["reserveCode"], "revDate": r["revDate"], "p_name": trim(r["p_name"]), "balance": val(r["balance"])}
        for r in data["top_balance"]
    ]
    return out


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, body, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = str(body)
    p.alignment = align
    for run in p.runs:
        run.font.name = KFONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def title(slide, heading: str, kicker: str):
    add_text(slide, Inches(0.65), Inches(0.32), Inches(11.5), Inches(0.28), kicker, 10, True, TEAL)
    add_text(slide, Inches(0.62), Inches(0.65), Inches(11.8), Inches(0.55), heading, 28, True, NAVY)


def footer(slide, source: str):
    add_text(slide, Inches(0.65), Inches(7.06), Inches(11.8), Inches(0.25), source, 8.5, False, MUTED)


def metric(slide, x, y, label, value, note="", color=NAVY):
    add_text(slide, x, y, Inches(2.7), Inches(0.24), label, 10.5, True, MUTED)
    add_text(slide, x, y + Inches(0.25), Inches(2.8), Inches(0.5), value, 24, True, color)
    if note:
        add_text(slide, x, y + Inches(0.78), Inches(2.8), Inches(0.3), note, 10.5, False, MUTED)


def bullets(slide, x, y, lines, size=14, width=5.6):
    body = "\n".join(f"- {line}" for line in lines)
    add_text(slide, x, y, Inches(width), Inches(2.6), body, size, False, INK)


def chart(slide, x, y, w, h, categories, data_series, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
    data = CategoryChartData()
    data.categories = categories
    for name, values in data_series:
        data.add_series(name, values)
    ch = slide.shapes.add_chart(chart_type, x, y, w, h, data).chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.value_axis.tick_labels.font.size = Pt(9)
    ch.category_axis.tick_labels.font.size = Pt(10)
    return ch


def table(slide, x, y, rows, widths, total_width=11.8, font_size=9):
    row_h = Inches(0.34)
    total_w = Inches(total_width)
    for r, row in enumerate(rows):
        cur = x
        for c, cell in enumerate(row):
            w = total_w * widths[c]
            shape = slide.shapes.add_shape(1, cur, y + row_h * r, w, row_h)
            shape.fill.solid()
            shape.fill.fore_color.rgb = NAVY if r == 0 else WHITE
            shape.line.color.rgb = LINE
            add_text(slide, cur + Inches(0.04), y + row_h * r + Inches(0.05), w - Inches(0.08), row_h, cell, font_size, r == 0, WHITE if r == 0 else INK, PP_ALIGN.RIGHT if c > 0 else PP_ALIGN.LEFT)
            cur += w


def build_deck(s: dict):
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    shape = slide.shapes.add_shape(1, 0, 0, Inches(4.1), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    add_text(slide, Inches(0.62), Inches(0.72), Inches(3.0), Inches(0.3), "ACCOUNT PAY LIST", 11, True, WHITE)
    add_text(slide, Inches(0.62), Inches(1.42), Inches(3.2), Inches(2.1), "예약별\n결제현황\n판매일자 기준", 29, True, WHITE)
    add_text(slide, Inches(0.62), Inches(5.9), Inches(3.2), Inches(0.5), "2025년 1-4월 vs 2026년 1-4월\nrevDate 기준 / 전체 예약", 11, False, WHITE)
    add_text(slide, Inches(4.8), Inches(1.0), Inches(7.0), Inches(0.45), "총 결제금액", 15, True, MUTED)
    add_text(slide, Inches(4.75), Inches(1.45), Inches(7.3), Inches(0.8), f"{money(s['total2025'])} → {money(s['total2026'])}", 40, True, NAVY)
    add_text(slide, Inches(4.85), Inches(2.35), Inches(6.5), Inches(0.35), f"전년 대비 {money(s['total_delta'])} ({pct(s['total_pct'])})", 16, True, RED if s["total_delta"] < 0 else GREEN)
    metric(slide, Inches(4.9), Inches(3.35), "받은금액", money(s["paid2026"]), f"YoY {pct(s['paid_pct'])}", BLUE)
    metric(slide, Inches(7.45), Inches(3.35), "잔액", money(s["balance2026"]), f"잔액률 {s['balance_rate2026']:.1f}%", ORANGE)
    metric(slide, Inches(9.95), Inches(3.35), "신규상품", money(s["new_total"]), f"기여율 {s['new_share']:.1f}%", TEAL)
    footer(slide, "출처: 인트라넷 메뉴 '예약별 결제현황' 판매일자 기준 DB 집계")

    slide = prs.slides.add_slide(BLANK)
    title(slide, "분석 기준과 핵심 변동", "기준")
    metric(slide, Inches(0.75), Inches(1.45), "총 결제금액", money(s["total_delta"]), pct(s["total_pct"]), RED if s["total_delta"] < 0 else GREEN)
    metric(slide, Inches(3.55), Inches(1.45), "받은금액", money(s["paid_delta"]), pct(s["paid_pct"]), RED if s["paid_delta"] < 0 else GREEN)
    metric(slide, Inches(6.35), Inches(1.45), "잔액", money(s["balance_delta"]), pct(s["balance_pct"]), ORANGE)
    metric(slide, Inches(9.15), Inches(1.45), "예약건수", num(s["bookings2026"]), pct(s["bookings_pct"]), TEAL)
    rows = [["항목", "2025 1-4월", "2026 1-4월", "변동"], ["예약건수", num(s["bookings2025"]), num(s["bookings2026"]), pct(s["bookings_pct"])], ["인원", num(s["pax2025"]), num(s["pax2026"]), pct(s["pax_pct"])], ["총 결제금액", money(s["total2025"]), money(s["total2026"]), pct(s["total_pct"])], ["받은금액", money(s["paid2025"]), money(s["paid2026"]), pct(s["paid_pct"])], ["잔액", money(s["balance2025"]), money(s["balance2026"]), pct(s["balance_pct"])], ["수금률", f"{s['collection2025']:.1f}%", f"{s['collection2026']:.1f}%", f"{s['collection_delta']:+.1f}p"]]
    table(slide, Inches(0.75), Inches(3.15), rows, [0.34, 0.22, 0.22, 0.22], total_width=6.9)
    bullets(slide, Inches(8.1), Inches(3.25), ["날짜 기준: reserve_info.revDate", "금액: last_total / last_total-last_bal / last_bal", "조건: parent='MAIN', 활성 예약상태", "신규상품: product_master.wdate가 2026년 1-4월"], width=4.4)
    footer(slide, "출처: 인트라넷 메뉴 '예약별 결제현황' 필터 및 합계 기준")

    slide = prs.slides.add_slide(BLANK)
    title(slide, "월별 결제금액 흐름", "매출 추이")
    chart(slide, Inches(0.75), Inches(1.45), Inches(7.5), Inches(4.9), s["month_names"], [("2025 총액", [v / 1_000_000 for v in s["total"][2025]]), ("2026 총액", [v / 1_000_000 for v in s["total"][2026]])])
    chart(slide, Inches(8.55), Inches(1.6), Inches(4.0), Inches(2.1), s["month_names"], [("2025 받은금액", [v / 1_000_000 for v in s["paid"][2025]]), ("2026 받은금액", [v / 1_000_000 for v in s["paid"][2026]])])
    rate25 = [(p / t * 100 if t else 0) for p, t in zip(s["paid"][2025], s["total"][2025])]
    rate26 = [(p / t * 100 if t else 0) for p, t in zip(s["paid"][2026], s["total"][2026])]
    chart(slide, Inches(8.55), Inches(4.2), Inches(4.0), Inches(2.05), s["month_names"], [("2025 수금률", rate25), ("2026 수금률", rate26)], XL_CHART_TYPE.LINE_MARKERS)
    footer(slide, "출처: 인트라넷 메뉴 '예약별 결제현황' 월별 DB 집계")

    slide = prs.slides.add_slide(BLANK)
    title(slide, "미수/수금 변동현황", "잔액")
    chart(slide, Inches(0.75), Inches(1.45), Inches(6.2), Inches(4.9), s["month_names"], [("2025 잔액", [v / 1_000_000 for v in s["balance"][2025]]), ("2026 잔액", [v / 1_000_000 for v in s["balance"][2026]])])
    rows = [["예약번호", "판매일자", "상품", "잔액"]]
    for r in s["top_balance"][:8]:
        rows.append([r["reserveCode"], r["revDate"], r["p_name"], money_k(r["balance"])])
    table(slide, Inches(7.05), Inches(1.55), rows, [0.25, 0.18, 0.39, 0.18], total_width=5.55, font_size=8)
    footer(slide, "출처: 인트라넷 메뉴 '예약별 결제현황' 미수/잔액 집계")

    slide = prs.slides.add_slide(BLANK)
    title(slide, "올해 신규생성상품은 별도 반영", "신규상품")
    metric(slide, Inches(0.75), Inches(1.42), "신규상품 총액", money(s["new_total"]), f"2026 총액의 {s['new_share']:.1f}%", TEAL)
    metric(slide, Inches(3.55), Inches(1.42), "신규상품 받은금액", money(s["new_paid"]), f"잔액 {money(s['new_balance'])}", BLUE)
    metric(slide, Inches(6.35), Inches(1.42), "신규상품 예약", num(s["new_bookings"]), "", NAVY)
    metric(slide, Inches(9.15), Inches(1.42), "기존상품 조정 YoY", pct(s["existing_total_pct"]), f"기존상품 {money(s['existing_total2026'])}", RED if (s["existing_total_pct"] or 0) < 0 else GREEN)
    rows = [["상품코드", "생성일", "상품명", "예약", "총액", "잔액"]]
    for r in s["new_products"][:8]:
        rows.append([r["p_code"], r["wdate"], r["p_name"], num(r["bookings"]), money(r["total"]), money_k(r["balance"])])
    if len(rows) == 1:
        rows.append(["-", "-", "2026년 1-4월 생성 후 결제금액 있는 상품 없음", "-", "$0.0M", "$0K"])
    table(slide, Inches(0.75), Inches(3.08), rows, [0.17, 0.14, 0.33, 0.09, 0.14, 0.13])
    footer(slide, "출처: 인트라넷 메뉴 '예약별 결제현황' 신규상품 별도 집계")

    slide = prs.slides.add_slide(BLANK)
    title(slide, "내부 믹스: 결제상태와 상품유형", "상태/상품")
    rows = [["결제상태", "2026 총액", "잔액", "예약"]]
    for r in s["pay2026"][:6]:
        rows.append([r["label"], money(r["total"]), money(r["balance"]), num(r["bookings"])])
    table(slide, Inches(0.75), Inches(1.55), rows, [0.28, 0.25, 0.25, 0.22])
    rows = [["상품유형", "2025", "2026", "변동"]]
    for r in s["ptype"][:6]:
        rows.append([r["label"], money(r["t2025"]), money(r["t2026"]), money(r["delta"])])
    table(slide, Inches(0.75), Inches(4.35), rows, [0.28, 0.24, 0.24, 0.24])
    footer(slide, "출처: 인트라넷 메뉴 '예약별 결제현황' 결제상태/상품유형 집계")

    slide = prs.slides.add_slide(BLANK)
    title(slide, "외부요인과 내부요인 분리", "원인 해석")
    add_text(slide, Inches(0.85), Inches(1.55), Inches(5.5), Inches(0.32), "외부요인", 17, True, TEAL)
    bullets(slide, Inches(0.85), Inches(2.0), ["미 교통 CPI는 2026년 3월 전년 대비 +5.0%로 공표", "NTTO는 2026년 미국 국제방문객 증가 전망", "환율 변동은 USD 결제 상품의 고객 체감가와 결제 시점에 영향"])
    add_text(slide, Inches(6.9), Inches(1.55), Inches(5.5), Inches(0.32), "내부요인", 17, True, ORANGE)
    bullets(slide, Inches(6.9), Inches(2.0), [f"판매일자 기준 총 결제금액 YoY {pct(s['total_pct'])}", f"신규상품 제외 기존상품 YoY {pct(s['existing_total_pct'])}", f"수금률은 {s['collection2025']:.1f}%에서 {s['collection2026']:.1f}%로 변화"])
    footer(slide, "출처: BTS Transportation CPI Mar 2026, NTTO Forecasts, Federal Reserve H.10, 인트라넷 '예약별 결제현황' 집계")

    slide = prs.slides.add_slide(BLANK)
    title(slide, "실행 제안", "액션")
    rows = [["우선순위", "실행", "확인 지표"], ["1", "판매일자 기준 미납/부분완납 예약 추적", "수금률, 잔액률"], ["2", "잔액 상위 예약을 담당자별로 매일 확인", "Top 잔액 감소액"], ["3", "올해 신규생성상품은 기존상품 YoY와 분리 보고", "신규상품 총액/기여율"], ["4", "결제상태와 예약상태 불일치 점검", "미지정/오류 상태 건수"], ["5", "상품유형별 하락 구간에 프로모션 또는 일정 재배치", "상품유형별 YoY 변동"]]
    table(slide, Inches(0.75), Inches(1.55), rows, [0.14, 0.58, 0.28])
    footer(slide, "출처: 인트라넷 메뉴 '예약별 결제현황'")

    prs.save(OUT)


def render_previews(s: dict) -> list[str]:
    try:
        font_big = ImageFont.truetype("C:/Windows/Fonts/malgunbd.ttf", 48)
        font_mid = ImageFont.truetype("C:/Windows/Fonts/malgun.ttf", 28)
        font_small = ImageFont.truetype("C:/Windows/Fonts/malgun.ttf", 20)
    except Exception:
        font_big = font_mid = font_small = ImageFont.load_default()
    slides = [
        ("예약별 결제현황 판매일자 기준", [f"총액 {money(s['total2025'])} → {money(s['total2026'])}", f"YoY {pct(s['total_pct'])}"]),
        ("분석 기준과 핵심 변동", [f"받은금액 {money(s['paid2026'])}", "revDate 기준"]),
        ("월별 결제금액 흐름", [f"총액 변동 {money(s['total_delta'])}", f"예약건수 YoY {pct(s['bookings_pct'])}"]),
        ("미수/수금 변동현황", [f"잔액률 {s['balance_rate2026']:.1f}%", f"잔액 YoY {pct(s['balance_pct'])}"]),
        ("올해 신규생성상품 별도 반영", [f"신규상품 총액 {money(s['new_total'])}", f"기여율 {s['new_share']:.1f}%"]),
        ("결제상태와 상품유형", [f"결제상태 {len(s['pay2026'])}개", f"상품유형 {len(s['ptype'])}개"]),
        ("외부요인과 내부요인", ["교통물가/환율/수요 전망", "결제 전환/신규상품/잔액 관리"]),
        ("실행 제안", ["판매일자 기준 수금 추적", "신규상품 분리 보고"]),
    ]
    paths = []
    for i, (head, lines) in enumerate(slides, 1):
        img = Image.new("RGB", (1600, 900), (244, 247, 251))
        d = ImageDraw.Draw(img)
        d.rectangle([0, 0, 1600, 82], fill=(19, 17, 118))
        d.text((70, 25), f"{i:02d}", font=font_small, fill=(255, 255, 255))
        d.text((150, 220), head, font=font_big, fill=(19, 17, 118))
        y = 330
        for line in lines:
            d.text((155, y), line, font=font_mid, fill=(29, 37, 52))
            y += 62
        d.text((150, 820), "인트라넷 메뉴 '예약별 결제현황' / 판매일자(revDate) 기준", font=font_small, fill=(96, 107, 121))
        path = PREVIEW_PREFIX.with_name(f"{PREVIEW_PREFIX.name}{i:02d}.png")
        img.save(path)
        paths.append(str(path))
    return paths


def main() -> None:
    summary = summarize(run_query())
    build_deck(summary)
    previews = render_previews(summary)
    print(json.dumps({
        "pptx": str(OUT),
        "previews": previews,
        "slides": len(prs.slides),
        "summary": {
            "total2025": summary["total2025"],
            "total2026": summary["total2026"],
            "total_pct": summary["total_pct"],
            "paid2025": summary["paid2025"],
            "paid2026": summary["paid2026"],
            "balance2025": summary["balance2025"],
            "balance2026": summary["balance2026"],
            "new_total": summary["new_total"],
            "new_share": summary["new_share"],
            "existing_total_pct": summary["existing_total_pct"],
        }
    }, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
