# -*- coding: utf-8 -*-
"""
admin/hotel_assign_m.php 사용자 설명 PPTX 생성.
스타일은 푸른투어 인트라넷 매뉴얼 계열과 맞춘다.
"""
from __future__ import annotations

import os
import time
from dataclasses import dataclass
from typing import Sequence, Tuple

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


NAVY = RGBColor(0x13, 0x11, 0x76)
NAVY_DARK = RGBColor(0x0A, 0x09, 0x55)
BLUE = RGBColor(0x0B, 0x5B, 0xD3)
GREEN = RGBColor(0x28, 0xA7, 0x45)
TEAL = RGBColor(0x20, 0xC9, 0x97)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)
RED = RGBColor(0xDC, 0x35, 0x45)
GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)
GRAY_TXT = RGBColor(0x6C, 0x75, 0x7D)
DARK_TXT = RGBColor(0x21, 0x25, 0x29)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE7, 0xF3, 0xFF)
LIGHT_GREEN = RGBColor(0xE8, 0xF7, 0xEE)
LIGHT_YELLOW = RGBColor(0xFF, 0xF8, 0xE1)

KFONT = "Malgun Gothic"
OUT = r"d:/www/prttour_myprt/hotel_assign_user_manual.pptx"
PREVIEW_PREFIX = r"d:/www/prttour_myprt/hotel_assign_user_manual_slide_"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def set_text(shape, text, size=14, bold=False, color=DARK_TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=KFONT):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = str(line)
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def add_text(slide, x, y, w, h, text, **kw):
    sh = slide.shapes.add_textbox(x, y, w, h)
    set_text(sh, text, **kw)
    return sh


def add_header(slide, title, subtitle, page, total):
    add_rect(slide, 0, 0, SW, Inches(0.9), NAVY)
    add_text(slide, Inches(0.4), Inches(0.10), Inches(10.4), Inches(0.46),
             title, size=22, bold=True, color=WHITE)
    add_text(slide, Inches(0.4), Inches(0.52), Inches(10.4), Inches(0.30),
             subtitle, size=11, color=RGBColor(0xCB, 0xD3, 0xE8))
    add_text(slide, Inches(11.5), Inches(0.28), Inches(1.55), Inches(0.34),
             f"{page} / {total}", size=11, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT)


def add_footer(slide, source="D:\\www\\hellousa\\admin\\hotel_assign_m.php"):
    add_rect(slide, 0, Inches(7.15), SW, Inches(0.35), GRAY_BG)
    add_text(slide, Inches(0.4), Inches(7.18), Inches(8), Inches(0.3),
             "푸른투어 인트라넷 · 호텔배정관리 사용자 매뉴얼",
             size=9, color=GRAY_TXT)
    add_text(slide, Inches(7.4), Inches(7.18), Inches(5.5), Inches(0.3),
             source, size=9, color=GRAY_TXT, align=PP_ALIGN.RIGHT)


def bullet(slide, x, y, w, items, size=12.2, bullet_color=NAVY, gap=0.42):
    cur_y = y
    for item in items:
        ic = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cur_y + Emu(45000),
                                    Inches(0.22), Inches(0.22))
        ic.fill.solid()
        ic.fill.fore_color.rgb = bullet_color
        ic.line.fill.background()
        set_text(ic, "✓", size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.32), cur_y, w - Inches(0.32), Inches(gap),
                 item, size=size, color=DARK_TXT)
        cur_y += Inches(gap)


def step_card(slide, x, y, w, h, no, title, body, color=NAVY):
    add_rect(slide, x, y, w, h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(slide, x, y, Inches(0.82), h, color)
    add_text(slide, x, y, Inches(0.82), h, no, size=18, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(1.0), y + Inches(0.18), w - Inches(1.2), Inches(0.38),
             title, size=13, bold=True, color=DARK_TXT)
    add_text(slide, x + Inches(1.0), y + Inches(0.66), w - Inches(1.2), h - Inches(0.82),
             body, size=10.5, color=GRAY_TXT, anchor=MSO_ANCHOR.TOP)


def table_like(slide, x, y, w, h, headers, rows, widths=None):
    add_rect(slide, x, y, w, h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
    if widths is None:
        widths = [1 / len(headers)] * len(headers)
    header_h = Inches(0.48)
    add_rect(slide, x, y, w, header_h, LIGHT_BLUE, line=RGBColor(0xDD, 0xDD, 0xDD))
    cur_x = x
    for hdr, ratio in zip(headers, widths):
        cw = int(w * ratio)
        add_text(slide, cur_x + Inches(0.06), y + Inches(0.05), cw - Inches(0.12), header_h - Inches(0.10),
                 hdr, size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        cur_x += cw
    row_h = int((h - header_h) / max(1, len(rows)))
    for r_idx, row in enumerate(rows):
        ry = y + header_h + row_h * r_idx
        if r_idx % 2 == 1:
            add_rect(slide, x, ry, w, row_h, RGBColor(0xFB, 0xFC, 0xFE))
        cur_x = x
        for txt, ratio in zip(row, widths):
            cw = int(w * ratio)
            add_text(slide, cur_x + Inches(0.06), ry + Inches(0.03), cw - Inches(0.12), row_h - Inches(0.06),
                     txt, size=8.9, color=DARK_TXT, align=PP_ALIGN.CENTER)
            cur_x += cw


def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, GRAY_BG)
    add_rect(s, 0, 0, SW, Inches(4.5), NAVY)
    for x, y, w in [(Inches(10.2), Inches(-1.5), Inches(5)), (Inches(-1.5), Inches(2.8), Inches(3.5))]:
        deco = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, w)
        deco.fill.solid()
        deco.fill.fore_color.rgb = NAVY_DARK
        deco.line.fill.background()
    add_text(s, Inches(0.8), Inches(0.7), Inches(7), Inches(0.4),
             "푸른투어 인트라넷 · 사용자 매뉴얼", size=14, bold=True,
             color=RGBColor(0xCB, 0xD3, 0xE8))
    add_text(s, Inches(0.8), Inches(1.35), Inches(11.5), Inches(1.15),
             "호텔배정관리", size=56, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(2.48), Inches(11), Inches(0.55),
             "Hotel Assignment Workflow Guide",
             size=21, color=RGBColor(0xCB, 0xD3, 0xE8))
    add_text(s, Inches(0.8), Inches(3.22), Inches(11.2), Inches(0.55),
             "행사별 차량/서브행사 기준으로 호텔, 방갯수, 호텔행사메모를 배정하는 화면",
             size=16, color=WHITE)
    cards = [
        ("대상 화면", "admin/hotel_assign_m.php", "호텔배정관리", NAVY),
        ("연동 파일", "admin/get_hotel.php", "호텔지역 선택 후 호텔명 조회", GREEN),
        ("업데이트", "admin/update_cnt.php", "방갯수 즉시 업데이트", ORANGE),
    ]
    for i, (eyebrow, title, desc, color) in enumerate(cards):
        x = Inches(0.8 + i * 4.05)
        add_rect(s, x, Inches(5.0), Inches(3.8), Inches(1.6), WHITE, line=color)
        add_text(s, x + Inches(0.18), Inches(5.1), Inches(3.45), Inches(0.35),
                 eyebrow, size=12, bold=True, color=color)
        add_text(s, x + Inches(0.18), Inches(5.45), Inches(3.45), Inches(0.45),
                 title, size=15, bold=True, color=DARK_TXT)
        add_text(s, x + Inches(0.18), Inches(5.95), Inches(3.45), Inches(0.4),
                 desc, size=11, color=GRAY_TXT)
    add_text(s, Inches(0.8), Inches(6.85), Inches(11), Inches(0.35),
             "푸른투어 인트라넷", size=11, bold=True, color=NAVY)


def slide_toc(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "목차 · CONTENTS", "호텔배정관리 화면 기준 사용자 흐름", 2, total)
    items = [
        ("01", "페이지 목적", "행사별 호텔 배정과 방갯수 관리"),
        ("02", "전체 업무 흐름", "행사 확인 → 차량 선택 → 호텔 추가/수정 → 저장"),
        ("03", "상단 행사 정보", "통합행사코드, 상품명, 출발일, 인원, 상태"),
        ("04", "좌측 차량/메모", "차량별 방갯수 확인과 호텔행사메모"),
        ("05", "우측 호텔배정", "일차별 호텔지역, 호텔명, 방갯수"),
        ("06", "호텔 추가", "호텔지역 선택, 호텔명 로딩, 추가"),
        ("07", "배정 목록 관리", "Excel, 업데이트, 삭제"),
        ("08", "저장/삭제 처리", "mode=save, mode=del"),
        ("09", "호텔명 로딩", "get_hotel.php와 chosen 업데이트"),
        ("10", "체크리스트", "저장 전후 확인사항"),
        ("11", "FAQ", "자주 확인할 내용"),
    ]
    box_w, box_h = Inches(6.0), Inches(0.66)
    sx, sy = Inches(0.5), Inches(1.10)
    for i, (no, title, sub) in enumerate(items):
        col, row = i % 2, i // 2
        x = sx + (box_w + Inches(0.3)) * col
        y = sy + (box_h + Inches(0.12)) * row
        add_rect(s, x, y, box_w, box_h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, x, y, Inches(0.85), box_h, NAVY)
        add_text(s, x, y, Inches(0.85), box_h, no, size=18, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(1.05), y + Inches(0.09), box_w - Inches(1.2), Inches(0.25),
                 title, size=12.5, bold=True, color=DARK_TXT)
        add_text(s, x + Inches(1.05), y + Inches(0.39), box_w - Inches(1.2), Inches(0.22),
                 sub, size=9.3, color=GRAY_TXT)
    add_footer(s)


def slide_purpose(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "01. 페이지 목적", "행사별 호텔과 방갯수를 배정하는 화면", 3, total)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.9), LIGHT_BLUE)
    add_text(s, Inches(0.7), Inches(1.23), Inches(11.9), Inches(0.58),
             "호텔배정관리 화면은 통합행사코드 기준으로 차량/서브행사를 선택하고, 일차별 호텔과 방갯수를 등록합니다.",
             size=15, bold=True, color=NAVY)
    cards = [
        ("행사 확인", "기본 정보", ["통합행사코드 확인", "상품명과 출발일 확인", "예약/행사 상태 확인"], NAVY),
        ("배정 입력", "일차별 호텔", ["호텔지역 선택", "호텔명 선택", "방갯수 입력 후 추가"], GREEN),
        ("관리", "수정/삭제/저장", ["기존 배정 목록 확인", "방갯수 업데이트", "삭제와 전체 저장"], ORANGE),
    ]
    for i, (label, title, items, color) in enumerate(cards):
        x = Inches(0.55 + i * 4.18)
        add_rect(s, x, Inches(2.35), Inches(3.9), Inches(4.35), WHITE, line=color)
        add_text(s, x + Inches(0.2), Inches(2.55), Inches(1.5), Inches(0.42),
                 label, size=14, bold=True, color=color)
        add_text(s, x + Inches(0.2), Inches(3.05), Inches(3.45), Inches(0.45),
                 title, size=20, bold=True, color=DARK_TXT)
        bullet(s, x + Inches(0.25), Inches(3.75), Inches(3.4), items,
               size=12.2, bullet_color=color, gap=0.48)
    add_footer(s)


def slide_flow(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "02. 전체 업무 흐름", "행사 확인부터 호텔배정 저장까지", 4, total)
    steps = [
        ("1", "행사 정보 확인", "통합행사코드, 상품명, 출발일, 투어정원, 예약인원을 확인합니다.", BLUE),
        ("2", "차량 선택", "좌측 차량/서브행사 라디오 버튼을 선택해 배정 대상을 바꿉니다.", GREEN),
        ("3", "호텔 메모", "필요 시 호텔행사메모에 운영 참고사항을 입력합니다.", TEAL),
        ("4", "호텔 추가", "일차별로 호텔지역, 호텔명, 방갯수를 선택하고 추가합니다.", ORANGE),
        ("5", "목록 관리", "추가된 호텔 배정의 방갯수를 업데이트하거나 삭제합니다.", RED),
        ("6", "최종 저장", "호텔배정저장 버튼으로 전체 내용을 저장합니다.", NAVY),
    ]
    for i, (no, title, body, color) in enumerate(steps):
        row, col = divmod(i, 3)
        x = Inches(0.7 + col * 4.15)
        y = Inches(1.35 + row * 2.35)
        step_card(s, x, y, Inches(3.65), Inches(1.75), no, title, body, color)
    add_footer(s)


def slide_top_info(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "03. 상단 행사 정보", "배정 전 반드시 확인하는 기준값", 5, total)
    table_like(
        s, Inches(0.65), Inches(1.25), Inches(12.0), Inches(2.65),
        ["항목", "화면 표시", "사용 목적"],
        [
            ["통합행사코드", "grand_eCode", "어떤 통합행사에 호텔을 배정하는지 확인"],
            ["상품명", "[p_code] p_name", "배정 대상 상품 확인"],
            ["출발일", "stDate", "일차별 숙박일 계산 기준"],
            ["투어정원/예약인원", "tour_pcnt / 예약 count", "방갯수 산정 참고"],
            ["예약상태/행사상태", "읽기 전용 라디오", "현재 행사 운영 상태 확인"],
        ],
        [0.20, 0.30, 0.50],
    )
    add_rect(s, Inches(0.65), Inches(4.55), Inches(12.0), Inches(1.10), LIGHT_YELLOW)
    add_text(s, Inches(0.9), Inches(4.72), Inches(11.4), Inches(0.32),
             "주의", size=15, bold=True, color=ORANGE)
    add_text(s, Inches(0.9), Inches(5.13), Inches(11.4), Inches(0.32),
             "상단 정보는 대부분 읽기 전용입니다. 잘못된 행사라면 이전 통합행사관리 화면에서 다시 선택해야 합니다.",
             size=12.6, color=DARK_TXT)
    add_footer(s)


def slide_left(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "04. 좌측 차량/메모 영역", "차량별 방갯수 확인과 호텔행사메모", 6, total)
    step_card(s, Inches(0.65), Inches(1.25), Inches(3.85), Inches(1.55),
              "A", "차량 선택", "차량 라디오 버튼을 클릭하면 gscode가 해당 서브행사코드로 바뀌고 화면이 다시 로드됩니다.", BLUE)
    step_card(s, Inches(4.75), Inches(1.25), Inches(3.85), Inches(1.55),
              "B", "방갯수 확인", "차량별로 hotelroom_assign 기준 방갯수가 표시됩니다.", GREEN)
    step_card(s, Inches(8.85), Inches(1.25), Inches(3.85), Inches(1.55),
              "C", "호텔행사메모", "호텔 운영 관련 메모를 입력하고 호텔배정저장 시 함께 저장합니다.", ORANGE)
    table_like(
        s, Inches(0.65), Inches(3.45), Inches(12.0), Inches(2.05),
        ["기능", "관련 코드", "설명"],
        [
            ["차량 목록", "buslist()", "차량/서브행사별 방갯수 라디오 출력"],
            ["차량 변경", "selectcar(subcode)", "gscode를 바꾸고 frmhotel 재전송"],
            ["메모", "hotelEventMemo", "hotel_assign 저장 시 h_memo 값으로 반영"],
        ],
        [0.20, 0.30, 0.50],
    )
    add_footer(s)


def slide_right(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "05. 우측 행사호텔배정", "일차별 호텔지역, 호텔명, 방갯수 입력", 7, total)
    table_like(
        s, Inches(0.65), Inches(1.25), Inches(12.0), Inches(2.35),
        ["화면 요소", "입력/표시", "역할"],
        [
            ["일차", "1일차, 2일차, 추가일차", "숙박일 구분"],
            ["호텔지역", "hoarea[]", "지역 선택 후 호텔명 목록 로딩"],
            ["호텔명", "hname[]", "get_hotel.php 결과에서 선택"],
            ["방갯수", "roompp[]", "해당 호텔에 배정할 방 수"],
            ["추가", "btnadd", "선택값을 폼으로 전송하여 저장 처리"],
        ],
        [0.20, 0.28, 0.52],
    )
    add_rect(s, Inches(0.65), Inches(4.25), Inches(12.0), Inches(0.78), LIGHT_GREEN)
    add_text(s, Inches(0.9), Inches(4.40), Inches(11.4), Inches(0.35),
             "ADD 상품은 product_details 일차 기준으로 여러 일차가 출력되고, 일반 상품은 추가일차 영역으로 호텔을 등록합니다.",
             size=12.7, bold=True, color=GREEN)
    add_footer(s)


def slide_add(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "06. 호텔 추가 방법", "호텔지역 선택 후 호텔명을 불러와 추가", 8, total)
    steps = [
        ("1", "호텔지역 선택", "호텔지역 드롭다운에서 지역 코드를 선택합니다.", BLUE),
        ("2", "호텔명 로딩", "get_hotel.php?code1=선택값 호출로 호텔명 목록이 채워집니다.", GREEN),
        ("3", "방갯수 입력", "roompp[]에 배정할 방갯수를 입력합니다.", TEAL),
        ("4", "추가 클릭", "추가 버튼을 누르면 frmhotel이 전송되어 저장 로직을 실행합니다.", ORANGE),
    ]
    for i, (no, title, body, color) in enumerate(steps):
        row, col = divmod(i, 2)
        x = Inches(0.65 + col * 6.15)
        y = Inches(1.35 + row * 2.25)
        step_card(s, x, y, Inches(5.75), Inches(1.65), no, title, body, color)
    add_footer(s, "D:\\www\\hellousa\\admin\\hotel_assign_m.php / get_hotel.php")


def slide_manage(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "07. 배정 목록 관리", "Excel, 업데이트, 삭제", 9, total)
    table_like(
        s, Inches(0.65), Inches(1.20), Inches(12.0), Inches(2.75),
        ["목록 열", "내용", "사용 방법"],
        [
            ["체크박스", "seq[] 선택", "전체선택 또는 개별 선택"],
            ["서브행사코드", "sub_eCode", "차량/서브행사 기준"],
            ["선택호텔명", "getHotelfInfo(hotel_code)", "저장된 호텔명 표시"],
            ["방갯수", "pcnt[]", "수량 변경 후 업데이트"],
            ["Action", "업데이트 / 삭제", "방갯수 업데이트 또는 행 삭제"],
        ],
        [0.18, 0.30, 0.52],
    )
    cards = [
        ("Excel", "DataTables 버튼으로 현재 목록을 엑셀로 내보냅니다.", BLUE),
        ("업데이트", "btnup 클릭 시 update_cnt.php로 방갯수와 일차 정보를 AJAX 전송합니다.", GREEN),
        ("삭제", "btndel 클릭 시 mode=del, no=seq_no로 폼을 전송해 삭제합니다.", RED),
    ]
    for i, (title, body, color) in enumerate(cards):
        step_card(s, Inches(0.65 + i * 4.1), Inches(4.55), Inches(3.85), Inches(1.55),
                  chr(65 + i), title, body, color)
    add_footer(s)


def slide_save(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "08. 저장/삭제 처리", "mode 값에 따라 저장과 삭제가 분기됩니다", 10, total)
    table_like(
        s, Inches(0.65), Inches(1.25), Inches(12.0), Inches(2.55),
        ["동작", "mode / 값", "처리"],
        [
            ["호텔배정저장", "mode=save", "선택된 호텔, 방갯수, 메모를 hotel_assign에 저장/갱신"],
            ["추가", "기본 submit", "호텔지역/호텔명/방갯수 입력값을 전송"],
            ["삭제", "mode=del, no=seq_no", "hotel_assign에서 해당 seq_no 삭제"],
            ["방갯수 업데이트", "AJAX update_cnt.php", "pcnt, day, sub_eCode 기준 수량 업데이트"],
        ],
        [0.24, 0.28, 0.48],
    )
    add_rect(s, Inches(0.65), Inches(4.60), Inches(12.0), Inches(0.95), LIGHT_YELLOW)
    add_text(s, Inches(0.9), Inches(4.78), Inches(11.4), Inches(0.30),
             "저장 전 확인", size=15, bold=True, color=ORANGE)
    add_text(s, Inches(0.9), Inches(5.16), Inches(11.4), Inches(0.28),
             "호텔배정저장은 confirm 창을 거친 뒤 실행됩니다. 삭제는 선택한 Action 행의 seq_no 기준으로 처리됩니다.",
             size=12.2, color=DARK_TXT)
    add_footer(s)


def slide_ajax(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "09. 호텔명 로딩 구조", "호텔지역 선택 시 get_hotel.php로 호텔명을 조회", 11, total)
    add_rect(s, Inches(0.65), Inches(1.15), Inches(5.85), Inches(4.95), WHITE, line=NAVY)
    add_text(s, Inches(0.9), Inches(1.38), Inches(5.25), Inches(0.42),
             "화면 동작", size=22, bold=True, color=NAVY)
    bullet(s, Inches(0.95), Inches(2.02), Inches(5.1), [
        "hoarea 변경 이벤트 발생",
        "선택값을 code1로 전달",
        "get_hotel.php?code1=값 호출",
        "응답 JSON으로 telarea 옵션 갱신",
        "chosen:updated 실행",
    ], size=12.5, bullet_color=NAVY, gap=0.50)
    add_rect(s, Inches(6.85), Inches(1.15), Inches(5.85), Inches(4.95), WHITE, line=GREEN)
    add_text(s, Inches(7.1), Inches(1.38), Inches(5.25), Inches(0.42),
             "데이터 기준", size=22, bold=True, color=GREEN)
    bullet(s, Inches(7.15), Inches(2.02), Inches(5.1), [
        "product_hotel 조회",
        "p_typem = code1",
        "u_type in ('1','3')",
        "h_name 오름차순",
        "h_code / h_name 반환",
    ], size=12.5, bullet_color=GREEN, gap=0.50)
    add_footer(s, "D:\\www\\hellousa\\admin\\hotel_assign_m.php / get_hotel.php")


def slide_checklist(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "10. 체크리스트", "저장 전후 확인 사항", 12, total)
    groups = [
        ("저장 전", ["통합행사코드 확인", "차량/서브행사 선택 확인", "호텔지역/호텔명 선택 확인"], BLUE),
        ("입력 중", ["방갯수 입력", "일차별 호텔 중복 확인", "호텔행사메모 입력 여부 확인"], GREEN),
        ("저장 후", ["선택호텔명 표시 확인", "방갯수 표시 확인", "Excel 내보내기 필요 여부 확인"], ORANGE),
        ("수정/삭제", ["업데이트 버튼 대상 행 확인", "삭제 버튼 대상 seq_no 확인", "삭제 후 목록 재확인"], NAVY),
    ]
    for i, (title, items, color) in enumerate(groups):
        row, col = divmod(i, 2)
        x = Inches(0.65 + col * 6.15)
        y = Inches(1.25 + row * 2.45)
        add_rect(s, x, y, Inches(5.75), Inches(2.0), WHITE, line=color)
        add_rect(s, x, y, Inches(5.75), Inches(0.45), color)
        add_text(s, x + Inches(0.18), y + Inches(0.08), Inches(5.4), Inches(0.25),
                 title, size=13.5, bold=True, color=WHITE)
        bullet(s, x + Inches(0.25), y + Inches(0.72), Inches(5.25), items,
               size=12.2, bullet_color=color, gap=0.40)
    add_footer(s)


def slide_faq(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "11. FAQ", "사용 중 자주 확인할 내용", 13, total)
    faqs = [
        ("Q1", "호텔명이 안 보일 때", "호텔지역 선택 후 get_hotel.php가 code1 값을 받아 product_hotel을 조회하는지 확인합니다."),
        ("Q2", "일차가 안 보일 때", "ADD 상품은 product_details의 day 기준으로 일차가 출력됩니다."),
        ("Q3", "방갯수 업데이트가 안 될 때", "pcnt 값이 1 이상인지 확인합니다. btnup은 update_cnt.php로 AJAX 전송합니다."),
        ("Q4", "삭제가 안 될 때", "삭제 버튼의 seq_no가 no hidden 값으로 들어가는지 확인합니다."),
        ("Q5", "저장 전 확인", "차량/서브행사 선택, 호텔명, 방갯수, 호텔행사메모를 다시 확인합니다."),
    ]
    y = Inches(1.15)
    for q, title, ans in faqs:
        add_rect(s, Inches(0.65), y, Inches(12.0), Inches(0.86), WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, Inches(0.65), y, Inches(0.85), Inches(0.86), NAVY)
        add_text(s, Inches(0.65), y, Inches(0.85), Inches(0.86),
                 q, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.7), y + Inches(0.08), Inches(10.6), Inches(0.30),
                 title, size=12.5, bold=True, color=DARK_TXT)
        add_text(s, Inches(1.7), y + Inches(0.44), Inches(10.6), Inches(0.25),
                 ans, size=10.4, color=GRAY_TXT)
        y += Inches(0.96)
    add_footer(s)


def slide_thanks(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    for x, y, w in [(Inches(-2), Inches(-2), Inches(6)), (Inches(9.5), Inches(4), Inches(6))]:
        deco = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, w)
        deco.fill.solid()
        deco.fill.fore_color.rgb = NAVY_DARK
        deco.line.fill.background()
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0),
             "Thank You", size=66, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.6),
             "호텔배정관리는 차량/서브행사 선택 후 일차별 호텔과 방갯수를 저장하는 흐름입니다.",
             size=19, color=RGBColor(0xCB, 0xD3, 0xE8))
    add_rect(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.45), WHITE)
    add_text(s, Inches(1.0), Inches(4.65), Inches(11), Inches(0.35),
             "확인 파일", size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(5.05), Inches(11), Inches(0.55),
             "D:\\www\\hellousa\\admin\\hotel_assign_m.php / get_hotel.php / update_cnt.php",
             size=14, bold=True, color=DARK_TXT)
    add_text(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
             "푸른투어 인트라넷 · hotel assignment manual",
             size=11, color=RGBColor(0xCB, 0xD3, 0xE8), align=PP_ALIGN.CENTER)


@dataclass
class PreviewSlide:
    title: str
    subtitle: str
    rows: Sequence[Tuple[str, str]]
    dark: bool = False


def rgb_tuple(color: RGBColor):
    return (color[0], color[1], color[2])


def load_font(size: int, bold: bool = False):
    path = r"C:/Windows/Fonts/malgunbd.ttf" if bold else r"C:/Windows/Fonts/malgun.ttf"
    if os.path.exists(path):
        return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def wrap_text(draw, text, font, max_width):
    words = str(text).split()
    if not words:
        return [""]
    lines, cur = [], words[0]
    for word in words[1:]:
        trial = cur + " " + word
        if draw.textbbox((0, 0), trial, font=font)[2] <= max_width:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    lines.append(cur)
    return lines


def draw_preview(idx: int, total: int, spec: PreviewSlide):
    W, H = 1920, 1080
    img = Image.new("RGB", (W, H), rgb_tuple(GRAY_BG))
    d = ImageDraw.Draw(img)
    if spec.dark:
        d.rectangle([0, 0, W, H], fill=rgb_tuple(NAVY))
        d.ellipse([1450, -220, 2200, 530], fill=rgb_tuple(NAVY_DARK))
        d.ellipse([-240, 420, 430, 1090], fill=rgb_tuple(NAVY_DARK))
        d.text((115, 280), spec.title, fill=rgb_tuple(WHITE), font=load_font(76, True))
        d.text((120, 415), spec.subtitle, fill=(203, 211, 232), font=load_font(25))
    else:
        d.rectangle([0, 0, W, 130], fill=rgb_tuple(NAVY))
        d.text((58, 28), spec.title, fill=rgb_tuple(WHITE), font=load_font(54, True))
        d.text((62, 86), spec.subtitle, fill=(203, 211, 232), font=load_font(18))
        d.text((1700, 48), f"{idx} / {total}", fill=rgb_tuple(WHITE), font=load_font(18))
        y = 170
        for i, (head, body) in enumerate(spec.rows):
            x = 90 + (i % 2) * 870
            if i and i % 2 == 0:
                y += 240
            d.rounded_rectangle([x, y, x + 790, y + 185], radius=8,
                                fill=rgb_tuple(WHITE), outline=(221, 221, 221), width=2)
            d.rectangle([x, y, x + 112, y + 185], fill=rgb_tuple(NAVY))
            d.text((x + 34, y + 58), str(i + 1), fill=rgb_tuple(WHITE), font=load_font(40, True))
            d.text((x + 145, y + 25), head, fill=rgb_tuple(DARK_TXT), font=load_font(27, True))
            line_y = y + 76
            for line in wrap_text(d, body, load_font(23), 590)[:3]:
                d.text((x + 145, line_y), line, fill=rgb_tuple(GRAY_TXT), font=load_font(23))
                line_y += 34
        d.rectangle([0, 1030, W, H], fill=rgb_tuple(GRAY_BG))
        d.text((58, 1040), "푸른투어 인트라넷 · 호텔배정관리 사용자 매뉴얼",
               fill=rgb_tuple(GRAY_TXT), font=load_font(18))
    out = f"{PREVIEW_PREFIX}{idx:02d}.png"
    img.save(out)
    return out


def render_previews(total: int):
    specs = [
        PreviewSlide("호텔배정관리", "Hotel Assignment Workflow Guide", [], True),
        PreviewSlide("목차 · CONTENTS", "호텔배정관리 화면 기준 사용자 흐름", [("페이지 목적", "행사별 호텔과 방갯수 관리"), ("업무 흐름", "행사 확인 → 차량 선택 → 호텔 저장")]),
        PreviewSlide("01. 페이지 목적", "행사별 호텔과 방갯수 배정", [("행사 확인", "통합행사코드와 출발일 확인"), ("배정 입력", "일차별 호텔과 방갯수"), ("관리", "업데이트/삭제/저장")]),
        PreviewSlide("02. 전체 업무 흐름", "행사 확인부터 저장까지", [("행사 정보", "상단 기준값 확인"), ("차량 선택", "좌측 서브행사 선택"), ("호텔 추가", "지역/호텔/방갯수 입력")]),
        PreviewSlide("03. 상단 행사 정보", "배정 전 확인하는 기준값", [("통합행사코드", "배정 대상 행사"), ("출발일", "숙박일 계산 기준"), ("상태", "예약/행사 상태 확인")]),
        PreviewSlide("04. 좌측 차량/메모 영역", "차량별 방갯수와 호텔행사메모", [("차량 선택", "gscode 변경"), ("방갯수", "차량별 방 수 확인"), ("메모", "호텔행사메모 저장")]),
        PreviewSlide("05. 우측 행사호텔배정", "일차별 호텔지역, 호텔명, 방갯수", [("일차", "숙박일 구분"), ("호텔 선택", "지역 선택 후 호텔명 로딩"), ("추가", "방갯수 입력 후 추가")]),
        PreviewSlide("06. 호텔 추가 방법", "호텔지역 선택 후 호텔명 조회", [("지역 선택", "hoarea 변경"), ("호텔명 로딩", "get_hotel.php 호출"), ("추가", "frmhotel 전송")]),
        PreviewSlide("07. 배정 목록 관리", "Excel, 업데이트, 삭제", [("Excel", "현재 목록 내보내기"), ("업데이트", "방갯수 AJAX 수정"), ("삭제", "seq_no 기준 삭제")]),
        PreviewSlide("08. 저장/삭제 처리", "mode 값에 따른 처리", [("저장", "mode=save"), ("삭제", "mode=del"), ("업데이트", "update_cnt.php")]),
        PreviewSlide("09. 호텔명 로딩 구조", "get_hotel.php와 chosen 업데이트", [("code1", "GET 파라미터"), ("product_hotel", "p_typem 기준 조회"), ("chosen", "옵션 갱신")]),
        PreviewSlide("10. 체크리스트", "저장 전후 확인사항", [("저장 전", "행사/차량/호텔 확인"), ("입력 중", "방갯수 확인"), ("저장 후", "목록 반영 확인")]),
        PreviewSlide("11. FAQ", "자주 확인할 내용", [("호텔명 없음", "code1/get_hotel 확인"), ("일차 없음", "product_details 확인"), ("수정 안 됨", "pcnt/update_cnt 확인")]),
        PreviewSlide("Thank You", "호텔배정관리 사용자 매뉴얼", [], True),
    ]
    return [draw_preview(i + 1, total, spec) for i, spec in enumerate(specs)]


def build():
    total = 14
    slide_cover()
    slide_toc(total)
    slide_purpose(total)
    slide_flow(total)
    slide_top_info(total)
    slide_left(total)
    slide_right(total)
    slide_add(total)
    slide_manage(total)
    slide_save(total)
    slide_ajax(total)
    slide_checklist(total)
    slide_faq(total)
    slide_thanks(total)
    try:
        prs.save(OUT)
        saved = OUT
    except PermissionError:
        saved = OUT.replace(".pptx", f"_v{int(time.time())}.pptx")
        prs.save(saved)
    previews = render_previews(total)
    print("SAVED:", saved)
    print("PREVIEWS:", len(previews))
    for path in previews:
        print(path)


if __name__ == "__main__":
    build()
