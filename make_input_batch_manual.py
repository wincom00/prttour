# -*- coding: utf-8 -*-
"""
input_batch.php (스마트 그룹 예약 등록) 순수 사용자 설명서 PPTX 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ===== 색상 팔레트 (input_batch.php 와 동일톤) =====
NAVY      = RGBColor(0x13, 0x11, 0x76)   # #131176
NAVY_DARK = RGBColor(0x0A, 0x09, 0x55)
GREEN     = RGBColor(0x28, 0xA7, 0x45)
TEAL      = RGBColor(0x20, 0xC9, 0x97)
ORANGE    = RGBColor(0xFD, 0x7E, 0x14)
RED       = RGBColor(0xDC, 0x35, 0x45)
GRAY_BG   = RGBColor(0xF8, 0xF9, 0xFA)
GRAY_TXT  = RGBColor(0x6C, 0x75, 0x7D)
DARK_TXT  = RGBColor(0x21, 0x25, 0x29)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE= RGBColor(0xE7, 0xF3, 0xFF)
LIGHT_GREEN=RGBColor(0xE8, 0xF7, 0xEE)
LIGHT_YEL = RGBColor(0xFF, 0xF8, 0xE1)

KFONT = "맑은 고딕"

# ===== 프레젠테이션 기본 설정 =====
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------- 헬퍼 ----------
def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh

def set_text(shape, text, size=14, bold=False, color=DARK_TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=KFONT):
    tf = shape.text_frame
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = str(line)
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color

def add_text(slide, x, y, w, h, text, **kw):
    sh = slide.shapes.add_textbox(x, y, w, h)
    set_text(sh, text, **kw)
    return sh

def add_header_bar(slide, title, subtitle=None, page=None, total=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), NAVY)
    add_text(slide, Inches(0.4), Inches(0.12), Inches(10), Inches(0.45),
             title, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.52), Inches(10), Inches(0.35),
                 subtitle, size=11, color=RGBColor(0xCB, 0xD3, 0xE8))
    if page is not None:
        add_text(slide, Inches(11.5), Inches(0.28), Inches(1.6), Inches(0.4),
                 f"{page} / {total}", size=11, color=WHITE,
                 align=PP_ALIGN.RIGHT, bold=True)

def add_footer(slide):
    add_rect(slide, 0, Inches(7.15), SW, Inches(0.35), GRAY_BG)
    add_text(slide, Inches(0.4), Inches(7.18), Inches(8), Inches(0.3),
             "푸른투어 Partner Portal · 스마트 그룹 예약 등록 사용 설명서",
             size=9, color=GRAY_TXT)
    add_text(slide, Inches(9), Inches(7.18), Inches(4), Inches(0.3),
             "input_batch.php", size=9, color=GRAY_TXT, align=PP_ALIGN.RIGHT)

def bullet(slide, x, y, w, h, items, size=14, color=DARK_TXT, bullet_color=NAVY, gap=0.32):
    """체크 아이콘 + 텍스트 형태"""
    cur_y = y
    for it in items:
        # 체크 아이콘
        ic = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cur_y + Emu(40000),
                                    Inches(0.22), Inches(0.22))
        ic.fill.solid(); ic.fill.fore_color.rgb = bullet_color
        ic.line.fill.background()
        set_text(ic, "✓", size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 본문
        tx = slide.shapes.add_textbox(x + Inches(0.32), cur_y,
                                      w - Inches(0.32), Inches(gap))
        set_text(tx, it, size=size, color=color)
        cur_y = cur_y + Inches(gap)

# =========================================================
# 슬라이드 1 — 표지
# =========================================================
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    # 배경 그라데이션 느낌(상단 NAVY, 하단 GRAY_BG)
    add_rect(s, 0, 0, SW, SH, GRAY_BG)
    add_rect(s, 0, 0, SW, Inches(4.5), NAVY)
    # 장식 도형
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(10.2), Inches(-1.5),
                              Inches(5), Inches(5))
    deco.fill.solid(); deco.fill.fore_color.rgb = NAVY_DARK
    deco.line.fill.background()
    deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(-1.5), Inches(2.8),
                               Inches(3.5), Inches(3.5))
    deco2.fill.solid(); deco2.fill.fore_color.rgb = NAVY_DARK
    deco2.line.fill.background()

    # 라벨
    add_text(s, Inches(0.8), Inches(0.7), Inches(6), Inches(0.4),
             "PARTNER PORTAL · 사용자 매뉴얼",
             size=14, bold=True, color=RGBColor(0xCB, 0xD3, 0xE8))
    # 메인 타이틀
    add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(1.2),
             "스마트 그룹 예약 등록",
             size=54, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(2.55), Inches(11), Inches(0.8),
             "Smart Group Reservation – Excel/CSV Batch Input",
             size=22, color=RGBColor(0xCB, 0xD3, 0xE8))
    add_text(s, Inches(0.8), Inches(3.3), Inches(11), Inches(0.6),
             "엑셀(.xlsx/.xls) · CSV 파일 업로드만으로 그룹 예약을 자동 등록",
             size=16, color=WHITE)

    # 정보카드
    add_rect(s, Inches(0.8), Inches(5.0), Inches(3.8), Inches(1.6),
             WHITE, line=NAVY)
    add_text(s, Inches(1.0), Inches(5.1), Inches(3.5), Inches(0.4),
             "📂 대상 페이지", size=12, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(5.45), Inches(3.5), Inches(0.5),
             "input_batch.php", size=18, bold=True, color=DARK_TXT)
    add_text(s, Inches(1.0), Inches(5.95), Inches(3.5), Inches(0.4),
             "협력사 전용 일괄 등록", size=12, color=GRAY_TXT)

    add_rect(s, Inches(4.85), Inches(5.0), Inches(3.8), Inches(1.6),
             WHITE, line=GREEN)
    add_text(s, Inches(5.05), Inches(5.1), Inches(3.5), Inches(0.4),
             "👥 사용 대상", size=12, bold=True, color=GREEN)
    add_text(s, Inches(5.05), Inches(5.45), Inches(3.5), Inches(0.5),
             "협력사 / 파트너", size=18, bold=True, color=DARK_TXT)
    add_text(s, Inches(5.05), Inches(5.95), Inches(3.5), Inches(0.4),
             "그룹 단체예약 담당자", size=12, color=GRAY_TXT)

    add_rect(s, Inches(8.9), Inches(5.0), Inches(3.8), Inches(1.6),
             WHITE, line=ORANGE)
    add_text(s, Inches(9.1), Inches(5.1), Inches(3.5), Inches(0.4),
             "📑 지원 파일", size=12, bold=True, color=ORANGE)
    add_text(s, Inches(9.1), Inches(5.45), Inches(3.5), Inches(0.5),
             ".xlsx · .xls · .csv", size=18, bold=True, color=DARK_TXT)
    add_text(s, Inches(9.1), Inches(5.95), Inches(3.5), Inches(0.4),
             "다중 파일 / 다중 시트", size=12, color=GRAY_TXT)

    add_text(s, Inches(0.8), Inches(6.85), Inches(11), Inches(0.4),
             "푸른투어 PARTNER PORTAL", size=11, bold=True, color=NAVY)

slide_cover()

# =========================================================
# 슬라이드 2 — 목차
# =========================================================
def slide_toc(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "목차 · CONTENTS", "본 매뉴얼 구성", page=2, total=total)

    items = [
        ("01", "이 페이지로 무엇을 할 수 있나요?", "스마트 그룹 예약 등록 개요"),
        ("02", "접속 방법 (로그인 · 권한)",       "협력사 계정 로그인 안내"),
        ("03", "전체 화면 구성",                  "메뉴/버튼 한눈에 보기"),
        ("04", "1단계 · 파일 준비하기",            "파일명 규칙과 양식"),
        ("05", "2단계 · 파일 업로드 & 분석",        "선택 → 분석 시작"),
        ("06", "3단계 · 시트별 그룹 예약 정보 확인", "자동 추출 결과 검토"),
        ("07", "4단계 · 여행자 목록 편집",          "추가 · 수정 · 삭제 · 복제"),
        ("08", "5단계 · 일괄(다중) 편집",           "성별/객실/선택 저장"),
        ("09", "6단계 · DB 저장",                  "시트별 저장 & 결과 확인"),
        ("10", "자주 묻는 질문 / 오류 대응",        "FAQ & Troubleshooting"),
    ]
    cols = 2
    box_w = Inches(6.0)
    box_h = Inches(0.95)
    sx = Inches(0.5)
    sy = Inches(1.15)
    gap_x = Inches(0.3)
    gap_y = Inches(0.15)

    for i, (no, t, sub) in enumerate(items):
        col = i % cols
        row = i // cols
        x = sx + (box_w + gap_x) * col
        y = sy + (box_h + gap_y) * row
        add_rect(s, x, y, box_w, box_h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        # 번호 박스
        add_rect(s, x, y, Inches(0.95), box_h, NAVY)
        add_text(s, x, y, Inches(0.95), box_h, no,
                 size=24, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(1.1), y + Inches(0.12), box_w - Inches(1.2), Inches(0.4),
                 t, size=14, bold=True, color=DARK_TXT)
        add_text(s, x + Inches(1.1), y + Inches(0.5), box_w - Inches(1.2), Inches(0.4),
                 sub, size=10, color=GRAY_TXT)

    add_footer(s)

# =========================================================
# 슬라이드 3 — 무엇을 할 수 있나
# =========================================================
def slide_overview(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "01. 이 페이지로 무엇을 할 수 있나요?",
                   "Excel/CSV 파일 한 번 업로드 → 그룹 예약 자동 등록", page=3, total=total)

    # 한 줄 요약 박스
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.9), LIGHT_BLUE)
    add_text(s, Inches(0.7), Inches(1.2), Inches(12), Inches(0.7),
             "💡 협력사가 받은 \"여행자 명단\" 엑셀을 그대로 올리면, "
             "시스템이 자동으로 그룹·여행자 정보를 추출해 예약을 등록해 줍니다.",
             size=15, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    # 카드 3개
    cards = [
        ("🤖", "자동 인식", NAVY, [
            "엑셀의 모든 시트를 자동으로 읽음",
            "병합된 셀 자동 처리",
            "다양한 날짜 형식 인식",
            "그룹 정보 / 여행자 정보 자동 추출",
        ]),
        ("✏️", "실시간 편집", GREEN, [
            "여행자 추가 · 수정 · 삭제 · 복제",
            "성별 / 객실타입 일괄 변경",
            "검색 · 통계 자동 업데이트",
            "필요 시 선택분만 별도 그룹 저장",
        ]),
        ("💾", "즉시 저장", ORANGE, [
            "시트별로 독립적인 그룹 예약 등록",
            "대표예약자 단위로 자동 묶음",
            "DB 저장 후 결과 즉시 확인",
            "협력사 ID 기준으로 자동 귀속",
        ]),
    ]
    cx = Inches(0.5)
    cy = Inches(2.25)
    cw = Inches(4.05)
    ch = Inches(4.55)
    gap = Inches(0.15)
    for i, (icon, title, color, lst) in enumerate(cards):
        x = cx + (cw + gap) * i
        add_rect(s, x, cy, cw, ch, WHITE, line=RGBColor(0xE0, 0xE0, 0xE0))
        add_rect(s, x, cy, cw, Inches(0.85), color)
        add_text(s, x, cy, cw, Inches(0.85),
                 f"{icon}  {title}", size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        bullet(s, x + Inches(0.25), cy + Inches(1.05),
               cw - Inches(0.4), ch - Inches(1.1),
               lst, size=13, bullet_color=color, gap=0.55)

    add_footer(s)

# =========================================================
# 슬라이드 4 — 접속 방법
# =========================================================
def slide_login(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "02. 접속 방법", "로그인 · 권한", page=4, total=total)

    # 좌측: 단계 안내
    add_text(s, Inches(0.5), Inches(1.1), Inches(7), Inches(0.5),
             "▍ 접속 절차", size=18, bold=True, color=NAVY)

    steps = [
        ("1", "Partner Portal 로그인",
         "협력사용 ID/PW 로 [login.php] 에서 로그인합니다."),
        ("2", "메뉴 진입",
         "좌측/상단 메뉴에서 \"스마트 그룹 예약 등록\"(input_batch) 선택."),
        ("3", "권한 확인",
         "로그인하지 않았거나 \"가이드(guide)\" 권한일 경우 자동으로 다른 페이지로 이동됩니다."),
        ("4", "페이지 진입 완료",
         "상단에 \"스마트 그룹 예약 등록\" 헤더가 보이면 정상 접속."),
    ]
    cy = Inches(1.65)
    for i, (n, t, d) in enumerate(steps):
        y = cy + Inches(1.15) * i
        add_rect(s, Inches(0.5), y, Inches(0.9), Inches(0.95), NAVY)
        add_text(s, Inches(0.5), y, Inches(0.9), Inches(0.95), n,
                 size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.45), y, Inches(5.6), Inches(0.95),
                 WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_text(s, Inches(1.6), y + Inches(0.08), Inches(5.4), Inches(0.4),
                 t, size=14, bold=True, color=DARK_TXT)
        add_text(s, Inches(1.6), y + Inches(0.45), Inches(5.4), Inches(0.55),
                 d, size=11, color=GRAY_TXT)

    # 우측: 권한/주의 안내
    add_rect(s, Inches(7.5), Inches(1.1), Inches(5.3), Inches(2.6),
             LIGHT_YEL, line=ORANGE)
    add_text(s, Inches(7.7), Inches(1.2), Inches(5), Inches(0.5),
             "⚠ 권한 안내", size=16, bold=True, color=ORANGE)
    bullet(s, Inches(7.7), Inches(1.7), Inches(5), Inches(2),
           [
            "협력사 계정으로만 사용 가능",
            "가이드(guide) 권한은 메모 페이지로 이동",
            "비로그인 시 로그인 페이지로 자동 이동",
            "등록된 예약은 로그인 협력사 ID로 귀속",
           ], size=12, bullet_color=ORANGE, gap=0.45)

    add_rect(s, Inches(7.5), Inches(3.85), Inches(5.3), Inches(3.0),
             LIGHT_GREEN, line=GREEN)
    add_text(s, Inches(7.7), Inches(3.95), Inches(5), Inches(0.5),
             "✅ 정상 진입 시 보이는 화면", size=16, bold=True, color=GREEN)
    bullet(s, Inches(7.7), Inches(4.45), Inches(5), Inches(2.4),
           [
            "상단 큰 제목: \"스마트 그룹 예약 등록\"",
            "안내 박스: \"스마트 처리 기능\"",
            "파일 업로드 영역 (파란색 두근두근 점선 박스)",
            "[파일 선택] 파란 버튼",
            "분석 시작 버튼은 파일 선택 후 표시",
           ], size=12, bullet_color=GREEN, gap=0.43)

    add_footer(s)

# =========================================================
# 슬라이드 5 — 화면 구성
# =========================================================
def slide_layout(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "03. 전체 화면 구성", "메뉴/버튼 한눈에 보기", page=5, total=total)

    # 큰 화면 모형
    canvas_x, canvas_y = Inches(0.5), Inches(1.1)
    canvas_w, canvas_h = Inches(8.5), Inches(5.7)
    add_rect(s, canvas_x, canvas_y, canvas_w, canvas_h,
             GRAY_BG, line=RGBColor(0xCC, 0xCC, 0xCC))

    # 1) 페이지 헤더
    add_rect(s, canvas_x + Inches(0.2), canvas_y + Inches(0.2),
             canvas_w - Inches(0.4), Inches(0.7), WHITE, line=NAVY)
    add_text(s, canvas_x + Inches(0.4), canvas_y + Inches(0.27),
             canvas_w - Inches(0.6), Inches(0.55),
             "🪄 스마트 그룹 예약 등록  (페이지 헤더)",
             size=12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    # 2) 안내(파란박스)
    add_rect(s, canvas_x + Inches(0.2), canvas_y + Inches(1.05),
             canvas_w - Inches(0.4), Inches(0.55),
             LIGHT_BLUE)
    add_text(s, canvas_x + Inches(0.4), canvas_y + Inches(1.1),
             canvas_w - Inches(0.6), Inches(0.5),
             "💡 스마트 처리 기능 안내 (자동 인식 / 실시간 편집)",
             size=11, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

    # 3) 업로드 카드
    add_rect(s, canvas_x + Inches(0.2), canvas_y + Inches(1.75),
             canvas_w - Inches(0.4), Inches(1.45), WHITE,
             line=NAVY)
    add_text(s, canvas_x + Inches(0.4), canvas_y + Inches(1.85),
             canvas_w - Inches(0.6), Inches(0.4),
             "📤 파일 업로드 영역", size=12, bold=True, color=NAVY)
    add_text(s, canvas_x + Inches(0.4), canvas_y + Inches(2.25),
             canvas_w - Inches(0.6), Inches(0.4),
             "•  [파일 선택] 버튼 → .xlsx / .xls / .csv (다중 선택 가능)",
             size=10, color=DARK_TXT)
    add_text(s, canvas_x + Inches(0.4), canvas_y + Inches(2.6),
             canvas_w - Inches(0.6), Inches(0.4),
             "•  파일 선택 후 [선택 파일 분석 시작] 버튼이 나타남",
             size=10, color=DARK_TXT)

    # 4) 시트별 아코디언
    add_rect(s, canvas_x + Inches(0.2), canvas_y + Inches(3.3),
             canvas_w - Inches(0.4), Inches(1.1),
             WHITE, line=GREEN)
    add_text(s, canvas_x + Inches(0.4), canvas_y + Inches(3.4),
             canvas_w - Inches(0.6), Inches(0.4),
             "📑 시트별 그룹 예약 정보 (아코디언)",
             size=12, bold=True, color=GREEN)
    add_text(s, canvas_x + Inches(0.4), canvas_y + Inches(3.78),
             canvas_w - Inches(0.6), Inches(0.6),
             "각 시트마다 그룹 정보 폼 + [DB 저장] 버튼이 표시됩니다.",
             size=10, color=DARK_TXT)

    # 5) 여행자 목록
    add_rect(s, canvas_x + Inches(0.2), canvas_y + Inches(4.5),
             canvas_w - Inches(0.4), Inches(1.05), WHITE, line=ORANGE)
    add_text(s, canvas_x + Inches(0.4), canvas_y + Inches(4.6),
             canvas_w - Inches(0.6), Inches(0.4),
             "👥 여행자 목록 테이블 (검색/추가/일괄작업)",
             size=12, bold=True, color=ORANGE)
    add_text(s, canvas_x + Inches(0.4), canvas_y + Inches(4.95),
             canvas_w - Inches(0.6), Inches(0.55),
             "한글성명 · 영문성명 · 여권번호 · 생년월일 · 성별 · 연락처 · 객실 · 비고",
             size=10, color=DARK_TXT)

    # 우측 설명 라벨
    labels = [
        ("①", "페이지 헤더",  Inches(1.4)),
        ("②", "안내 박스",    Inches(2.05)),
        ("③", "파일 업로드",   Inches(2.85)),
        ("④", "시트별 그룹",   Inches(4.45)),
        ("⑤", "여행자 목록",   Inches(5.6)),
    ]
    add_rect(s, Inches(9.2), Inches(1.1), Inches(3.6), Inches(5.7),
             WHITE, line=NAVY)
    add_text(s, Inches(9.4), Inches(1.2), Inches(3.3), Inches(0.5),
             "🧭 화면 구성 요약", size=16, bold=True, color=NAVY)
    cur_y = Inches(1.75)
    descs = [
        "현재 페이지 제목",
        "자동인식/편집 기능 설명",
        "엑셀 · CSV 업로드 영역",
        "분석 결과(시트별 그룹폼+저장)",
        "여행자 추가/수정/일괄편집",
    ]
    for i, ((mark, name, _y), d) in enumerate(zip(labels, descs)):
        y = cur_y + Inches(0.85) * i
        add_text(s, Inches(9.4), y, Inches(0.6), Inches(0.4),
                 mark, size=22, bold=True, color=NAVY)
        add_text(s, Inches(10.0), y, Inches(2.7), Inches(0.4),
                 name, size=13, bold=True, color=DARK_TXT)
        add_text(s, Inches(10.0), y + Inches(0.35), Inches(2.7), Inches(0.5),
                 d, size=10, color=GRAY_TXT)

    add_footer(s)

# =========================================================
# 슬라이드 6 — 1단계 파일 준비
# =========================================================
def slide_step1(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "04. 1단계 · 파일 준비하기",
                   "올바른 파일명/양식이 자동인식 정확도를 높입니다", page=6, total=total)

    # 핵심 파일명 규칙
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.05),
             RGBColor(0xFF, 0xEC, 0xEC), line=RED)
    add_text(s, Inches(0.7), Inches(1.18), Inches(12), Inches(0.4),
             "🔥 가장 중요한 규칙 — 파일명",
             size=15, bold=True, color=RED)
    add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.55),
             "파일명을  XXX_{상품코드}.xlsx  형태로 변경해 주세요.   예) 단체명단_KR12345.xlsx",
             size=14, bold=True, color=DARK_TXT)

    # 좌: 권장 양식
    add_text(s, Inches(0.5), Inches(2.35), Inches(6), Inches(0.5),
             "▍ 권장 시트 구성", size=16, bold=True, color=NAVY)
    add_rect(s, Inches(0.5), Inches(2.85), Inches(6), Inches(3.9),
             WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))

    headers = ["한글성명", "영문성명", "여권번호", "생년월일", "성별", "연락처", "객실"]
    rows = [
        ["홍길동", "HONG GILDONG", "M12345678", "1980-01-15", "남", "010-0000-0000", "2r1p"],
        ["김영희", "KIM YOUNGHEE", "M87654321", "1985-03-22", "여", "010-1111-2222", "2r1p"],
        ["이철수", "LEE CHEOLSU", "M11112222", "1979-07-04", "남", "010-3333-4444", "1r1p"],
    ]
    cell_w = Inches(0.83)
    cell_h = Inches(0.4)
    tx = Inches(0.6)
    ty = Inches(2.95)
    # 헤더
    for ci, hd in enumerate(headers):
        add_rect(s, tx + cell_w * ci, ty, cell_w, cell_h, NAVY)
        add_text(s, tx + cell_w * ci, ty, cell_w, cell_h, hd,
                 size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 데이터
    for ri, row in enumerate(rows):
        for ci, v in enumerate(row):
            add_rect(s, tx + cell_w * ci, ty + cell_h * (ri + 1),
                     cell_w, cell_h, WHITE,
                     line=RGBColor(0xE0, 0xE0, 0xE0))
            add_text(s, tx + cell_w * ci, ty + cell_h * (ri + 1),
                     cell_w, cell_h, v, size=9, color=DARK_TXT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.6), Inches(5.0), Inches(6), Inches(1.7),
             [
              "• 첫 행은 컬럼명, 둘째 행부터 여행자 데이터",
              "• 한 시트 = 한 개의 그룹예약 으로 처리됨",
              "• 시트가 여러 개면 각 시트가 각각의 그룹예약이 됩니다",
              "• 셀 병합/빈 행/요약 행 모두 자동 보정",
             ], size=11, color=DARK_TXT)

    # 우: 인식되는 항목
    add_text(s, Inches(7), Inches(2.35), Inches(6), Inches(0.5),
             "▍ 자동 인식되는 정보", size=16, bold=True, color=GREEN)
    add_rect(s, Inches(7), Inches(2.85), Inches(5.8), Inches(3.9),
             LIGHT_GREEN, line=GREEN)
    bullet(s, Inches(7.2), Inches(3.0), Inches(5.5), Inches(3.5),
           [
            "그룹 정보  : 상품코드 / 상품명 / 출발일 / 도착일",
            "대표자       : 그룹대표자명 / 연락처 / 이메일",
            "여행자정보 : 한글·영문성명 / 여권번호 / 생년월일",
            "그 외          : 성별 / 연락처 / 이메일",
            "객실 정보    : 객실타입(2r1p/1r1p 등) / 객실번호",
            "메모           : 비고/특이사항",
            "파일명에서 상품코드 자동 추출 (XXX_코드.xlsx)",
           ], size=12, bullet_color=GREEN, gap=0.45)

    add_footer(s)

# =========================================================
# 슬라이드 7 — 2단계 업로드 & 분석
# =========================================================
def slide_step2(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "05. 2단계 · 파일 업로드 & 분석",
                   "선택 → 분석 시작 까지", page=7, total=total)

    # 단계
    flow = [
        ("①", "파일 선택",
         "[파일 선택] 버튼을 누르고 .xlsx / .xls / .csv 파일을 고릅니다.\n"
         "여러 파일을 동시에 선택할 수 있습니다."),
        ("②", "선택 파일 확인",
         "업로드 영역 아래에 \"선택된 파일 목록\"이 표시되며,\n"
         "필요 시 다시 [파일 선택] 으로 변경 가능합니다."),
        ("③", "분석 시작",
         "[선택 파일 분석 시작] 버튼을 클릭합니다.\n"
         "AI 분석 화면이 나타나고 잠시 후 결과가 표시됩니다."),
        ("④", "결과 확인",
         "오류가 있으면 빨간 박스, 성공 시 초록 박스로 안내됩니다.\n"
         "성공 시 자동으로 \"시트별 그룹 예약 정보\" 가 펼쳐집니다."),
    ]
    cy = Inches(1.15)
    cw = Inches(3.0)
    ch = Inches(3.4)
    gap = Inches(0.16)
    sx = Inches(0.5)
    colors = [NAVY, GREEN, ORANGE, RGBColor(0x6F, 0x42, 0xC1)]
    for i, (n, t, d) in enumerate(flow):
        x = sx + (cw + gap) * i
        add_rect(s, x, cy, cw, ch, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, x, cy, cw, Inches(0.7), colors[i])
        add_text(s, x, cy, cw, Inches(0.7), f"{n}  {t}",
                 size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.2), cy + Inches(0.85),
                 cw - Inches(0.4), ch - Inches(0.95),
                 d, size=11, color=DARK_TXT)
        # 화살표
        if i < 3:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       x + cw + Emu(20000),
                                       cy + ch / 2 - Inches(0.18),
                                       Inches(0.16), Inches(0.36))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = NAVY
            arrow.line.fill.background()

    # 하단 팁
    add_rect(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.05),
             LIGHT_YEL, line=ORANGE)
    add_text(s, Inches(0.7), Inches(4.9), Inches(12), Inches(0.5),
             "💡 분석 중 화면 안내", size=14, bold=True, color=ORANGE)
    bullet(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.5),
           [
            "분석 중에는 \"AI가 파일을 분석하고 있습니다…\" 로딩 메시지가 표시됩니다.",
            "오류 박스에 \"파일 형식이 올바르지 않습니다\" 가 뜨면 파일 확장자/내용 확인.",
            "성공 박스에 \"처리 완료\" 가 뜨면 시트별 결과 영역으로 자동 스크롤됩니다.",
           ], size=12, bullet_color=ORANGE, gap=0.42)

    add_footer(s)

# =========================================================
# 슬라이드 8 — 3단계 시트별 그룹 정보 확인
# =========================================================
def slide_step3(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "06. 3단계 · 시트별 그룹 예약 정보 확인",
                   "자동 추출된 결과를 검토 · 보완하세요", page=8, total=total)

    # 좌측 - 그룹 정보 폼 mockup
    fx, fy = Inches(0.5), Inches(1.1)
    fw, fh = Inches(7.5), Inches(5.7)
    add_rect(s, fx, fy, fw, fh, WHITE, line=GREEN)
    add_rect(s, fx, fy, fw, Inches(0.55), GREEN)
    add_text(s, fx + Inches(0.2), fy, fw - Inches(0.4), Inches(0.55),
             "📑  시트1   ▾   (예: 단체명단_KR12345)",
             size=13, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    fields = [
        ("상품코드",     "KR12345"),
        ("상품명/투어명", "한국 단풍여행 9일"),
        ("출발일",       "2026-10-12"),
        ("도착일",       "2026-10-20"),
        ("대표자",       "홍길동"),
        ("연락처",       "010-1234-5678"),
        ("대표 이메일",   "leader@partner.com"),
        ("총 예약금액",   "12,000.00"),
        ("메모",         "단체 식사 한식 위주"),
    ]
    cur_y = fy + Inches(0.8)
    col_h = Inches(0.5)
    for i, (lbl, v) in enumerate(fields):
        col = i % 2
        row = i // 2
        x = fx + Inches(0.25) + Inches(3.6) * col
        y = cur_y + (col_h + Inches(0.1)) * row
        add_text(s, x, y, Inches(1.4), col_h, lbl,
                 size=11, bold=True, color=GRAY_TXT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, x + Inches(1.45), y, Inches(2.05), col_h,
                 GRAY_BG, line=RGBColor(0xCC, 0xCC, 0xCC))
        add_text(s, x + Inches(1.45), y, Inches(2.05), col_h, v,
                 size=11, color=DARK_TXT,
                 anchor=MSO_ANCHOR.MIDDLE)

    # 저장 버튼 mockup
    btn_y = fy + fh - Inches(0.85)
    add_rect(s, fx + Inches(5.2), btn_y, Inches(2.05), Inches(0.55), GREEN)
    add_text(s, fx + Inches(5.2), btn_y, Inches(2.05), Inches(0.55),
             "💾  DB 저장",
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 우측 - 안내
    rx = Inches(8.2)
    add_text(s, rx, Inches(1.1), Inches(4.7), Inches(0.5),
             "▍ 확인/수정 포인트", size=16, bold=True, color=NAVY)

    tips = [
        ("상품코드 / 상품명",
         "파일명 또는 시트 내용에서 추출. 비어 있거나 다르면 직접 입력. "
         "[상품코드] 입력 시 자동으로 상품명을 조회합니다."),
        ("출발일 / 도착일",
         "다양한 날짜형식을 인식. 한 행에서만 못 찾으면 시트 전체를 다시 탐색."),
        ("대표자 · 연락처 · 이메일",
         "1) 시트 내 명시값 → 2) 첫 여행자 → 3) 빈 값 순으로 자동 입력."),
        ("총 예약금액",
         "자동 인식되지 않은 경우 수동 입력. 통화는 시스템 기본값 사용."),
        ("메모",
         "그룹 단위 특이사항을 자유롭게 추가 가능."),
    ]
    cy = Inches(1.65)
    for t, d in tips:
        add_rect(s, rx, cy, Inches(4.7), Inches(0.95),
                 LIGHT_BLUE, line=NAVY)
        add_text(s, rx + Inches(0.2), cy + Inches(0.05),
                 Inches(4.4), Inches(0.4),
                 t, size=12, bold=True, color=NAVY)
        add_text(s, rx + Inches(0.2), cy + Inches(0.4),
                 Inches(4.4), Inches(0.55),
                 d, size=10, color=DARK_TXT)
        cy = cy + Inches(1.05)

    add_footer(s)

# =========================================================
# 슬라이드 9 — 4단계 여행자 목록 편집
# =========================================================
def slide_step4(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "07. 4단계 · 여행자 목록 편집",
                   "추가 · 수정 · 삭제 · 복제", page=9, total=total)

    # 통계 mockup
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.85),
             NAVY)
    stats = [("총원", "24"), ("남성", "12"), ("여성", "12"), ("여권보유", "22")]
    for i, (k, v) in enumerate(stats):
        x = Inches(0.5) + Inches(3.075) * i
        add_text(s, x, Inches(1.15), Inches(3), Inches(0.35),
                 k, size=11, color=RGBColor(0xCB, 0xD3, 0xE8),
                 align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(1.4), Inches(3), Inches(0.5),
                 v, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # 테이블 mockup
    tx, ty = Inches(0.5), Inches(2.15)
    tw, th = Inches(8.6), Inches(4.65)
    add_rect(s, tx, ty, tw, th, WHITE,
             line=RGBColor(0xCC, 0xCC, 0xCC))
    cols = ["#", "한글성명", "영문성명", "여권번호", "성별", "객실", "작업"]
    widths = [0.4, 1.4, 1.6, 1.6, 0.7, 1.2, 1.7]
    cx = tx
    add_rect(s, tx, ty, tw, Inches(0.45), NAVY)
    for i, c in enumerate(cols):
        add_text(s, cx, ty, Inches(widths[i]), Inches(0.45),
                 c, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(widths[i])

    sample = [
        ["1", "홍길동", "HONG GILDONG", "M12345678", "남", "2r1p", "✏ 📋 🗑"],
        ["2", "김영희", "KIM YOUNGHEE", "M87654321", "여", "2r1p", "✏ 📋 🗑"],
        ["3", "이철수", "LEE CHEOLSU",  "M11112222", "남", "1r1p", "✏ 📋 🗑"],
        ["4", "박지민", "PARK JIMIN",   "M33334444", "여", "2r1p", "✏ 📋 🗑"],
    ]
    for ri, row in enumerate(sample):
        cx = tx
        rh = Inches(0.5)
        ry = ty + Inches(0.45) + rh * ri
        for i, v in enumerate(row):
            add_rect(s, cx, ry, Inches(widths[i]), rh,
                     WHITE if ri % 2 == 0 else GRAY_BG,
                     line=RGBColor(0xEE, 0xEE, 0xEE))
            add_text(s, cx, ry, Inches(widths[i]), rh,
                     v, size=10, color=DARK_TXT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += Inches(widths[i])

    # 우측 행동 가이드
    rx = Inches(9.3)
    add_text(s, rx, Inches(2.15), Inches(3.5), Inches(0.5),
             "▍ 행 단위 동작", size=16, bold=True, color=NAVY)
    actions = [
        ("✏", "수정",  NAVY,
         "모달이 열리고 11개 필드를 직접 편집"),
        ("📋", "복제",  GREEN,
         "동일 정보 한 명 더 추가 (+ \"(복사본)\")"),
        ("🗑", "삭제",  RED,
         "해당 여행자 한 명만 목록에서 제거"),
        ("➕", "추가",  ORANGE,
         "[여행자 추가] 버튼으로 새 행 입력"),
    ]
    cy = Inches(2.7)
    for icon, t, color, d in actions:
        add_rect(s, rx, cy, Inches(3.5), Inches(0.95),
                 WHITE, line=color)
        add_rect(s, rx, cy, Inches(0.85), Inches(0.95), color)
        add_text(s, rx, cy, Inches(0.85), Inches(0.95), icon,
                 size=22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, rx + Inches(0.95), cy + Inches(0.07),
                 Inches(2.4), Inches(0.35),
                 t, size=13, bold=True, color=color)
        add_text(s, rx + Inches(0.95), cy + Inches(0.4),
                 Inches(2.4), Inches(0.55),
                 d, size=10, color=GRAY_TXT)
        cy = cy + Inches(1.05)

    add_footer(s)

# =========================================================
# 슬라이드 10 — 5단계 일괄 편집 (예약번호별 설명)
# =========================================================
def slide_step5(total):
    """
    8번 단계를 예약번호별로 설명한다.
    - 시스템은 '대표예약자' 가 같은 여행자끼리 '예약그룹(=예약번호)' 단위로 자동 묶음
    - 각 예약그룹마다 배경색이 다르게 표시
    - 일괄편집은 (a) 같은 예약번호 내, (b) 여러 예약번호에 걸쳐 모두 가능
    - 선택분만 별도 새 예약번호로 분리도 가능
    """
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "08. 5단계 · 일괄(다중) 편집",
                   "예약번호(=예약그룹) 단위로 어떻게 동작하는지 살펴보세요",
                   page=10, total=total)

    # ── 핵심 개념 박스 ─────────────────────────────────────
    add_rect(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.95),
             LIGHT_BLUE, line=NAVY)
    add_text(s, Inches(0.7), Inches(1.13), Inches(12), Inches(0.4),
             "🧩 핵심 개념 — 예약번호는 \"대표예약자\" 단위로 자동 분리됩니다",
             size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.7), Inches(1.48), Inches(12), Inches(0.5),
             "같은 시트라도  대표예약자가 \"홍길동\" / \"김영희\" 로 나뉘면 → 저장 시 두 개의 예약번호로 분리됩니다.   "
             "여행자 목록에서 [예약그룹 1] · [예약그룹 2] … 헤더로 표시되며 행 배경색도 그룹마다 다릅니다.",
             size=11, color=DARK_TXT)

    # ── 좌측: 예약그룹 mockup (테이블 미리보기) ───────────
    tx, ty = Inches(0.5), Inches(2.15)
    tw     = Inches(7.5)
    add_text(s, tx, ty, tw, Inches(0.4),
             "▍ 화면에서 보이는 모습", size=14, bold=True, color=NAVY)

    # 표 본체
    table_y = ty + Inches(0.45)
    add_rect(s, tx, table_y, tw, Inches(4.5), WHITE,
             line=RGBColor(0xCC, 0xCC, 0xCC))

    # 헤더
    cols = ["☐", "#", "한글성명", "영문성명", "여권번호", "객실"]
    widths = [0.45, 0.45, 1.3, 1.6, 1.7, 2.0]
    cx = tx
    add_rect(s, tx, table_y, tw, Inches(0.4), NAVY)
    for i, c in enumerate(cols):
        add_text(s, cx, table_y, Inches(widths[i]), Inches(0.4),
                 c, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(widths[i])

    # 그룹 1 (연파랑) ─ 대표: 홍길동, 3명
    grp1_color = RGBColor(0xE8, 0xF4, 0xFD)
    grp1_y = table_y + Inches(0.4)
    add_rect(s, tx, grp1_y, tw, Inches(0.4), grp1_color)
    add_text(s, tx + Inches(0.1), grp1_y, tw - Inches(0.2), Inches(0.4),
             "👥  예약그룹 1   ─   대표예약자: 홍길동   (3명)",
             size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    g1_rows = [
        ["☐", "1", "홍길동", "HONG GILDONG", "M12345678", "2r1p"],
        ["☐", "2", "홍순이", "HONG SOONI",   "M12345679", "2r1p"],
        ["☐", "3", "홍철수", "HONG CHEOLSU", "M12345680", "2r1p"],
    ]
    for ri, row in enumerate(g1_rows):
        cx = tx
        ry = grp1_y + Inches(0.4) + Inches(0.38) * ri
        for i, v in enumerate(row):
            add_rect(s, cx, ry, Inches(widths[i]), Inches(0.38),
                     grp1_color, line=RGBColor(0xE0, 0xE0, 0xE0))
            add_text(s, cx, ry, Inches(widths[i]), Inches(0.38),
                     v, size=9, color=DARK_TXT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += Inches(widths[i])

    # 그룹 2 (연노랑) ─ 대표: 김영희, 2명
    grp2_color = RGBColor(0xFE, 0xF9, 0xE7)
    grp2_y = grp1_y + Inches(0.4) + Inches(0.38) * 3
    add_rect(s, tx, grp2_y, tw, Inches(0.4), grp2_color)
    add_text(s, tx + Inches(0.1), grp2_y, tw - Inches(0.2), Inches(0.4),
             "👥  예약그룹 2   ─   대표예약자: 김영희   (2명)",
             size=10, bold=True, color=ORANGE, anchor=MSO_ANCHOR.MIDDLE)
    g2_rows = [
        ["☐", "4", "김영희", "KIM YOUNGHEE", "M87654321", "2r1p"],
        ["☐", "5", "김미영", "KIM MIYOUNG",  "M87654322", "1r1p"],
    ]
    for ri, row in enumerate(g2_rows):
        cx = tx
        ry = grp2_y + Inches(0.4) + Inches(0.38) * ri
        for i, v in enumerate(row):
            add_rect(s, cx, ry, Inches(widths[i]), Inches(0.38),
                     grp2_color, line=RGBColor(0xE0, 0xE0, 0xE0))
            add_text(s, cx, ry, Inches(widths[i]), Inches(0.38),
                     v, size=9, color=DARK_TXT,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += Inches(widths[i])

    # 캡션
    add_text(s, tx, table_y + Inches(4.55), tw, Inches(0.35),
             "↑ 그룹마다 배경색이 자동으로 달라지고, 그룹 헤더에 \"예약그룹 N — 대표 / 인원수\" 가 표시됩니다.",
             size=9, color=GRAY_TXT)

    # ── 우측: 예약번호별 일괄편집 동작 ──────────────────
    rx = Inches(8.2)
    add_text(s, rx, Inches(2.15), Inches(4.7), Inches(0.4),
             "▍ 예약번호 단위로 본 일괄 편집",
             size=14, bold=True, color=NAVY)

    # 케이스 카드 4개
    cases = [
        ("CASE 1",
         "한 예약번호 안에서만",  NAVY,
         "예) 예약그룹 1(홍길동 외 2)만 선택  →  성별·객실 일괄 변경.",
         "동일 예약번호 내 인원만 한꺼번에 바뀝니다. 예약번호는 그대로 유지."),
        ("CASE 2",
         "여러 예약번호에 걸쳐",   GREEN,
         "예) 그룹1+그룹2 행을 함께 체크  →  객실타입 일괄 \"2r1p\".",
         "원래 속한 예약번호는 변하지 않고 값만 일괄 적용됩니다."),
        ("CASE 3",
         "예약번호별로 별도 저장", ORANGE,
         "예) 그룹1만 체크  →  [DB 저장] / 그룹2만 체크  →  [DB 저장].",
         "예약번호가 명확히 분리되어 저장 결과 메시지에도 따로 표시."),
        ("CASE 4",
         "선택 → 새 예약번호 분리", RGBColor(0x6F,0x42,0xC1),
         "예) 그룹1 중 일부만 체크  →  [선택 → 새 그룹] 클릭.",
         "선택분이 새로운 대표(=새 예약번호) 로 분리되어 저장됩니다."),
    ]
    cy = Inches(2.65)
    for tag, ttl, color, ex, eff in cases:
        add_rect(s, rx, cy, Inches(4.7), Inches(1.0),
                 WHITE, line=color)
        add_rect(s, rx, cy, Inches(0.95), Inches(1.0), color)
        add_text(s, rx, cy + Inches(0.07), Inches(0.95), Inches(0.35),
                 tag, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, rx, cy + Inches(0.45), Inches(0.95), Inches(0.5),
                 "✓", size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(s, rx + Inches(1.05), cy + Inches(0.05),
                 Inches(3.55), Inches(0.32),
                 ttl, size=12, bold=True, color=color)
        add_text(s, rx + Inches(1.05), cy + Inches(0.35),
                 Inches(3.55), Inches(0.32),
                 ex, size=9, color=DARK_TXT)
        add_text(s, rx + Inches(1.05), cy + Inches(0.66),
                 Inches(3.55), Inches(0.32),
                 "→ " + eff, size=9, color=GRAY_TXT)
        cy = cy + Inches(1.05)

    # ── 하단 팁 ─────────────────────────────────────────
    add_rect(s, Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.18),
             NAVY)

    add_footer(s)

# =========================================================
# 슬라이드 11 — 6단계 DB 저장
# =========================================================
def slide_step6(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "09. 6단계 · DB 저장",
                   "시트별 저장 · 결과 확인", page=11, total=total)

    # 좌: 저장 절차
    add_text(s, Inches(0.5), Inches(1.1), Inches(7), Inches(0.5),
             "▍ 저장 방법", size=18, bold=True, color=NAVY)

    steps = [
        ("1", "시트별 [DB 저장] 클릭",
         "각 시트(아코디언) 우측의 초록색 [DB 저장] 버튼을 누릅니다."),
        ("2", "확인 메시지",
         "\"이 시트의 그룹 예약을 저장하시겠습니까?\" → [확인]."),
        ("3", "처리 결과 표시",
         "성공 시 → 초록 박스로 그룹예약번호 안내.\n실패 시 → 빨간 박스로 사유 안내."),
        ("4", "선택분 별도 저장",
         "여행자 목록에서 일부만 체크 후\n[선택 저장]을 누르면 새 예약번호로 저장."),
    ]
    cy = Inches(1.7)
    for i, (n, t, d) in enumerate(steps):
        y = cy + Inches(1.15) * i
        add_rect(s, Inches(0.5), y, Inches(0.9), Inches(0.95), GREEN)
        add_text(s, Inches(0.5), y, Inches(0.9), Inches(0.95), n,
                 size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.45), y, Inches(5.6), Inches(0.95),
                 WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_text(s, Inches(1.6), y + Inches(0.07), Inches(5.4), Inches(0.4),
                 t, size=13, bold=True, color=DARK_TXT)
        add_text(s, Inches(1.6), y + Inches(0.45), Inches(5.4), Inches(0.55),
                 d, size=10, color=GRAY_TXT)

    # 우: 결과 박스 mockup + 주의
    rx = Inches(7.5)
    add_text(s, rx, Inches(1.1), Inches(5.3), Inches(0.5),
             "▍ 저장 결과", size=18, bold=True, color=NAVY)

    add_rect(s, rx, Inches(1.7), Inches(5.3), Inches(1.5),
             LIGHT_GREEN, line=GREEN)
    add_text(s, rx + Inches(0.2), Inches(1.78), Inches(5), Inches(0.4),
             "✅ 성공 메시지 예시", size=12, bold=True, color=GREEN)
    add_text(s, rx + Inches(0.2), Inches(2.15), Inches(5), Inches(1.0),
             "그룹 예약이 저장되었습니다.\n  • 예약번호: GR-2026-0428-001\n  • 등록 인원: 24명",
             size=11, color=DARK_TXT)

    add_rect(s, rx, Inches(3.3), Inches(5.3), Inches(1.5),
             RGBColor(0xFD, 0xEC, 0xEA), line=RED)
    add_text(s, rx + Inches(0.2), Inches(3.38), Inches(5), Inches(0.4),
             "❌ 실패 메시지 예시", size=12, bold=True, color=RED)
    add_text(s, rx + Inches(0.2), Inches(3.75), Inches(5), Inches(1.0),
             "필수 항목이 누락되었습니다.\n  • 한글성명, 출발일, 상품명/코드 확인 후 다시 시도",
             size=11, color=DARK_TXT)

    add_rect(s, rx, Inches(4.9), Inches(5.3), Inches(1.95),
             LIGHT_YEL, line=ORANGE)
    add_text(s, rx + Inches(0.2), Inches(4.98), Inches(5), Inches(0.4),
             "⚠ 저장 시 유의사항", size=12, bold=True, color=ORANGE)
    bullet(s, rx + Inches(0.2), Inches(5.4), Inches(5), Inches(1.4),
           [
            "동일 시트를 두 번 저장하면 중복 등록될 수 있음",
            "대표예약자 단위로 묶어서 저장됨",
            "저장된 예약은 협력사 ID 기준으로 귀속",
           ], size=11, bullet_color=ORANGE, gap=0.4)

    add_footer(s)

# =========================================================
# 슬라이드 12 — FAQ / 오류 대응
# =========================================================
def slide_faq(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "10. 자주 묻는 질문 · 오류 대응",
                   "FAQ & Troubleshooting", page=12, total=total)

    items = [
        ("Q1.", "업로드 후 \"파일을 처리할 수 없습니다\" 오류가 떠요.",
         "→ 파일 확장자 확인 (.xlsx / .xls / .csv 만 가능). 시트가 비어 있거나 암호가 걸려 있으면 처리 불가."),
        ("Q2.", "여행자가 한 명도 인식되지 않아요.",
         "→ 첫 행에 컬럼명(한글성명/영문성명 등)이 있어야 합니다. 시트가 합계/요약 행만 있는 경우 인식이 안 됩니다."),
        ("Q3.", "출발일이 비어 있거나 이상하게 표시돼요.",
         "→ 자동 인식이 실패한 경우. 시트별 [출발일/도착일] 입력란에서 직접 수정 후 [DB 저장]."),
        ("Q4.", "상품코드를 입력해도 상품명이 안 나와요.",
         "→ 등록되지 않은 상품코드일 수 있습니다. 담당자에게 정확한 코드 확인 요청."),
        ("Q5.", "같은 여행자가 두 번 보입니다.",
         "→ 자동 중복제거 기능이 있지만 이름/여권번호가 다르면 다른 사람으로 인식됩니다. 직접 [삭제] 또는 [수정]."),
        ("Q6.", "선택 저장과 DB 저장의 차이는?",
         "→ [DB 저장] = 시트 전체 / [선택 저장] = 체크된 여행자만 별도 새 그룹예약으로 저장."),
    ]
    cy = Inches(1.15)
    for q, t, a in items:
        add_rect(s, Inches(0.5), cy, Inches(12.3), Inches(0.85),
                 WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, Inches(0.5), cy, Inches(0.9), Inches(0.85), NAVY)
        add_text(s, Inches(0.5), cy, Inches(0.9), Inches(0.85),
                 q, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(1.55), cy + Inches(0.05),
                 Inches(11), Inches(0.38),
                 t, size=12, bold=True, color=DARK_TXT)
        add_text(s, Inches(1.55), cy + Inches(0.42),
                 Inches(11), Inches(0.4),
                 a, size=10, color=GRAY_TXT)
        cy = cy + Inches(0.93)

    add_footer(s)

# =========================================================
# 슬라이드 13 — 마무리
# =========================================================
def slide_thanks(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    deco = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(-2), Inches(-2),
                              Inches(6), Inches(6))
    deco.fill.solid(); deco.fill.fore_color.rgb = NAVY_DARK
    deco.line.fill.background()
    deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(9.5), Inches(4),
                               Inches(6), Inches(6))
    deco2.fill.solid(); deco2.fill.fore_color.rgb = NAVY_DARK
    deco2.line.fill.background()

    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1),
             "Thank You", size=66, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.6),
             "수고하셨습니다.  스마트 그룹 예약 등록을 즐겨보세요!",
             size=22, color=RGBColor(0xCB, 0xD3, 0xE8))

    add_rect(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.7),
             WHITE)
    add_text(s, Inches(1.0), Inches(4.6), Inches(11), Inches(0.5),
             "📞  도움이 필요하시면", size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(5.05), Inches(11), Inches(1.1),
             [
              "• 운영팀 / 담당자에게 문의 주세요.",
              "• 오류 메시지가 있다면 화면 캡처와 사용 파일을 함께 보내 주시면 빠르게 도와드릴 수 있습니다.",
             ], size=13, color=DARK_TXT)

    add_text(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
             "푸른투어 PARTNER PORTAL · input_batch.php",
             size=11, color=RGBColor(0xCB, 0xD3, 0xE8), align=PP_ALIGN.CENTER)

# ===== 빌드 =====
TOTAL = 13
slide_toc(TOTAL)
slide_overview(TOTAL)
slide_login(TOTAL)
slide_layout(TOTAL)
slide_step1(TOTAL)
slide_step2(TOTAL)
slide_step3(TOTAL)
slide_step4(TOTAL)
slide_step5(TOTAL)
slide_step6(TOTAL)
slide_faq(TOTAL)
slide_thanks(TOTAL)

import os, time
OUT = r"d:/www/prttour_myprt/input_batch_user_manual.pptx"
try:
    prs.save(OUT)
    print("SAVED:", OUT)
except PermissionError:
    alt = OUT.replace(".pptx", f"_v{int(time.time())}.pptx")
    prs.save(alt)
    print("SAVED (locked, used new name):", alt)
