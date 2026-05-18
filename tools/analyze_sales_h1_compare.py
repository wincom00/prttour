# -*- coding: utf-8 -*-
"""
2025 vs 2026 H1(1~4월, 판매일 기준) 매출 비교 분석 PPTX 생성기.
prod_month.php 의 쿼리 로직을 그대로 사용한다.
"""
import os, sys, io, datetime, math
import pymysql
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ----- CONFIG -----
DB = dict(
    host="98.91.65.48", port=3306, user="wincom00",
    password="Lee10011!", database="prtadmindb", charset="utf8mb4",
)
YEARS = [2025, 2026]
MONTHS = [1, 2, 3, 4]                     # H1 현재월(2026-04-30)까지
PROJECT_ROOT = r"d:\www\prttour_myprt"
TOOLS_DIR = os.path.join(PROJECT_ROOT, "tools")
CHARTS_DIR = os.path.join(TOOLS_DIR, "_charts")
OUT_PPTX = os.path.join(PROJECT_ROOT, "sales_h1_compare_2025_2026.pptx")
RAW_CSV = os.path.join(TOOLS_DIR, "_h1_compare_raw.csv")

NAVY = RGBColor(0x13, 0x11, 0x76)
SKY = RGBColor(0x39, 0x93, 0xBA)
GRAY = RGBColor(0x66, 0x66, 0x66)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)

NAVY_HEX = "#131176"
SKY_HEX = "#3993ba"
GREEN_HEX = "#27ae60"
RED_HEX = "#c0392b"

os.makedirs(CHARTS_DIR, exist_ok=True)

# 한글 폰트
rcParams["font.family"] = "Malgun Gothic"
rcParams["axes.unicode_minus"] = False


# ----- DATA -----
def fetch_period(year):
    sql = """
    SELECT a.p_day, b.p_code, b.p_name,
           MONTH(b.wdate) AS m,
           SUM(b.last_total) AS revenue,
           SUM(b.p_cnt)      AS headcount
    FROM   product_master a
    INNER JOIN reserve_info b ON a.p_code = b.p_code
    WHERE  b.rev_status = 'DONE'
      AND  a.p_code NOT LIKE %s
      AND  a.p_code NOT LIKE %s
      AND  b.parent = 'MAIN'
      AND  YEAR(b.wdate)  = %s
      AND  MONTH(b.wdate) BETWEEN %s AND %s
    GROUP BY a.p_day, b.p_code, b.p_name, m
    """
    conn = pymysql.connect(**DB, connect_timeout=15)
    try:
        df = pd.read_sql(
            sql, conn,
            params=("%PICKUP%", "%SENDING%", year, MONTHS[0], MONTHS[-1]),
        )
    finally:
        conn.close()
    df["year"] = year
    df["revenue"] = df["revenue"].astype(float)
    df["headcount"] = df["headcount"].astype(float)
    return df


def build_compare(df25, df26):
    # 월별 집계
    monthly = pd.DataFrame({"m": MONTHS})
    for y, df in [(2025, df25), (2026, df26)]:
        g = df.groupby("m", as_index=False).agg(rev=("revenue", "sum"),
                                                 hc=("headcount", "sum"))
        g = g.rename(columns={"rev": f"rev_{y}", "hc": f"hc_{y}"})
        monthly = monthly.merge(g, on="m", how="left")
    monthly = monthly.fillna(0.0)

    # KPI
    tot_rev = {y: float(df["revenue"].sum()) for y, df in [(2025, df25), (2026, df26)]}
    tot_hc  = {y: float(df["headcount"].sum()) for y, df in [(2025, df25), (2026, df26)]}
    arpu    = {y: (tot_rev[y] / tot_hc[y]) if tot_hc[y] else 0.0 for y in YEARS}
    n_prod  = {y: df["p_code"].nunique() for y, df in [(2025, df25), (2026, df26)]}

    def yoy(a, b):
        if b == 0: return None
        return (a - b) / b * 100.0

    kpi = dict(
        rev_25=tot_rev[2025], rev_26=tot_rev[2026], rev_yoy=yoy(tot_rev[2026], tot_rev[2025]),
        hc_25=tot_hc[2025], hc_26=tot_hc[2026], hc_yoy=yoy(tot_hc[2026], tot_hc[2025]),
        arpu_25=arpu[2025], arpu_26=arpu[2026], arpu_yoy=yoy(arpu[2026], arpu[2025]),
        np_25=n_prod[2025], np_26=n_prod[2026],
    )

    # 상품 단위 비교 (양년도 outer merge)
    p25 = df25.groupby(["p_code", "p_name", "p_day"], as_index=False).agg(
        rev_25=("revenue", "sum"), hc_25=("headcount", "sum"))
    p26 = df26.groupby(["p_code", "p_name", "p_day"], as_index=False).agg(
        rev_26=("revenue", "sum"), hc_26=("headcount", "sum"))
    prods = p25.merge(p26, on=["p_code", "p_name", "p_day"], how="outer").fillna(0.0)
    prods["rev_diff"] = prods["rev_26"] - prods["rev_25"]
    prods["rev_pct"] = prods.apply(
        lambda r: (r["rev_diff"] / r["rev_25"] * 100.0) if r["rev_25"] else (100.0 if r["rev_26"] > 0 else 0.0),
        axis=1,
    )
    prods["status"] = prods.apply(
        lambda r: "신규(올해만)" if r["rev_25"] == 0 and r["rev_26"] > 0
        else ("단종(작년만)" if r["rev_26"] == 0 and r["rev_25"] > 0
              else "동시운영"),
        axis=1,
    )

    # p_day 그룹화: 당일/1박/2~3박/4박+
    def bucket(d):
        d = int(d) if pd.notna(d) else 0
        if d <= 1: return "당일"
        if d == 2: return "1박"
        if d <= 4: return "2~3박"
        return "4박+"
    prods["pday_bucket"] = prods["p_day"].apply(bucket)

    return monthly, kpi, prods


# ----- SIGNALS -----
def detect_signals(monthly, kpi, prods):
    # 외부요인 시그널
    common = prods[(prods["status"] == "동시운영") & (prods["rev_25"] > 0)]
    if len(common) > 0:
        same_dir = ((common["rev_diff"] > 0).sum(), (common["rev_diff"] < 0).sum())
        majority = max(same_dir) / max(1, len(common)) * 100.0
        majority_dir = "증가" if same_dir[0] >= same_dir[1] else "감소"
    else:
        majority, majority_dir = 0.0, "-"

    # 월별 변동률 분산도
    monthly["rev_diff"] = monthly["rev_2026"] - monthly["rev_2025"]
    monthly["rev_pct"] = monthly.apply(
        lambda r: (r["rev_diff"] / r["rev_2025"] * 100.0) if r["rev_2025"] else 0.0,
        axis=1,
    )
    season_std = float(monthly["rev_pct"].std() or 0.0)
    season_max_month = int(monthly.loc[monthly["rev_diff"].abs().idxmax(), "m"]) if len(monthly) else 0

    external = dict(
        common_count=len(common),
        same_dir_pct=majority,
        same_dir_label=majority_dir,
        season_std=season_std,
        season_max_month=season_max_month,
    )

    # 내부요인 시그널
    prods_ranked = prods.assign(abs_diff=prods["rev_diff"].abs()).sort_values("abs_diff", ascending=False)
    top5 = prods_ranked.head(5)
    total_abs = float(prods["rev_diff"].abs().sum() or 0.0)
    top5_share = float(top5["abs_diff"].sum()) / total_abs * 100.0 if total_abs else 0.0

    new_prods = prods[prods["status"] == "신규(올해만)"]
    dropped = prods[prods["status"] == "단종(작년만)"]

    # p_day 믹스 (매출 기준)
    bucket_25 = prods.groupby("pday_bucket")["rev_25"].sum()
    bucket_26 = prods.groupby("pday_bucket")["rev_26"].sum()
    mix_25 = (bucket_25 / bucket_25.sum() * 100.0) if bucket_25.sum() else bucket_25 * 0
    mix_26 = (bucket_26 / bucket_26.sum() * 100.0) if bucket_26.sum() else bucket_26 * 0

    internal = dict(
        top5_share=top5_share,
        top5=top5[["p_name", "rev_25", "rev_26", "rev_diff", "rev_pct"]].to_dict("records"),
        new_count=len(new_prods),
        new_rev_26=float(new_prods["rev_26"].sum()),
        dropped_count=len(dropped),
        dropped_rev_25=float(dropped["rev_25"].sum()),
        mix_25=mix_25.to_dict(),
        mix_26=mix_26.to_dict(),
        arpu_change=kpi["arpu_yoy"],
    )
    return external, internal


# ----- CHARTS -----
def chart_monthly_bar(monthly, value_prefix, ylabel, fname, fmt="{:,.0f}"):
    fig, ax = plt.subplots(figsize=(9, 5))
    x = list(range(len(monthly)))
    w = 0.38
    v25 = monthly[f"{value_prefix}_2025"].values
    v26 = monthly[f"{value_prefix}_2026"].values
    b1 = ax.bar([i - w/2 for i in x], v25, w, label="2025", color=NAVY_HEX)
    b2 = ax.bar([i + w/2 for i in x], v26, w, label="2026", color=SKY_HEX)
    ax.set_xticks(x)
    ax.set_xticklabels([f"{int(m)}월" for m in monthly["m"]])
    ax.set_ylabel(ylabel)
    ax.legend()
    ax.grid(axis="y", linestyle="--", alpha=0.4)
    # 증감률 라벨
    for i, (a, b) in enumerate(zip(v25, v26)):
        if a > 0:
            pct = (b - a) / a * 100.0
            color = GREEN_HEX if pct >= 0 else RED_HEX
            ax.text(i, max(a, b) * 1.02, f"{pct:+.1f}%",
                    ha="center", va="bottom", color=color, fontsize=10, fontweight="bold")
    fig.tight_layout()
    out = os.path.join(CHARTS_DIR, fname)
    fig.savefig(out, dpi=150)
    plt.close(fig)
    return out


def chart_arpu(monthly, fname):
    fig, ax = plt.subplots(figsize=(9, 5))
    arpu25 = [(r / h) if h else 0 for r, h in zip(monthly["rev_2025"], monthly["hc_2025"])]
    arpu26 = [(r / h) if h else 0 for r, h in zip(monthly["rev_2026"], monthly["hc_2026"])]
    ax.plot(monthly["m"], arpu25, marker="o", label="2025", color=NAVY_HEX, linewidth=2)
    ax.plot(monthly["m"], arpu26, marker="s", label="2026", color=SKY_HEX, linewidth=2)
    ax.set_xticks(monthly["m"])
    ax.set_xticklabels([f"{int(m)}월" for m in monthly["m"]])
    ax.set_ylabel("객단가 (USD)")
    ax.legend()
    ax.grid(linestyle="--", alpha=0.4)
    fig.tight_layout()
    out = os.path.join(CHARTS_DIR, fname)
    fig.savefig(out, dpi=150)
    plt.close(fig)
    return out


def chart_top_products(prods, mode, n, fname):
    sub = prods[prods["status"] == "동시운영"].copy()
    if mode == "up":
        sub = sub.sort_values("rev_diff", ascending=False).head(n)
        title_color = GREEN_HEX
    else:
        sub = sub.sort_values("rev_diff", ascending=True).head(n)
        title_color = RED_HEX
    sub = sub.iloc[::-1]  # bar 차트 위쪽 표시
    names = [str(s)[:30] for s in sub["p_name"].values]
    fig, ax = plt.subplots(figsize=(11, 6))
    y = list(range(len(sub)))
    h = 0.4
    ax.barh([i + h/2 for i in y], sub["rev_25"].values, h, label="2025", color=NAVY_HEX)
    ax.barh([i - h/2 for i in y], sub["rev_26"].values, h, label="2026", color=SKY_HEX)
    ax.set_yticks(y)
    ax.set_yticklabels(names)
    ax.set_xlabel("매출 (USD)")
    ax.legend(loc="lower right")
    for i, d in enumerate(sub["rev_diff"].values):
        ax.text(max(sub["rev_25"].iloc[i], sub["rev_26"].iloc[i]) * 1.02, i,
                f"{d:+,.0f}", va="center", color=title_color, fontsize=9, fontweight="bold")
    ax.grid(axis="x", linestyle="--", alpha=0.4)
    fig.tight_layout()
    out = os.path.join(CHARTS_DIR, fname)
    fig.savefig(out, dpi=150)
    plt.close(fig)
    return out


def chart_pday_mix(prods, fname):
    order = ["당일", "1박", "2~3박", "4박+"]
    b25 = prods.groupby("pday_bucket")["rev_25"].sum().reindex(order).fillna(0.0)
    b26 = prods.groupby("pday_bucket")["rev_26"].sum().reindex(order).fillna(0.0)
    fig, axes = plt.subplots(1, 2, figsize=(11, 5))
    colors = [NAVY_HEX, SKY_HEX, "#7fb069", "#f4a261"]
    for ax, vals, title in [(axes[0], b25, "2025"), (axes[1], b26, "2026")]:
        if vals.sum() > 0:
            ax.pie(vals.values, labels=[f"{n}\n${v:,.0f}" for n, v in zip(order, vals.values)],
                   colors=colors, autopct="%1.1f%%", startangle=90, wedgeprops=dict(width=0.45))
        ax.set_title(title, fontweight="bold")
    fig.tight_layout()
    out = os.path.join(CHARTS_DIR, fname)
    fig.savefig(out, dpi=150)
    plt.close(fig)
    return out


# ----- PPTX -----
def add_title_bar(slide, prs, text):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.7)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.runs[0].font.size = Pt(22)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_text_box(slide, left, top, width, height, text, size=14, bold=False, color=None, align=None):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if align is not None:
            p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            if color is not None:
                run.font.color.rgb = color
    return tx


def fmt_money(v):
    return f"${v:,.0f}"


def fmt_int(v):
    return f"{v:,.0f}"


def fmt_pct(v):
    if v is None or (isinstance(v, float) and (math.isnan(v) or math.isinf(v))):
        return "N/A"
    return f"{v:+.1f}%"


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_text_box(s, Inches(0.7), Inches(2.2), Inches(12), Inches(1.0),
                 "푸른투어", size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(s, Inches(0.7), Inches(2.8), Inches(12), Inches(1.4),
                 "2025 vs 2026 상반기 매출 비교 분석", size=44, bold=True,
                 color=RGBColor(0xFF, 0xFF, 0xFF))
    add_text_box(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.8),
                 "기준: 판매일(wdate)  /  비교 기간: 양 년도 1~4월", size=18,
                 color=RGBColor(0xC8, 0xCD, 0xE8))
    add_text_box(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
                 f"출력일: {datetime.date.today().isoformat()}", size=12,
                 color=RGBColor(0xC8, 0xCD, 0xE8))


def slide_summary(prs, kpi):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, prs, "Executive Summary  —  KPI 한눈에")
    cards = [
        ("총 매출 (USD)", fmt_money(kpi["rev_26"]), f"vs 2025 {fmt_money(kpi['rev_25'])}", kpi["rev_yoy"]),
        ("총 인원수 (명)", fmt_int(kpi["hc_26"]), f"vs 2025 {fmt_int(kpi['hc_25'])}", kpi["hc_yoy"]),
        ("객단가 (USD)", fmt_money(kpi["arpu_26"]), f"vs 2025 {fmt_money(kpi['arpu_25'])}", kpi["arpu_yoy"]),
        ("활성 상품 수", f"{kpi['np_26']}", f"vs 2025 {kpi['np_25']}",
         (kpi["np_26"] - kpi["np_25"]) / max(1, kpi["np_25"]) * 100.0),
    ]
    left0 = Inches(0.5); top = Inches(1.4); w = Inches(3.05); h = Inches(2.6); gap = Inches(0.1)
    for i, (label, val, sub, yoy) in enumerate(cards):
        left = Emu(left0 + (w + gap) * i)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = NAVY
        add_text_box(s, left, Emu(top + Inches(0.2)), w, Inches(0.5),
                     label, size=14, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        add_text_box(s, left, Emu(top + Inches(0.7)), w, Inches(0.9),
                     val, size=26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(s, left, Emu(top + Inches(1.6)), w, Inches(0.4),
                     sub, size=11, color=GRAY, align=PP_ALIGN.CENTER)
        yoy_color = GREEN if (yoy or 0) >= 0 else RED
        if yoy is None:
            yoy_text = "변화 없음"
        else:
            yoy_text = f"작년 대비 {abs(yoy):.1f}% 증가" if yoy >= 0 else f"작년 대비 {abs(yoy):.1f}% 감소"
        add_text_box(s, left, Emu(top + Inches(2.0)), w, Inches(0.5),
                     yoy_text, size=14, bold=True, color=yoy_color, align=PP_ALIGN.CENTER)

    add_text_box(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.5),
                 "비교 기간: 2025년 1~4월 과 2026년 1~4월 (판매일 기준)\n"
                 "포함 데이터: 예약 상태가 확정(DONE), 메인 예약, 픽업과 샌딩 상품은 제외.\n"
                 "통화: 시스템 기록 금액(달러 기준). 객단가 = 총매출 ÷ 총인원수.",
                 size=12, color=GRAY)


def slide_chart(prs, title, chart_path, footer=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, prs, title)
    s.shapes.add_picture(chart_path, Inches(0.5), Inches(1.0),
                         width=Inches(12.3), height=Inches(5.8))
    if footer:
        add_text_box(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
                     footer, size=11, color=GRAY)


def slide_external(prs, ext, monthly):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, prs, "외부요인 분석  —  데이터로 본 시그널과 가설")
    # 좌측 데이터 시그널
    left_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.0),
                                   Inches(6.2), Inches(6.0))
    left_box.fill.solid(); left_box.fill.fore_color.rgb = LIGHT
    left_box.line.color.rgb = NAVY
    add_text_box(s, Inches(0.6), Inches(1.1), Inches(6.0), Inches(0.5),
                 "데이터로 본 시그널", size=16, bold=True, color=NAVY)
    if ext["same_dir_label"] == "감소":
        same_dir_text = f"매출이 감소한 상품이 약 {ext['same_dir_pct']:.0f}%"
    else:
        same_dir_text = f"매출이 증가한 상품이 약 {ext['same_dir_pct']:.0f}%"
    if ext["same_dir_pct"] >= 70:
        verdict = "대다수 상품이 같은 방향이라 외부 환경 영향이 클 가능성이 높습니다."
    else:
        verdict = "상품마다 방향이 갈려 시장 전반보다 상품별 사정이 더 큽니다."
    lines = [
        f"두 해 모두 판매된 상품 {ext['common_count']}개 중,",
        f"올해 {same_dir_text} 입니다.",
        verdict,
        "",
        f"월별로 매출 증감폭이 들쭉날쭉합니다. 가장 크게 움직인 달은 {ext['season_max_month']}월입니다.",
        "",
        "월별 매출 증감 (2026 기준, 2025 대비):",
    ]
    add_text_box(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(2.6),
                 "\n".join(lines), size=13, color=GRAY)
    # 월별 표 — 증가/감소 자연어
    monthly_pct_lines = []
    for _, r in monthly.iterrows():
        pct = r["rev_pct"]
        diff = r["rev_diff"]
        if diff >= 0:
            monthly_pct_lines.append(f"   {int(r['m'])}월 : {abs(pct):.1f}% 증가 ($ {abs(diff):,.0f} 늘어남)")
        else:
            monthly_pct_lines.append(f"   {int(r['m'])}월 : {abs(pct):.1f}% 감소 ($ {abs(diff):,.0f} 줄어듦)")
    add_text_box(s, Inches(0.6), Inches(4.4), Inches(6.0), Inches(2.0),
                 "\n".join(monthly_pct_lines), size=12, color=NAVY)

    # 우측 가설(검증 필요)
    right_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.0),
                                    Inches(6.2), Inches(6.0))
    right_box.fill.solid(); right_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFB, 0xEC)
    right_box.line.color.rgb = RGBColor(0xCB, 0x9B, 0x2A)
    add_text_box(s, Inches(7.0), Inches(1.1), Inches(6.0), Inches(0.5),
                 "원인 가설 (별도 확인이 필요한 항목)", size=16, bold=True, color=RGBColor(0x88, 0x6A, 0x16))
    hyp = [
        "환율 (원/달러): 달러가 비싸지면 미주 패키지 가격 부담이 커져 예약이 줄어듭니다.",
        "거시 경기와 소비심리: 가계 여유가 줄면 해외여행 지출이 우선 줄어듭니다.",
        "미주 여행 수요 자체의 흐름: 미국 인기 회복 또는 둔화 시기 영향.",
        "항공 운임과 좌석 공급: 항공권이 비싸지거나 좌석이 줄면 패키지 가격이 따라 오릅니다.",
        "주요 이벤트와 정책: 올림픽, 박람회, 비자나 입국 정책 변화가 수요를 바꿉니다.",
        "경쟁사 프로모션과 가격: 경쟁사가 싸게 팔면 우리 매출이 줄어들 수 있습니다.",
        "날씨와 자연재해: 산불, 한파, 폭우 등 출발지 또는 도착지 이슈.",
    ]
    add_text_box(s, Inches(7.0), Inches(1.6), Inches(6.0), Inches(5.0),
                 "\n".join(hyp), size=13, color=GRAY)
    add_text_box(s, Inches(7.0), Inches(6.5), Inches(6.0), Inches(0.5),
                 "환율, 항공 운임 같은 외부 수치를 함께 보면 원인 추정이 더 정확해집니다.",
                 size=11, color=GRAY)


def slide_internal(prs, intl):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, prs, "내부요인 분석  —  데이터로 본 시그널과 가설")
    # 좌측: 데이터 시그널
    left_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.0),
                                   Inches(6.2), Inches(6.0))
    left_box.fill.solid(); left_box.fill.fore_color.rgb = LIGHT
    left_box.line.color.rgb = NAVY
    add_text_box(s, Inches(0.6), Inches(1.1), Inches(6.0), Inches(0.5),
                 "데이터로 본 시그널", size=16, bold=True, color=NAVY)
    if intl["top5_share"] >= 60:
        verdict = "특정 상위 5개 상품이 매출 흐름을 좌우하고 있어 내부 요인(상품 전략)의 영향이 큽니다."
    else:
        verdict = "변동이 여러 상품에 흩어져 있어 우리 상품 전략보다 시장 환경의 영향도 같이 의심됩니다."
    arpu_pct = intl["arpu_change"] or 0
    if arpu_pct >= 0:
        arpu_text = f"객단가는 작년보다 {abs(arpu_pct):.1f}% 증가했습니다."
    else:
        arpu_text = f"객단가는 작년보다 {abs(arpu_pct):.1f}% 감소했습니다."
    lines = [
        f"매출이 가장 많이 변한 5개 상품이 전체 변동의 약 {intl['top5_share']:.0f}%를 설명합니다.",
        verdict,
        "",
        f"올해 새로 판 상품: {intl['new_count']}개, 매출 $ {intl['new_rev_26']:,.0f} 발생.",
        f"올해는 안 판 상품(단종): {intl['dropped_count']}개, 작년 매출 $ {intl['dropped_rev_25']:,.0f}이 사라짐.",
        "",
        arpu_text,
        "",
        "매출이 가장 많이 변한 상품 TOP 5",
    ]
    add_text_box(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(2.8),
                 "\n".join(lines), size=13, color=GRAY)

    rows = []
    for r in intl["top5"]:
        name = (r["p_name"] or "")[:34]
        diff = r["rev_diff"]
        pct = r["rev_pct"] or 0
        if diff >= 0:
            rows.append(f"   {name}  →  $ {abs(diff):,.0f} 증가 ({abs(pct):.1f}% 증가)")
        else:
            rows.append(f"   {name}  →  $ {abs(diff):,.0f} 감소 ({abs(pct):.1f}% 감소)")
    add_text_box(s, Inches(0.6), Inches(4.6), Inches(6.0), Inches(2.3),
                 "\n".join(rows), size=11, color=NAVY)

    # 우측: 가설
    right_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.0),
                                    Inches(6.2), Inches(6.0))
    right_box.fill.solid(); right_box.fill.fore_color.rgb = RGBColor(0xEC, 0xF7, 0xFF)
    right_box.line.color.rgb = SKY
    add_text_box(s, Inches(7.0), Inches(1.1), Inches(6.0), Inches(0.5),
                 "원인 가설 (별도 확인이 필요한 항목)", size=16, bold=True, color=NAVY)
    hyp = [
        "상품 라인업 변경: 신상품을 새로 넣고 잘 안 팔리던 상품을 뺀 영향.",
        "가격과 할인 정책 변경: 가격을 올리거나 내리면 객단가가 따라 움직입니다.",
        "판매 채널 변화: 직판, 대리점, OTA 등 어디서 더 많이 팔렸는지의 변화.",
        "마케팅 예산과 시즌 캠페인 강도: 광고를 많이 한 상품이 더 잘 팔립니다.",
        "운영 가용성: 가이드, 버스, 호텔 확보 가능 여부에 따라 판매량이 달라집니다.",
        "영업과 고객센터 인력 변동: 응대 속도와 품질이 예약 전환에 영향을 줍니다.",
        "예약 화면과 시스템 사용성 개선 효과.",
    ]
    add_text_box(s, Inches(7.0), Inches(1.6), Inches(6.0), Inches(5.0),
                 "\n".join(hyp), size=13, color=GRAY)
    mix_text = ", ".join(
        [f"{k} {intl['mix_25'].get(k,0):.0f}% → {intl['mix_26'].get(k,0):.0f}%"
         for k in ["당일", "1박", "2~3박", "4박+"]]
    )
    add_text_box(s, Inches(7.0), Inches(6.5), Inches(6.0), Inches(0.5),
                 f"투어 일수별 매출 비중 (2025 → 2026): {mix_text}",
                 size=10, color=GRAY)


def slide_conclusion(prs, kpi, ext, intl):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, prs, "종합 결론과 권고")

    def updown(v):
        if v is None:
            return "변화 없음"
        return f"{abs(v):.1f}% 증가" if v >= 0 else f"{abs(v):.1f}% 감소"

    bullets = []
    bullets.append(f"1. 매출은 작년보다 {updown(kpi['rev_yoy'])}, 인원은 {updown(kpi['hc_yoy'])}, "
                   f"객단가는 {updown(kpi['arpu_yoy'])} 했습니다.")
    if ext["same_dir_pct"] >= 70:
        bullets.append(f"2. 두 해 모두 운영된 상품 중 약 {ext['same_dir_pct']:.0f}%가 같은 방향으로 움직였습니다. "
                       "외부 환경의 영향이 클 가능성이 높아 환율과 수요 지표를 함께 보는 것이 좋습니다.")
    else:
        bullets.append(f"2. 두 해 모두 운영된 상품 중 같은 방향으로 움직인 비율이 약 {ext['same_dir_pct']:.0f}% 입니다. "
                       "상품마다 사정이 갈려 있어 상품 단위 점검이 필요합니다.")
    if intl["top5_share"] >= 60:
        bullets.append(f"3. 매출 변동 상위 5개 상품이 전체 변동의 {intl['top5_share']:.0f}%를 차지합니다. "
                       "특정 상품 집중도가 높아 우리 내부 전략의 영향이 큽니다.")
    else:
        bullets.append(f"3. 매출 변동이 여러 상품에 흩어져 있고 상위 5개의 비중이 {intl['top5_share']:.0f}%에 그칩니다. "
                       "특정 상품 집중도가 낮아 외부 환경의 영향이 함께 의심됩니다.")
    if intl["new_count"] > 0:
        bullets.append(f"4. 올해 새로 판매를 시작한 상품 {intl['new_count']}개가 $ {intl['new_rev_26']:,.0f} 의 매출을 "
                       "추가로 만들었습니다. 신규 상품의 수익성과 광고비 회수율을 점검할 시점입니다.")
    if intl["dropped_count"] > 0:
        bullets.append(f"5. 작년에는 팔렸으나 올해는 빠진 상품 {intl['dropped_count']}개의 작년 매출 "
                       f"$ {intl['dropped_rev_25']:,.0f} 이 사라졌습니다. 이 매출이 다른 상품으로 흡수되었는지 확인이 필요합니다.")
    bullets.append("")
    bullets.append("권고 사항")
    bullets.append("환율, 항공 운임, 미주 여행 수요와 같은 외부 지표를 함께 결합해 매출 변동의 원인을 더 정확히 파악합니다.")
    bullets.append("매출 변동 상위 상품을 채널, 가격, 마케팅 단위로 더 깊이 들여다 봅니다.")
    bullets.append("새로 추가한 상품과 빠진 상품의 효과를 별도로 추적합니다.")
    add_text_box(s, Inches(0.6), Inches(1.0), Inches(12.3), Inches(6.0),
                 "\n".join(bullets), size=14, color=NAVY)


def slide_external_indicators(prs, kpi, ext):
    """첨부 권고('환율, 항공 운임, 미주 여행 트래픽 결합 분석') 상세 슬라이드"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(s, prs, "외부 지표 결합 분석 가이드  —  환율 · 항공 운임 · 여행 트래픽")

    intro = (
        "현재 매출 데이터만으로는 매출이 늘었는지 줄었는지는 알 수 있지만, "
        "왜 그렇게 되었는지는 알기 어렵습니다. 아래 세 가지 외부 지표를 우리 매출 데이터와 함께 보면, "
        "원인이 외부 환경 때문인지 우리 안의 문제인지 더 분명하게 가릴 수 있습니다."
    )
    add_text_box(s, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.8),
                 intro, size=12, color=GRAY)

    def card(left, color_main, title, body_lines):
        w = Inches(4.05); h = Inches(5.4); top = Inches(1.85)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        box.fill.solid(); box.fill.fore_color.rgb = LIGHT
        box.line.color.rgb = color_main
        add_text_box(s, left, Emu(top + Inches(0.15)), w, Inches(0.5),
                     title, size=15, bold=True, color=color_main, align=PP_ALIGN.CENTER)
        add_text_box(s, Emu(left + Inches(0.2)), Emu(top + Inches(0.7)),
                     Emu(w - Inches(0.4)), Inches(4.6),
                     "\n".join(body_lines), size=11, color=GRAY)

    fx_lines = [
        "확인 데이터: 한국은행 또는 네이버 환율 화면에서 월평균 USD/KRW 환율.",
        "",
        "보는 방법:",
        "1) 환율이 오른 달의 매출 변화와 환율이 내린 달의 매출 변화를 비교합니다.",
        "2) 환율 상승과 매출 감소가 같이 움직이면 환율 영향이 큰 것입니다.",
        "",
        "왜 영향을 주나:",
        "달러가 비싸지면 같은 미국 패키지가 원화로 더 비싸 보여 예약을 미루거나 줄입니다.",
        "",
        "활용 예: 4월 환율이 작년보다 5% 올랐는데 매출이 같이 5% 감소했다면, "
        "감소 원인의 상당 부분이 환율로 설명됩니다.",
    ]
    card(Inches(0.4), NAVY, "환율 (USD / KRW)", fx_lines)

    air_lines = [
        "확인 데이터: 인천-LA, 인천-뉴욕 등 주요 노선의 월평균 항공권 가격, 좌석 공급량.",
        "출처 후보: 항공사 발권 데이터, 스카이스캐너, 인천공항 통계.",
        "",
        "보는 방법:",
        "1) 항공권 가격이 오른 달은 패키지 원가도 함께 오릅니다.",
        "2) 좌석 공급이 줄면 판매할 수 있는 자리 자체가 줄어듭니다.",
        "",
        "왜 영향을 주나:",
        "항공권은 패키지 가격에서 가장 큰 비중을 차지하므로, "
        "운임이 오르면 가격 경쟁력이 떨어지고 예약이 감소합니다.",
        "",
        "활용 예: 항공권이 10% 오른 달에 객단가가 같이 7% 오르고 인원이 5% 감소했다면, "
        "운임 상승이 가격 인상과 수요 감소로 이어졌다고 해석할 수 있습니다.",
    ]
    card(Inches(4.65), SKY, "항공 운임과 좌석 공급", air_lines)

    traffic_lines = [
        "확인 데이터: 미주행 한국인 출국자 수, 미국 여행 검색량, 입국 통계.",
        "출처 후보: 한국관광공사, 인천공항 출국 통계, 구글 트렌드 (검색어: 미국여행, LA여행 등).",
        "",
        "보는 방법:",
        "1) 시장 전체 여행 수요가 늘었는지 줄었는지 확인합니다.",
        "2) 시장은 늘었는데 우리 매출이 줄었다면 우리 점유율이 떨어진 것입니다.",
        "3) 시장이 줄었는데 우리는 유지했다면 상대적으로 잘 한 것입니다.",
        "",
        "왜 영향을 주나:",
        "시장 전체의 흐름은 어느 회사도 피하기 어려운 외부 환경입니다.",
        "",
        "활용 예: 미주행 출국자가 작년보다 15% 증가했는데 우리 매출이 20% 증가했다면, "
        "시장 흐름 + 우리 노력 둘 다 기여한 것입니다.",
    ]
    card(Inches(8.9), GREEN, "미주 여행 트래픽 (시장 수요)", traffic_lines)

    add_text_box(
        s, Inches(0.5), Inches(7.3), Inches(12.3), Inches(0.3),
        "세 지표를 함께 보면, 매출 변화가 외부 환경 때문인지 우리 상품 전략 때문인지 더 명확히 구분할 수 있습니다.",
        size=11, color=NAVY,
    )


def render_pptx(monthly, kpi, prods, ext, intl, charts):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_summary(prs, kpi)
    slide_chart(prs, "월별 매출 추이 (1~4월)",
                charts["monthly_rev"], "단위: USD. 라벨은 YoY 증감률.")
    slide_chart(prs, "월별 인원수 추이 (1~4월)",
                charts["monthly_hc"], "단위: 명. 라벨은 YoY 증감률.")
    slide_chart(prs, "객단가(ARPU) 월별 추이",
                charts["arpu"], "객단가 = 매출 / 인원. 단위: USD.")
    slide_chart(prs, "매출 증가 TOP 10 상품 (동시운영)",
                charts["top_up"], "양 년도 모두 판매된 상품 중 절대 증가액 상위.")
    slide_chart(prs, "매출 감소 TOP 10 상품 (동시운영)",
                charts["top_dn"], "양 년도 모두 판매된 상품 중 절대 감소액 상위.")
    slide_chart(prs, "투어 일수(p_day)별 매출 믹스",
                charts["pday"], "당일 / 1박 / 2~3박 / 4박+ 그룹화 비중.")
    slide_external(prs, ext, monthly)
    slide_internal(prs, intl)
    slide_external_indicators(prs, kpi, ext)
    slide_conclusion(prs, kpi, ext, intl)
    prs.save(OUT_PPTX)


def main():
    print(">> Fetching 2025 ...")
    df25 = fetch_period(2025)
    print(f"   rows={len(df25)}, total_rev={df25['revenue'].sum():,.0f}, total_hc={df25['headcount'].sum():,.0f}")
    print(">> Fetching 2026 ...")
    df26 = fetch_period(2026)
    print(f"   rows={len(df26)}, total_rev={df26['revenue'].sum():,.0f}, total_hc={df26['headcount'].sum():,.0f}")

    monthly, kpi, prods = build_compare(df25, df26)
    ext, intl = detect_signals(monthly, kpi, prods)

    print("\n>> KPI ----------")
    for k, v in kpi.items():
        if isinstance(v, float):
            print(f"   {k:>10s} = {v:,.2f}")
        else:
            print(f"   {k:>10s} = {v}")

    print("\n>> External signals", ext)
    print(">> Internal signals (excerpt) top5_share=%.1f%%, new=%d, dropped=%d" %
          (intl["top5_share"], intl["new_count"], intl["dropped_count"]))

    # raw csv (검증용)
    prods.to_csv(RAW_CSV, index=False, encoding="utf-8-sig")
    print(f">> Raw CSV: {RAW_CSV}")

    # charts
    charts = {}
    charts["monthly_rev"] = chart_monthly_bar(monthly, "rev", "매출 (USD)", "monthly_rev.png")
    charts["monthly_hc"]  = chart_monthly_bar(monthly, "hc",  "인원수 (명)", "monthly_hc.png")
    charts["arpu"]        = chart_arpu(monthly, "arpu.png")
    charts["top_up"]      = chart_top_products(prods, "up", 10, "top_up.png")
    charts["top_dn"]      = chart_top_products(prods, "dn", 10, "top_dn.png")
    charts["pday"]        = chart_pday_mix(prods, "pday.png")

    render_pptx(monthly, kpi, prods, ext, intl, charts)
    print(f"\n>> PPTX saved: {OUT_PPTX}")


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    main()
