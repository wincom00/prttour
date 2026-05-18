# -*- coding: utf-8 -*-
"""
event_assign_user_guide_detail.pptx 내용을 푸른투어 인트라넷 매뉴얼 형식으로 재생성.
기준 스타일: estimate_user_manual_with_view_intranet.pptx
"""
from __future__ import annotations

import os
import time
from dataclasses import dataclass
from typing import Sequence, Tuple

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


NAVY = RGBColor(0x13, 0x11, 0x76)
NAVY_DARK = RGBColor(0x0A, 0x09, 0x55)
BLUE = RGBColor(0x0B, 0x5B, 0xD3)
GREEN = RGBColor(0x28, 0xA7, 0x45)
TEAL = RGBColor(0x20, 0xC9, 0x97)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)
RED = RGBColor(0xDC, 0x35, 0x45)
GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)
GRAY_TXT = RGBColor(0x6C, 0x75, 0x7D)
DARK_TXT = RGBColor(0x21, 0x25, 0x29)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE7, 0xF3, 0xFF)
LIGHT_GREEN = RGBColor(0xE8, 0xF7, 0xEE)
LIGHT_YELLOW = RGBColor(0xFF, 0xF8, 0xE1)

KFONT = "Malgun Gothic"
OUT = r"d:/www/prttour_myprt/event_assign_user_guide_detail_intranet_buttons.pptx"
PREVIEW_PREFIX = r"d:/www/prttour_myprt/event_assign_user_guide_detail_intranet_slide_"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def set_text(shape, text, size=14, bold=False, color=DARK_TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, font=KFONT):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = str(line)
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def add_text(slide, x, y, w, h, text, **kw):
    sh = slide.shapes.add_textbox(x, y, w, h)
    set_text(sh, text, **kw)
    return sh


def add_header(slide, title, subtitle, page, total):
    add_rect(slide, 0, 0, SW, Inches(0.9), NAVY)
    add_text(slide, Inches(0.4), Inches(0.10), Inches(10.4), Inches(0.46),
             title, size=22, bold=True, color=WHITE)
    add_text(slide, Inches(0.4), Inches(0.52), Inches(10.4), Inches(0.30),
             subtitle, size=11, color=RGBColor(0xCB, 0xD3, 0xE8))
    add_text(slide, Inches(11.5), Inches(0.28), Inches(1.55), Inches(0.34),
             f"{page} / {total}", size=11, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT)


def add_footer(slide, source="D:\\www\\hellousa\\admin"):
    add_rect(slide, 0, Inches(7.15), SW, Inches(0.35), GRAY_BG)
    add_text(slide, Inches(0.4), Inches(7.18), Inches(8), Inches(0.3),
             "푸른투어 인트라넷 · 행사배정 사용자 매뉴얼",
             size=9, color=GRAY_TXT)
    add_text(slide, Inches(7.4), Inches(7.18), Inches(5.5), Inches(0.3),
             source, size=9, color=GRAY_TXT, align=PP_ALIGN.RIGHT)


def bullet(slide, x, y, w, items, size=12.2, bullet_color=NAVY, gap=0.40):
    cur_y = y
    for item in items:
        ic = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cur_y + Emu(45000),
                                    Inches(0.22), Inches(0.22))
        ic.fill.solid()
        ic.fill.fore_color.rgb = bullet_color
        ic.line.fill.background()
        set_text(ic, "✓", size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.32), cur_y, w - Inches(0.32), Inches(gap),
                 item, size=size, color=DARK_TXT, anchor=MSO_ANCHOR.MIDDLE)
        cur_y += Inches(gap)


def step_card(slide, x, y, w, h, no, title, body, color=NAVY):
    add_rect(slide, x, y, w, h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(slide, x, y, Inches(0.82), h, color)
    add_text(slide, x, y, Inches(0.82), h, no, size=18, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(1.0), y + Inches(0.18), w - Inches(1.2), Inches(0.38),
             title, size=13, bold=True, color=DARK_TXT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(1.0), y + Inches(0.66), w - Inches(1.2), h - Inches(0.82),
             body, size=10.5, color=GRAY_TXT, anchor=MSO_ANCHOR.TOP)


def table_like(slide, x, y, w, h, headers, rows, widths=None):
    add_rect(slide, x, y, w, h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
    if widths is None:
        widths = [1 / len(headers)] * len(headers)
    header_h = Inches(0.48)
    add_rect(slide, x, y, w, header_h, LIGHT_BLUE, line=RGBColor(0xDD, 0xDD, 0xDD))
    cur_x = x
    for hdr, ratio in zip(headers, widths):
        cw = int(w * ratio)
        add_text(slide, cur_x + Inches(0.06), y + Inches(0.05), cw - Inches(0.12), header_h - Inches(0.10),
                 hdr, size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        cur_x += cw
    row_h = int((h - header_h) / max(1, len(rows)))
    for r_idx, row in enumerate(rows):
        ry = y + header_h + row_h * r_idx
        if r_idx % 2 == 1:
            add_rect(slide, x, ry, w, row_h, RGBColor(0xFB, 0xFC, 0xFE))
        cur_x = x
        for txt, ratio in zip(row, widths):
            cw = int(w * ratio)
            add_text(slide, cur_x + Inches(0.06), ry + Inches(0.03), cw - Inches(0.12), row_h - Inches(0.06),
                     txt, size=8.9, color=DARK_TXT, align=PP_ALIGN.CENTER)
            cur_x += cw


def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, GRAY_BG)
    add_rect(s, 0, 0, SW, Inches(4.5), NAVY)
    for x, y, w in [(Inches(10.2), Inches(-1.5), Inches(5)), (Inches(-1.5), Inches(2.8), Inches(3.5))]:
        deco = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, w)
        deco.fill.solid()
        deco.fill.fore_color.rgb = NAVY_DARK
        deco.line.fill.background()
    add_text(s, Inches(0.8), Inches(0.7), Inches(7), Inches(0.4),
             "푸른투어 인트라넷 · 사용자 매뉴얼", size=14, bold=True,
             color=RGBColor(0xCB, 0xD3, 0xE8))
    add_text(s, Inches(0.8), Inches(1.35), Inches(11.5), Inches(1.15),
             "행사 기본관리 · 통합행사배정", size=47, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(2.48), Inches(11), Inches(0.55),
             "Event Base & Assignment Workflow Guide",
             size=21, color=RGBColor(0xCB, 0xD3, 0xE8))
    add_text(s, Inches(0.8), Inches(3.22), Inches(11.2), Inches(0.55),
             "행사코드 생성, 상태 저장, 배정 대상 조회, 상세 배정 허브 사용 순서",
             size=16, color=WHITE)
    cards = [
        ("행사 준비", "admin/eventbs_list.php", "검색 · 통합행사코드 · 상태 저장", NAVY),
        ("배정 목록", "admin/assign_list.php", "배정 대상 조회 · 현황 배지", GREEN),
        ("상세 배정", "admin/assign_m.php", "차량 · 호텔 · 가이드 업무 이동", ORANGE),
    ]
    for i, (eyebrow, title, desc, color) in enumerate(cards):
        x = Inches(0.8 + i * 4.05)
        add_rect(s, x, Inches(5.0), Inches(3.8), Inches(1.6), WHITE, line=color)
        add_text(s, x + Inches(0.18), Inches(5.1), Inches(3.45), Inches(0.35),
                 eyebrow, size=12, bold=True, color=color)
        add_text(s, x + Inches(0.18), Inches(5.45), Inches(3.45), Inches(0.45),
                 title, size=17, bold=True, color=DARK_TXT)
        add_text(s, x + Inches(0.18), Inches(5.95), Inches(3.45), Inches(0.4),
                 desc, size=11, color=GRAY_TXT)
    add_text(s, Inches(0.8), Inches(6.85), Inches(11), Inches(0.35),
             "푸른투어 인트라넷", size=11, bold=True, color=NAVY)


def slide_toc(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "목차 · CONTENTS", "행사 기본관리와 통합행사배정 업무 흐름", 2, total)
    items = [
        ("01", "페이지 목적", "행사 생성/상태 관리와 배정 실행"),
        ("02", "전체 업무 흐름", "기본관리 → 코드/상태 → 배정목록 → 상세배정"),
        ("03", "행사 기본관리 조회", "검색 조건과 행사 목록"),
        ("04", "상세화면별 설명", "검색, 목록, 코드 생성, 상태 저장"),
        ("05", "통합행사배정 목록", "배정 대상 조회와 현황 배지"),
        ("06", "배정 상세 허브", "예약자보기, 차량/호텔/가이드 배정"),
        ("07", "상태 저장 기준", "예약상태와 행사상태"),
        ("08", "배정현황 확인", "차량, 호텔, 가이드 누락 확인"),
        ("09", "상세 배정 버튼", "예약자보기, 차량/호텔/가이드, 히스토리, 명단"),
        ("10", "운영 표준 순서", "저장 순서와 완료 확인"),
        ("11", "체크리스트", "작업 전후 확인 사항"),
        ("12", "FAQ", "사용 중 자주 확인할 내용"),
    ]
    box_w, box_h = Inches(6.0), Inches(0.66)
    sx, sy = Inches(0.5), Inches(1.10)
    for i, (no, title, sub) in enumerate(items):
        col, row = i % 2, i // 2
        x = sx + (box_w + Inches(0.3)) * col
        y = sy + (box_h + Inches(0.12)) * row
        add_rect(s, x, y, box_w, box_h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, x, y, Inches(0.85), box_h, NAVY)
        add_text(s, x, y, Inches(0.85), box_h, no, size=18, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(1.05), y + Inches(0.09), box_w - Inches(1.2), Inches(0.25),
                 title, size=12.5, bold=True, color=DARK_TXT)
        add_text(s, x + Inches(1.05), y + Inches(0.39), box_w - Inches(1.2), Inches(0.22),
                 sub, size=9.3, color=GRAY_TXT)
    add_footer(s)


def slide_overview(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "01. 페이지 목적", "행사 생성/상태 관리와 배정 실행을 분리", 3, total)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.9), LIGHT_BLUE)
    add_text(s, Inches(0.7), Inches(1.23), Inches(11.9), Inches(0.58),
             "admin/eventbs_list.php에서 행사 기준을 만들고, admin/assign_list.php에서 배정 대상을 열어 상세 배정 업무로 이동합니다.",
             size=15, bold=True, color=NAVY)
    cards = [
        ("행사 기본관리", "코드 · 상태", ["상품과 출발일 검색", "통합행사코드 일괄생성", "예약상태/행사상태 저장"], NAVY),
        ("통합행사배정", "대상 조회", ["배정할 행사 목록 확인", "예약/대기 인원 확인", "차량/호텔/가이드 배지 확인"], GREEN),
        ("상세 배정", "업무 이동", ["예약자보기", "차량/호텔/가이드 배정", "명단/메일 팝업 이동"], ORANGE),
    ]
    for i, (label, title, items, color) in enumerate(cards):
        x = Inches(0.55 + i * 4.18)
        add_rect(s, x, Inches(2.35), Inches(3.9), Inches(4.35), WHITE, line=color)
        add_text(s, x + Inches(0.2), Inches(2.55), Inches(1.45), Inches(0.42),
                 label, size=14, bold=True, color=color)
        add_text(s, x + Inches(0.2), Inches(3.05), Inches(3.45), Inches(0.45),
                 title, size=20, bold=True, color=DARK_TXT)
        bullet(s, x + Inches(0.25), Inches(3.75), Inches(3.4), items,
               size=12.2, bullet_color=color, gap=0.48)
    add_footer(s)


def slide_flow(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "02. 전체 업무 흐름", "기본관리에서 배정 완료 확인까지", 4, total)
    steps = [
        ("1", "행사 조회", "상품명/상품코드, 출발일, 행사상태로 대상을 검색합니다.", BLUE),
        ("2", "코드 생성", "체크한 행에 통합행사코드를 일괄생성합니다.", GREEN),
        ("3", "상태 저장", "예약상태와 행사상태를 선택 행 기준으로 저장합니다.", TEAL),
        ("4", "배정 목록", "admin/assign_list.php에서 통합행사코드와 배정현황을 확인합니다.", ORANGE),
        ("5", "상세 진입", "목록 행 클릭으로 admin/assign_m.php 화면을 엽니다.", RED),
        ("6", "배정 완료", "차량, 호텔, 가이드 배지가 채워졌는지 확인합니다.", NAVY),
    ]
    for i, (no, title, body, color) in enumerate(steps):
        row, col = divmod(i, 3)
        x = Inches(0.7 + col * 4.15)
        y = Inches(1.35 + row * 2.35)
        step_card(s, x, y, Inches(3.65), Inches(1.75), no, title, body, color)
    add_footer(s)


def slide_event_list(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "03. 행사 기본관리 조회", "eventbs_list.php: 검색 조건과 행사 목록", 5, total)
    add_rect(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.55), LIGHT_BLUE)
    add_text(s, Inches(0.7), Inches(1.14), Inches(11.8), Inches(0.35),
             "조회 대상, 상품명/코드, 출발일 범위, 행사상태를 입력한 뒤 검색합니다.",
             size=13, bold=True, color=NAVY)
    table_like(
        s, Inches(0.65), Inches(1.9), Inches(12.0), Inches(2.35),
        ["통합행사코드", "상품분류", "상품코드", "상품명", "출발일", "정원", "예약수", "예약상태", "행사상태"],
        [
            ["비어 있으면 생성 대상", "분류명", "p_code", "p_name", "stDate", "p_cnt", "예약/대기", "r_status", "ev_status"],
            ["생성 완료 시 코드 표시", "지역/구분", "상품코드", "상품명 링크", "출발 날짜", "정원", "예약 인원", "예약 상태", "행사 상태"],
        ],
        [0.15, 0.10, 0.12, 0.20, 0.11, 0.07, 0.10, 0.08, 0.07],
    )
    step_card(s, Inches(0.65), Inches(4.65), Inches(3.85), Inches(1.55),
              "A", "검색 조건", "예약된 상품만 보기 또는 전체 시즌 보기를 선택하고 기간/상태를 입력합니다.", BLUE)
    step_card(s, Inches(4.75), Inches(4.65), Inches(3.85), Inches(1.55),
              "B", "목록 확인", "통합행사코드, 상품명, 출발일, 정원, 예약수를 확인합니다.", GREEN)
    step_card(s, Inches(8.85), Inches(4.65), Inches(3.85), Inches(1.55),
              "C", "행 선택", "체크박스 선택 행이 이후 코드 생성/상태 저장 범위입니다.", ORANGE)
    add_footer(s, "D:\\www\\hellousa\\admin\\eventbs_list.php")


def slide_event_detail(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "04. 상세화면별 설명", "검색, 목록, 코드 생성, 상태 저장", 6, total)
    blocks = [
        ("A", "검색 조건", "조회 대상, 상품명 또는 상품코드, 출발일 시작/종료일, 행사상태를 입력한 뒤 검색합니다.", BLUE),
        ("B", "행사 목록", "왼쪽 체크박스가 저장 범위입니다. 코드가 없는 행은 생성 대상, 코드가 있는 행은 기존 행사 기준입니다.", GREEN),
        ("C", "통합행사코드 일괄생성", "선택 후 버튼을 누르면 mode=save로 제출됩니다. 신규 행은 tour_master에 생성됩니다.", ORANGE),
        ("D", "상태 일괄저장", "예약상태 저장은 r_status, 행사상태 저장은 ev_status를 선택 행 기준으로 업데이트합니다.", RED),
    ]
    for i, (no, title, body, color) in enumerate(blocks):
        row, col = divmod(i, 2)
        x = Inches(0.65 + col * 6.15)
        y = Inches(1.25 + row * 2.45)
        step_card(s, x, y, Inches(5.75), Inches(1.85), no, title, body, color)
    add_rect(s, Inches(0.65), Inches(6.20), Inches(12.0), Inches(0.55), LIGHT_YELLOW)
    add_text(s, Inches(0.9), Inches(6.30), Inches(11.4), Inches(0.28),
             "저장 버튼은 기능별로 분리되어 있으므로 코드 생성, 예약상태, 행사상태 버튼을 혼동하지 않아야 합니다.",
             size=12.3, bold=True, color=ORANGE)
    add_footer(s, "D:\\www\\hellousa\\admin\\eventbs_list.php")


def slide_assign_list(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "05. 통합행사배정 목록", "admin/assign_list.php: 배정 대상 조회와 현황 배지", 7, total)
    table_like(
        s, Inches(0.65), Inches(1.25), Inches(12.0), Inches(2.25),
        ["통합행사코드", "상품분류", "지역분류", "상품코드", "상품명", "출발일", "정원", "예약/대기", "배정현황"],
        [
            ["grand_eCode", "c_code1", "c_code2", "p_code", "p_name", "stDate", "p_cnt", "예약/대기", "차량/호텔/가이드"],
            ["클릭 대상", "분류", "지역", "상품코드", "행사명", "출발일", "정원", "인원", "카운트 배지"],
        ],
        [0.14, 0.10, 0.10, 0.12, 0.18, 0.10, 0.07, 0.10, 0.19],
    )
    cards = [
        ("조회", "상품명, 출발일 범위, 행사상태로 배정 대상 행사를 검색합니다.", BLUE),
        ("행 클릭", "목록 행을 클릭하면 상세 배정 화면으로 이동합니다.", GREEN),
        ("현황 확인", "배정현황 배지로 차량/호텔/가이드 누락을 확인합니다.", ORANGE),
    ]
    for i, (title, body, color) in enumerate(cards):
        step_card(s, Inches(0.65 + i * 4.1), Inches(4.05), Inches(3.85), Inches(1.55),
                  chr(65 + i), title, body, color)
    add_footer(s, "D:\\www\\hellousa\\admin\\assign_list.php")


def slide_assign_detail(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "06. 배정 상세 허브", "admin/assign_m.php: 상세 배정 업무로 이동", 8, total)
    add_rect(s, Inches(0.65), Inches(1.15), Inches(5.85), Inches(4.95), WHITE, line=NAVY)
    add_text(s, Inches(0.9), Inches(1.38), Inches(5.25), Inches(0.42),
             "상단 확인 정보", size=22, bold=True, color=NAVY)
    bullet(s, Inches(0.95), Inches(2.02), Inches(5.1), [
        "통합행사코드",
        "상품명과 상품코드",
        "출발일",
        "투어정원과 예약인원",
        "예약상태와 행사상태",
        "기타 상태 메모",
    ], size=12.5, bullet_color=NAVY, gap=0.45)
    add_rect(s, Inches(6.85), Inches(1.15), Inches(5.85), Inches(4.95), WHITE, line=ORANGE)
    add_text(s, Inches(7.1), Inches(1.38), Inches(5.25), Inches(0.42),
             "하단 업무 버튼", size=22, bold=True, color=ORANGE)
    bullet(s, Inches(7.15), Inches(2.02), Inches(5.1), [
        "예약자보기",
        "차량배정",
        "호텔배정",
        "가이드배정",
        "히스토리등록",
        "명단 / 메일보내기",
    ], size=12.5, bullet_color=ORANGE, gap=0.45)
    add_footer(s, "D:\\www\\hellousa\\admin\\assign_m.php")


def slide_status(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "07. 상태 저장 기준", "예약상태와 행사상태를 구분해서 저장", 9, total)
    items = [
        ("예약상태", "r_status", "예약 관련 운영 상태를 선택 행 기준으로 저장합니다.", BLUE),
        ("행사상태", "ev_status", "미확정, 확정, 만차, 취소, 기타 상태를 선택 행 기준으로 저장합니다.", GREEN),
        ("저장 범위", "체크박스", "목록에서 체크한 행만 코드 생성 또는 상태 저장 대상입니다.", ORANGE),
        ("저장 후", "조회조건 유지", "eventbs_list.php는 저장 후 기존 조회조건을 유지하는 쿼리로 돌아갑니다.", NAVY),
    ]
    for i, (title, key, desc, color) in enumerate(items):
        y = Inches(1.15 + i * 1.10)
        add_rect(s, Inches(0.65), y, Inches(12.0), Inches(0.86), WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, Inches(0.65), y, Inches(1.55), Inches(0.86), color)
        add_text(s, Inches(0.65), y, Inches(1.55), Inches(0.86),
                 title, size=13.2, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.4), y + Inches(0.10), Inches(2.0), Inches(0.28),
                 key, size=12.0, bold=True, color=DARK_TXT)
        add_text(s, Inches(2.4), y + Inches(0.45), Inches(9.7), Inches(0.25),
                 desc, size=10.5, color=GRAY_TXT)
    add_footer(s, "D:\\www\\hellousa\\admin\\eventbs_list.php")


def slide_badges(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "08. 배정현황 확인", "차량, 호텔, 가이드 누락 확인", 10, total)
    badges = [
        ("차량", "tour_car 카운트 기준", BLUE),
        ("호텔", "hotel_assign 카운트 기준", GREEN),
        ("가이드", "tour_guide 카운트 기준", ORANGE),
    ]
    for i, (title, desc, color) in enumerate(badges):
        x = Inches(1.10 + i * 4.05)
        add_rect(s, x, Inches(1.65), Inches(3.25), Inches(2.20), WHITE, line=color)
        add_rect(s, x, Inches(1.65), Inches(3.25), Inches(0.52), color)
        add_text(s, x, Inches(1.73), Inches(3.25), Inches(0.30),
                 title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.18), Inches(2.45), Inches(2.88), Inches(0.65),
                 desc, size=11.2, color=DARK_TXT, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.18), Inches(3.16), Inches(2.88), Inches(0.28),
                 "있음/없음 배지로 표시", size=9.8, color=GRAY_TXT, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.65), Inches(4.65), Inches(12.0), Inches(1.05), LIGHT_GREEN)
    add_text(s, Inches(0.9), Inches(4.82), Inches(11.4), Inches(0.32),
             "운영 기준", size=15, bold=True, color=GREEN)
    add_text(s, Inches(0.9), Inches(5.23), Inches(11.4), Inches(0.32),
             "admin/assign_list.php 목록에서 차량없음, 호텔없음, 가이드없음 배지를 먼저 확인하고 상세 배정 화면으로 이동합니다.",
             size=12.6, color=DARK_TXT)
    add_footer(s, "D:\\www\\hellousa\\admin\\assign_list.php")


def slide_standard_order(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "10. 운영 표준 순서", "저장 순서를 지키는 것이 중요합니다", 12, total)
    steps = [
        ("1", "기본 조회", "eventbs_list.php에서 기간과 상품을 검색합니다.", BLUE),
        ("2", "코드 생성", "코드 없는 행을 선택하고 통합행사코드를 생성합니다.", GREEN),
        ("3", "상태 확정", "예약상태와 행사상태를 선택 행 기준으로 저장합니다.", TEAL),
        ("4", "배정 진행", "admin/assign_list.php에서 행사를 열고 상세 배정으로 이동합니다.", ORANGE),
        ("5", "누락 확인", "배정현황 배지에서 차량/호텔/가이드 누락을 확인합니다.", RED),
        ("6", "후속 처리", "예약자보기, 명단, 메일 발송 등 필요한 팝업 업무를 진행합니다.", NAVY),
    ]
    for i, (no, title, body, color) in enumerate(steps):
        row, col = divmod(i, 3)
        x = Inches(0.7 + col * 4.15)
        y = Inches(1.35 + row * 2.35)
        step_card(s, x, y, Inches(3.65), Inches(1.75), no, title, body, color)
    add_footer(s)


def slide_button_roles(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "09. 상세 배정 버튼 역할", "admin/assign_m.php 하단 아이콘 메뉴", 11, total)

    add_rect(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.62), LIGHT_BLUE)
    add_text(s, Inches(0.72), Inches(1.16), Inches(11.85), Inches(0.34),
             "상세페이지 하단 아이콘 메뉴에서 예약자보기, 차량/호텔/가이드 배정, 히스토리등록, 명단을 엽니다.",
             size=13.2, bold=True, color=NAVY)

    rows = [
        ["예약자보기", "예약자 상세 목록 확인", "event_reservation_detail.php", "예약 인원과 예약자 정보를 확인합니다."],
        ["차량배정", "차량 배정 화면 이동", "car_assign_m.php", "일반 상품은 차량배정 화면으로 이동하며, PICUP 상품은 링크가 비활성 표시됩니다."],
        ["호텔배정", "호텔 배정 화면 이동", "hotel_assign_m.php", "행사별 호텔 배정 내용을 등록하거나 확인합니다."],
        ["가이드배정", "가이드 배정 화면 이동", "guide_assign_m.php", "메인/로컬 가이드 배정을 등록하거나 확인합니다."],
        ["히스토리등록", "가이드 배정 히스토리 저장", "assign_m.php?action=guide_history", "확인창 후 가이드 배정 히스토리를 메모에 등록합니다."],
        ["명단", "명단 팝업 열기", "short_customer.php", "상품코드, 통합행사코드, 출발일 기준 명단 팝업을 엽니다."],
    ]
    table_like(
        s, Inches(0.65), Inches(1.95), Inches(12.0), Inches(3.35),
        ["버튼", "역할", "연결 화면/처리", "사용 시점"],
        rows,
        [0.14, 0.24, 0.28, 0.34],
    )

    step_card(s, Inches(0.65), Inches(5.70), Inches(3.85), Inches(1.05),
              "A", "먼저 확인", "상단의 통합행사코드, 상품명, 출발일이 맞는지 확인합니다.", BLUE)
    step_card(s, Inches(4.75), Inches(5.70), Inches(3.85), Inches(1.05),
              "B", "배정 작업", "차량, 호텔, 가이드 버튼으로 각각의 세부 배정 화면을 엽니다.", GREEN)
    step_card(s, Inches(8.85), Inches(5.70), Inches(3.85), Inches(1.05),
              "C", "후속 처리", "히스토리등록과 명단 팝업은 배정 확인 후 필요한 경우 실행합니다.", ORANGE)
    add_footer(s, "D:\\www\\hellousa\\admin\\assign_m.php")


def slide_checklist(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "12. 체크리스트", "작업 전후 확인 사항", 13, total)
    groups = [
        ("조회 전", ["출발일 범위 확인", "상품명/상품코드 확인", "행사상태 조건 확인"], BLUE),
        ("저장 전", ["체크박스 선택 행 확인", "코드 있는 행/없는 행 구분", "예약상태와 행사상태 버튼 구분"], GREEN),
        ("배정 전", ["통합행사코드 확인", "예약/대기 인원 확인", "상세 화면 업무 버튼 확인"], ORANGE),
        ("완료 후", ["차량 배지 확인", "호텔 배지 확인", "가이드 배지 확인"], NAVY),
    ]
    for i, (title, items, color) in enumerate(groups):
        row, col = divmod(i, 2)
        x = Inches(0.65 + col * 6.15)
        y = Inches(1.25 + row * 2.45)
        add_rect(s, x, y, Inches(5.75), Inches(2.0), WHITE, line=color)
        add_rect(s, x, y, Inches(5.75), Inches(0.45), color)
        add_text(s, x + Inches(0.18), y + Inches(0.08), Inches(5.4), Inches(0.25),
                 title, size=13.5, bold=True, color=WHITE)
        bullet(s, x + Inches(0.25), y + Inches(0.72), Inches(5.25), items,
               size=12.2, bullet_color=color, gap=0.40)
    add_footer(s)


def slide_faq(total):
    s = prs.slides.add_slide(BLANK)
    add_header(s, "13. FAQ", "사용 중 자주 확인할 내용", 14, total)
    faqs = [
        ("Q1", "통합행사코드가 비어 있을 때", "eventbs_list.php에서 행을 체크한 뒤 통합행사코드 일괄생성을 실행합니다."),
        ("Q2", "상태가 반영되지 않을 때", "선택한 체크박스 행과 예약상태/행사상태 저장 버튼을 다시 확인합니다."),
        ("Q3", "배정 상세로 이동할 때", "admin/assign_list.php 목록 행을 클릭해 admin/assign_m.php 화면을 엽니다."),
        ("Q4", "배정 누락을 찾을 때", "배정현황의 차량, 호텔, 가이드 배지를 기준으로 확인합니다."),
        ("Q5", "상세 배정 버튼 위치", "admin/assign_m.php 하단 아이콘 메뉴에서 예약자보기, 차량/호텔/가이드 배정, 명단을 엽니다."),
    ]
    y = Inches(1.15)
    for q, title, ans in faqs:
        add_rect(s, Inches(0.65), y, Inches(12.0), Inches(0.86), WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, Inches(0.65), y, Inches(0.85), Inches(0.86), NAVY)
        add_text(s, Inches(0.65), y, Inches(0.85), Inches(0.86),
                 q, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.7), y + Inches(0.08), Inches(10.6), Inches(0.30),
                 title, size=12.5, bold=True, color=DARK_TXT)
        add_text(s, Inches(1.7), y + Inches(0.44), Inches(10.6), Inches(0.25),
                 ans, size=10.4, color=GRAY_TXT)
        y += Inches(0.96)
    add_footer(s)


def slide_thanks(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    for x, y, w in [(Inches(-2), Inches(-2), Inches(6)), (Inches(9.5), Inches(4), Inches(6))]:
        deco = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, w)
        deco.fill.solid()
        deco.fill.fore_color.rgb = NAVY_DARK
        deco.line.fill.background()
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0),
             "Thank You", size=66, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.6),
             "행사 기본관리 → 통합행사배정 → 상세 배정 순서로 사용합니다.",
             size=20, color=RGBColor(0xCB, 0xD3, 0xE8))
    add_rect(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.45), WHITE)
    add_text(s, Inches(1.0), Inches(4.65), Inches(11), Inches(0.35),
             "확인 파일", size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(5.05), Inches(11), Inches(0.55),
             "eventbs_list.php / assign_list.php / D:\\www\\hellousa\\admin\\assign_m.php",
             size=15, bold=True, color=DARK_TXT)
    add_text(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
             "푸른투어 인트라넷 · event assignment manual",
             size=11, color=RGBColor(0xCB, 0xD3, 0xE8), align=PP_ALIGN.CENTER)


@dataclass
class PreviewSlide:
    title: str
    subtitle: str
    rows: Sequence[Tuple[str, str]]
    dark: bool = False


def rgb_tuple(color: RGBColor):
    return (color[0], color[1], color[2])


def load_font(size: int, bold: bool = False):
    path = r"C:/Windows/Fonts/malgunbd.ttf" if bold else r"C:/Windows/Fonts/malgun.ttf"
    if os.path.exists(path):
        return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def wrap_text(draw, text, font, max_width):
    words = str(text).split()
    if not words:
        return [""]
    lines, cur = [], words[0]
    for word in words[1:]:
        trial = cur + " " + word
        if draw.textbbox((0, 0), trial, font=font)[2] <= max_width:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    lines.append(cur)
    return lines


def draw_preview(idx: int, total: int, spec: PreviewSlide):
    W, H = 1920, 1080
    img = Image.new("RGB", (W, H), rgb_tuple(GRAY_BG))
    d = ImageDraw.Draw(img)
    if spec.dark:
        d.rectangle([0, 0, W, H], fill=rgb_tuple(NAVY))
        d.ellipse([1450, -220, 2200, 530], fill=rgb_tuple(NAVY_DARK))
        d.ellipse([-240, 420, 430, 1090], fill=rgb_tuple(NAVY_DARK))
        d.text((115, 280), spec.title, fill=rgb_tuple(WHITE), font=load_font(72, True))
        d.text((120, 415), spec.subtitle, fill=(203, 211, 232), font=load_font(25))
    else:
        d.rectangle([0, 0, W, 130], fill=rgb_tuple(NAVY))
        d.text((58, 28), spec.title, fill=rgb_tuple(WHITE), font=load_font(54, True))
        d.text((62, 86), spec.subtitle, fill=(203, 211, 232), font=load_font(18))
        d.text((1700, 48), f"{idx} / {total}", fill=rgb_tuple(WHITE), font=load_font(18))
        y = 170
        for i, (head, body) in enumerate(spec.rows):
            x = 90 + (i % 2) * 870
            if i and i % 2 == 0:
                y += 240
            d.rounded_rectangle([x, y, x + 790, y + 185], radius=8,
                                fill=rgb_tuple(WHITE), outline=(221, 221, 221), width=2)
            d.rectangle([x, y, x + 112, y + 185], fill=rgb_tuple(NAVY))
            d.text((x + 34, y + 58), str(i + 1), fill=rgb_tuple(WHITE), font=load_font(40, True))
            d.text((x + 145, y + 25), head, fill=rgb_tuple(DARK_TXT), font=load_font(27, True))
            line_y = y + 76
            for line in wrap_text(d, body, load_font(23), 590)[:3]:
                d.text((x + 145, line_y), line, fill=rgb_tuple(GRAY_TXT), font=load_font(23))
                line_y += 34
        d.rectangle([0, 1030, W, H], fill=rgb_tuple(GRAY_BG))
        d.text((58, 1040), "푸른투어 인트라넷 · 행사배정 사용자 매뉴얼",
               fill=rgb_tuple(GRAY_TXT), font=load_font(18))
    out = f"{PREVIEW_PREFIX}{idx:02d}.png"
    img.save(out)
    return out


def render_previews(total: int):
    specs = [
        PreviewSlide("행사 기본관리 · 통합행사배정", "Event Base & Assignment Workflow Guide", [], True),
        PreviewSlide("목차 · CONTENTS", "행사 기본관리와 통합행사배정 업무 흐름", [("페이지 목적", "행사 생성/상태 관리와 배정 실행"), ("전체 업무 흐름", "기본관리 → 배정 목록 → 상세 배정")]),
        PreviewSlide("01. 페이지 목적", "행사 생성/상태 관리와 배정 실행을 분리", [("행사 기본관리", "코드 생성과 상태 저장"), ("통합행사배정", "대상 조회와 현황 배지"), ("상세 배정", "차량/호텔/가이드 업무 이동")]),
        PreviewSlide("02. 전체 업무 흐름", "기본관리에서 배정 완료 확인까지", [("행사 조회", "상품명/출발일/상태 검색"), ("코드 생성", "통합행사코드 일괄생성"), ("배정 완료", "현황 배지 확인")]),
        PreviewSlide("03. 행사 기본관리 조회", "eventbs_list.php 검색 조건과 행사 목록", [("검색 조건", "상품명/코드, 기간, 행사상태"), ("목록 확인", "통합행사코드와 예약 수"), ("행 선택", "체크박스 저장 범위")]),
        PreviewSlide("04. 상세화면별 설명", "검색, 목록, 코드 생성, 상태 저장", [("검색 조건", "조회 대상과 기간 입력"), ("행사 목록", "체크박스 기준 저장"), ("상태 저장", "예약상태와 행사상태 구분")]),
        PreviewSlide("05. 통합행사배정 목록", "admin/assign_list.php 배정 대상 조회와 현황 배지", [("조회", "배정 대상 행사 검색"), ("행 클릭", "admin/assign_m.php 이동"), ("현황 확인", "차량/호텔/가이드")]),
        PreviewSlide("06. 배정 상세 허브", "admin/assign_m.php 상세 배정 업무로 이동", [("상단 정보", "통합행사코드와 출발일"), ("업무 버튼", "예약자/차량/호텔/가이드"), ("후속 업무", "명단 팝업")]),
        PreviewSlide("07. 상태 저장 기준", "예약상태와 행사상태 구분", [("예약상태", "r_status 저장"), ("행사상태", "ev_status 저장"), ("저장 범위", "체크박스 선택 행")]),
        PreviewSlide("08. 배정현황 확인", "차량, 호텔, 가이드 누락 확인", [("차량", "tour_car 기준"), ("호텔", "hotel_assign 기준"), ("가이드", "tour_guide 기준")]),
        PreviewSlide("09. 상세 배정 버튼 역할", "admin/assign_m.php 하단 아이콘 메뉴", [("예약자보기", "예약자 상세 목록 확인"), ("차량/호텔/가이드", "세부 배정 화면 이동"), ("히스토리/명단", "메모 등록과 명단 팝업")]),
        PreviewSlide("10. 운영 표준 순서", "저장 순서를 지키는 것이 중요", [("기본 조회", "기간과 상품 검색"), ("상태 확정", "예약/행사 상태 저장"), ("배정 진행", "누락 배지 확인")]),
        PreviewSlide("11. 체크리스트", "작업 전후 확인 사항", [("조회 전", "기간/상품/상태 확인"), ("저장 전", "선택 행 확인"), ("완료 후", "현황 배지 확인")]),
        PreviewSlide("12. FAQ", "사용 중 자주 확인할 내용", [("코드 없음", "통합행사코드 생성"), ("상태 미반영", "체크박스와 버튼 확인"), ("배정 누락", "현황 배지 기준 확인")]),
        PreviewSlide("Thank You", "행사 기본관리 → 통합행사배정 → 상세 배정", [], True),
    ]
    return [draw_preview(i + 1, total, spec) for i, spec in enumerate(specs)]


def build():
    total = 15
    slide_cover()
    slide_toc(total)
    slide_overview(total)
    slide_flow(total)
    slide_event_list(total)
    slide_event_detail(total)
    slide_assign_list(total)
    slide_assign_detail(total)
    slide_status(total)
    slide_badges(total)
    slide_button_roles(total)
    slide_standard_order(total)
    slide_checklist(total)
    slide_faq(total)
    slide_thanks(total)
    try:
        prs.save(OUT)
        saved = OUT
    except PermissionError:
        saved = OUT.replace(".pptx", f"_v{int(time.time())}.pptx")
        prs.save(saved)
    previews = render_previews(total)
    print("SAVED:", saved)
    print("PREVIEWS:", len(previews))
    for path in previews:
        print(path)


if __name__ == "__main__":
    build()
