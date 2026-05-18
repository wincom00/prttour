# -*- coding: utf-8 -*-
"""
예약상담등록 사용자 매뉴얼 PPTX 생성.
기준 화면:
- base_consult.php: 상품 검색/선택 화면
- base_conslut_m.php: 상담예약 입력 화면
"""
from __future__ import annotations

import os
from dataclasses import dataclass
from typing import Iterable, Sequence

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = r"d:/www/prttour_myprt/consult_user_manual.pptx"
PREVIEW_PREFIX = r"d:/www/prttour_myprt/consult_user_manual_slide_"
KFONT = "Malgun Gothic"

NAVY = RGBColor(18, 20, 105)
NAVY_DARK = RGBColor(10, 12, 67)
BLUE = RGBColor(38, 103, 232)
CYAN = RGBColor(229, 241, 255)
GREEN = RGBColor(37, 145, 90)
ORANGE = RGBColor(238, 132, 37)
RED = RGBColor(202, 59, 78)
BG = RGBColor(246, 248, 252)
LINE = RGBColor(214, 221, 232)
TEXT = RGBColor(32, 37, 45)
MUTED = RGBColor(96, 105, 118)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.8)
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, value, size=14, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    sh = slide.shapes.add_textbox(x, y, w, h)
    tf = sh.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    lines = value if isinstance(value, list) else [value]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = str(line)
        r.font.name = KFONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return sh


def header(slide, title, subtitle, page, total):
    rect(slide, 0, 0, SW, Inches(0.88), NAVY)
    text(slide, Inches(0.45), Inches(0.08), Inches(10.3), Inches(0.48),
         title, size=22, bold=True, color=WHITE)
    text(slide, Inches(0.48), Inches(0.50), Inches(10.4), Inches(0.28),
         subtitle, size=10.5, color=RGBColor(208, 216, 238))
    text(slide, Inches(11.3), Inches(0.28), Inches(1.55), Inches(0.28),
         f"{page} / {total}", size=11, bold=True, color=WHITE,
         align=PP_ALIGN.RIGHT)


def footer(slide):
    text(slide, Inches(0.45), Inches(7.12), Inches(12.4), Inches(0.25),
         "출처: base_consult.php, base_conslut_m.php", size=8.5,
         color=MUTED, align=PP_ALIGN.RIGHT)


def chip(slide, x, y, label, color=BLUE):
    rect(slide, x, y, Inches(1.35), Inches(0.32), color, radius=True)
    text(slide, x, y + Inches(0.01), Inches(1.35), Inches(0.25),
         label, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def bullets(slide, x, y, w, items: Iterable[str], size=12, gap=0.42, color=BLUE):
    yy = y
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, yy + Inches(0.12),
                                     Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        text(slide, x + Inches(0.25), yy, w - Inches(0.25), Inches(0.32),
             item, size=size, color=TEXT)
        yy += Inches(gap)


def card(slide, x, y, w, h, title, body, color=BLUE, no=None):
    rect(slide, x, y, w, h, WHITE, LINE, radius=True)
    if no is not None:
        rect(slide, x, y, Inches(0.62), h, color)
        text(slide, x, y, Inches(0.62), h, str(no), size=17,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x2 = x + Inches(0.78)
        w2 = w - Inches(0.95)
    else:
        x2 = x + Inches(0.18)
        w2 = w - Inches(0.36)
    text(slide, x2, y + Inches(0.14), w2, Inches(0.32),
         title, size=13.5, bold=True, color=TEXT)
    text(slide, x2, y + Inches(0.54), w2, h - Inches(0.65),
         body, size=10.8, color=MUTED, anchor=MSO_ANCHOR.TOP)


def simple_table(slide, x, y, w, headers, rows, widths):
    header_h = Inches(0.38)
    row_h = Inches(0.42)
    total_h = header_h + row_h * len(rows)
    rect(slide, x, y, w, total_h, WHITE, LINE)
    rect(slide, x, y, w, header_h, CYAN, LINE)
    xx = x
    for header_txt, ratio in zip(headers, widths):
        cw = int(w * ratio)
        text(slide, xx, y + Inches(0.05), cw, Inches(0.22), header_txt,
             size=8.8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        xx += cw
    for r, row in enumerate(rows):
        yy = y + header_h + row_h * r
        if r % 2:
            rect(slide, x, yy, w, row_h, RGBColor(250, 252, 255))
        xx = x
        for cell, ratio in zip(row, widths):
            cw = int(w * ratio)
            text(slide, xx + Inches(0.04), yy + Inches(0.06),
                 cw - Inches(0.08), Inches(0.22), cell, size=8.2,
                 color=TEXT, align=PP_ALIGN.CENTER)
            xx += cw


def slide_cover(total):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    deco1 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.8), Inches(-1.6), Inches(5.3), Inches(5.3))
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = NAVY_DARK
    deco1.line.fill.background()
    deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.4), Inches(4.4), Inches(3.8), Inches(3.8))
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = NAVY_DARK
    deco2.line.fill.background()
    text(s, Inches(0.78), Inches(0.95), Inches(10.5), Inches(0.35),
         "푸른투어 인트라넷", size=15, bold=True, color=RGBColor(208, 216, 238))
    text(s, Inches(0.74), Inches(1.62), Inches(11.5), Inches(1.0),
         "예약상담등록 사용자 매뉴얼", size=42, bold=True, color=WHITE)
    text(s, Inches(0.80), Inches(2.78), Inches(10.8), Inches(0.52),
         "상품 검색/선택 화면 + 상담예약 입력 화면", size=18,
         color=RGBColor(208, 216, 238))
    card(s, Inches(0.82), Inches(4.78), Inches(3.7), Inches(1.1),
         "첨부 1", "검색 조건으로 상품을 찾고 선택 버튼으로 상담 입력 화면으로 이동", BLUE)
    card(s, Inches(4.82), Inches(4.78), Inches(3.7), Inches(1.1),
         "첨부 2", "예약자 정보와 상담내용을 입력하고 상담등록으로 저장", GREEN)
    card(s, Inches(8.82), Inches(4.78), Inches(3.7), Inches(1.1),
         "연결 업무", "저장 후 상담이력 확인, 필요 시 견적등록으로 예약 업무 전환", ORANGE)
    text(s, Inches(0.78), Inches(6.75), Inches(11.8), Inches(0.3),
         f"총 {total} 슬라이드", size=10.5, color=RGBColor(208, 216, 238),
         align=PP_ALIGN.RIGHT)


def slide_flow(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "01. 업무 흐름", "예약상담등록은 상품 선택 후 상담 정보를 저장하는 순서입니다.", 2, total)
    steps = [
        ("메뉴 진입", "예약상담관리 > 예약상담등록 메뉴로 이동"),
        ("상품 검색", "투어분류, 지역분류, 상품명, 출발일 조건 입력"),
        ("상품 선택", "목록에서 대상 상품의 선택 버튼 클릭"),
        ("상담 입력", "여행기간, 전달직원, 방갯수, 여행인원, 예약자 정보 입력"),
        ("저장/전환", "상담등록 저장 후 상담이력 확인 또는 견적등록 진행"),
    ]
    x = Inches(0.55)
    for i, (title, body) in enumerate(steps, start=1):
        card(s, x, Inches(1.45), Inches(2.35), Inches(3.1), title, body, BLUE, i)
        if i < len(steps):
            text(s, x + Inches(2.36), Inches(2.75), Inches(0.28), Inches(0.28),
                 ">", size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        x += Inches(2.52)
    rect(s, Inches(0.75), Inches(5.25), Inches(11.8), Inches(0.9), RGBColor(238, 245, 255), LINE, radius=True)
    text(s, Inches(1.0), Inches(5.42), Inches(11.3), Inches(0.45),
         "핵심: 첨부 1에서 상품을 정확히 선택해야 첨부 2의 상품명/상품코드가 자동으로 연결됩니다.",
         size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    footer(s)


def slide_search_screen(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "02. 첨부 1 화면: 상품 검색/선택", "base_consult.php 화면 구성", 3, total)
    rect(s, Inches(0.55), Inches(1.15), Inches(12.25), Inches(5.45), WHITE, LINE, radius=True)
    text(s, Inches(0.8), Inches(1.38), Inches(11.7), Inches(0.32),
         "검색 조건 영역", size=15, bold=True, color=NAVY)
    labels = [("투어분류", "로컬/인바운드/인센티브/아웃바운드"),
              ("지역분류", "1차 지역 + 2차 지역"),
              ("상품명", "키워드 입력"),
              ("출발일", "시작일 ~ 마지막일")]
    y = Inches(1.95)
    for i, (label, val) in enumerate(labels):
        row = i // 2
        col = i % 2
        x = Inches(0.95 + col * 5.85)
        yy = y + Inches(row * 0.78)
        rect(s, x, yy, Inches(1.25), Inches(0.42), RGBColor(239, 242, 247), LINE)
        text(s, x, yy + Inches(0.06), Inches(1.25), Inches(0.22), label,
             size=9.5, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        rect(s, x + Inches(1.3), yy, Inches(4.2), Inches(0.42), WHITE, LINE)
        text(s, x + Inches(1.42), yy + Inches(0.06), Inches(4.0), Inches(0.22), val,
             size=9.5, color=MUTED)
    rect(s, Inches(6.05), Inches(3.50), Inches(1.2), Inches(0.38), BLUE, radius=True)
    text(s, Inches(6.05), Inches(3.56), Inches(1.2), Inches(0.2), "검색", size=10,
         bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    simple_table(
        s, Inches(0.95), Inches(4.25), Inches(11.65),
        ["투어분류", "지역분류", "상품코드", "상품명", "소유사", "예약구분"],
        [
            ["인센티브", "미국/서부", "INC1001", "[인센티브] 상품명 예시", "푸른투어", "선택"],
            ["인센티브", "미국/동부", "INC1002", "[인센티브] 투어상품 예시", "푸른투어", "선택"],
            ["로컬상품", "캐나다/동부", "LOC2100", "[로컬] 캐나다 상품 예시", "푸른투어", "선택"],
        ],
        [0.12, 0.17, 0.13, 0.38, 0.12, 0.08],
    )
    chip(s, Inches(10.95), Inches(1.36), "Copy/CSV/Excel", ORANGE)
    footer(s)


def slide_search_fields(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "03. 검색 조건 입력 방법", "검색 결과는 입력 조건에 따라 상품 목록으로 표시됩니다.", 4, total)
    data = [
        ("투어분류", "상품의 대분류를 선택합니다. 화면 값은 p_type 조건으로 사용됩니다."),
        ("지역분류", "1차 지역(area1)과 2차 지역(area2)을 선택하면 지역 조건이 추가됩니다."),
        ("상품명", "공백 기준 키워드를 나누어 p_name LIKE 검색을 수행합니다."),
        ("출발일", "기본 시작일은 오늘, 종료일은 3개월 후로 설정됩니다."),
        ("검색 버튼", "POST mode=search로 현재 조건을 제출합니다."),
        ("출력 목록", "DataTable 버튼으로 Copy, CSV, Excel, Print를 사용할 수 있습니다."),
    ]
    for i, (title, body) in enumerate(data):
        x = Inches(0.70 + (i % 2) * 6.18)
        y = Inches(1.18 + (i // 2) * 1.28)
        card(s, x, y, Inches(5.75), Inches(1.02), title, body, BLUE)
    rect(s, Inches(0.7), Inches(5.35), Inches(11.95), Inches(0.72),
         RGBColor(255, 249, 231), RGBColor(245, 215, 120), radius=True)
    text(s, Inches(0.95), Inches(5.48), Inches(11.45), Inches(0.34),
         "주의: 상품 선택 전 출발일 조건이 맞는지 확인해야 다음 화면의 여행기간 기준이 흔들리지 않습니다.",
         size=13.5, bold=True, color=RGBColor(124, 88, 20), align=PP_ALIGN.CENTER)
    footer(s)


def slide_select_product(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "04. 상품 목록과 선택 버튼", "목록에서 상품을 고르면 상담예약 입력 화면으로 이동합니다.", 5, total)
    simple_table(
        s, Inches(0.65), Inches(1.2), Inches(12.0),
        ["컬럼", "사용자가 확인할 내용", "코드 기준"],
        [
            ["투어분류", "로컬상품/인바운드/인센티브/아웃바운드 구분", "p_type"],
            ["지역분류", "상품의 1차/2차 지역", "c_code1, c_code2"],
            ["상품코드", "상담과 예약에 연결될 상품 고유 코드", "p_code"],
            ["상품명", "고객에게 안내할 실제 상품명", "p_name"],
            ["소유사", "상품 담당/소유 회사", "p_own"],
            ["선택", "base_conslut_m.php로 이동", "pcode, st 전달"],
        ],
        [0.18, 0.52, 0.30],
    )
    bullets(s, Inches(0.85), Inches(5.0), Inches(11.6), [
        "선택 버튼은 `base_conslut_m.php?...&pcode=상품코드&st=출발일` 형태로 이동합니다.",
        "선택 후에는 첨부 2 화면에서 상품명과 상품코드가 상담예약기본정보에 표시됩니다.",
        "같은 이름의 유사 상품이 있을 수 있으므로 상품코드와 출발일을 같이 확인합니다.",
    ], size=12.3, color=GREEN)
    footer(s)


def slide_form_screen(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "05. 첨부 2 화면: 상담예약 입력", "base_conslut_m.php 화면 구성", 6, total)
    rect(s, Inches(0.55), Inches(1.12), Inches(12.25), Inches(5.55), WHITE, LINE, radius=True)
    text(s, Inches(0.85), Inches(1.35), Inches(11.7), Inches(0.3),
         "상담예약기본정보", size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    fields = [
        ("투어분류", "로컬상품"), ("상품명", "[LAS출발] 4대캐년 3박4일"),
        ("상품코드", "SIN922"), ("상담예약번호", "저장후에 생성"),
        ("여행기간", "시작날짜 ~ 종료날짜"), ("전달직원", "직원 선택"),
        ("방갯수", "0 개"), ("예약인원", "0 명"),
        ("접수일", "오늘 날짜"), ("여행인원", "1 명"),
    ]
    for i, (label, value) in enumerate(fields):
        col = i % 2
        row = i // 2
        x = Inches(0.95 + col * 5.78)
        y = Inches(1.88 + row * 0.58)
        rect(s, x, y, Inches(1.2), Inches(0.34), RGBColor(239, 242, 247), LINE)
        text(s, x, y + Inches(0.05), Inches(1.2), Inches(0.18), label,
             size=8.4, bold=True, align=PP_ALIGN.CENTER)
        rect(s, x + Inches(1.22), y, Inches(4.24), Inches(0.34), WHITE, LINE)
        text(s, x + Inches(1.32), y + Inches(0.05), Inches(4.0), Inches(0.18), value,
             size=8.4, color=MUTED)
    text(s, Inches(0.85), Inches(5.02), Inches(11.7), Inches(0.26),
         "상담정보", size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s, Inches(0.95), Inches(5.38), Inches(2.6), Inches(0.38), WHITE, LINE)
    text(s, Inches(1.05), Inches(5.45), Inches(2.4), Inches(0.18), "이름 / 연락처 / 이메일", size=8.4, color=MUTED)
    rect(s, Inches(3.8), Inches(5.38), Inches(8.7), Inches(0.78), WHITE, LINE)
    text(s, Inches(3.9), Inches(5.48), Inches(8.5), Inches(0.45), "상담내용 입력 영역", size=9.5, color=MUTED)
    footer(s)


def slide_required_input(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "06. 필수 입력 항목", "상담등록 버튼을 누르면 필수값을 먼저 검사합니다.", 7, total)
    data = [
        ("여행기간", "여행시작날짜가 비어 있으면 저장하지 않고 입력 요청 알림을 표시합니다."),
        ("예약자 이름", "r_name 값이 필수입니다. 고객 식별 기준으로 사용됩니다."),
        ("예약자 전화번호", "r_phone 값이 필수입니다. 후속 상담 연락 기준입니다."),
        ("예약자 이메일", "r_email 값이 필수입니다. 견적/안내 발송 기준입니다."),
        ("방갯수/여행인원", "기본값은 방갯수 0개, 여행인원 1명으로 보입니다."),
        ("상담내용", "고객 요청사항, 안내 내용, 후속 처리 메모를 남깁니다."),
    ]
    for i, (title, body) in enumerate(data):
        x = Inches(0.68 + (i % 2) * 6.15)
        y = Inches(1.18 + (i // 2) * 1.18)
        card(s, x, y, Inches(5.75), Inches(0.95), title, body, GREEN)
    bullets(s, Inches(0.9), Inches(5.15), Inches(11.3), [
        "상담예약번호는 신규 저장 전에는 '저장후에 생성'으로 표시됩니다.",
        "기존 상담에서 다시 저장하는 경우 consultCode 기준으로 순번이 이어집니다.",
    ], size=12.3, color=ORANGE)
    footer(s)


def slide_buttons(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "07. 상단 버튼 역할", "상담등록, 상담삭제, 견적등록 버튼의 동작입니다.", 8, total)
    buttons = [
        ("상담등록", "필수값 확인 후 저장 확인창을 띄우고 frmreserve를 제출합니다. mode 기본값은 save입니다.", GREEN),
        ("상담삭제", "삭제 확인 후 mode 값을 delete로 바꾸고 현재 상담을 삭제 요청합니다.", RED),
        ("견적등록", "상담코드, 상품코드, 출발일을 가지고 base_reservation_m.php 견적/예약 화면으로 이동합니다.", ORANGE),
    ]
    for i, (title, body, color) in enumerate(buttons):
        card(s, Inches(0.9 + i * 4.1), Inches(1.55), Inches(3.65), Inches(2.6),
             title, body, color, i + 1)
    rect(s, Inches(1.0), Inches(4.82), Inches(11.3), Inches(0.95),
         RGBColor(238, 245, 255), LINE, radius=True)
    text(s, Inches(1.25), Inches(5.0), Inches(10.8), Inches(0.42),
         "실무 순서: 상담등록으로 고객 문의를 먼저 저장한 뒤, 확정 또는 금액 안내가 필요하면 견적등록을 진행합니다.",
         size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    footer(s)


def slide_save_logic(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "08. 저장 후 처리", "상담정보는 consult_info 테이블에 저장됩니다.", 9, total)
    steps = [
        ("1", "상담코드 생성", "신규 저장 시 C + 날짜/시간 + 순번 형태의 상담코드를 만듭니다."),
        ("2", "상담정보 저장", "register, 전달직원, 예약자, 상품, 여행기간, 인원, 상담내용을 저장합니다."),
        ("3", "저장 완료 이동", "저장 성공 후 base_conslut_mylist.php로 이동합니다."),
        ("4", "삭제 처리", "mode=delete이면 seq_no 기준으로 consult_info 데이터를 삭제합니다."),
    ]
    for i, (no, title, body) in enumerate(steps):
        x = Inches(0.78 + (i % 2) * 6.15)
        y = Inches(1.35 + (i // 2) * 1.55)
        card(s, x, y, Inches(5.75), Inches(1.18), title, body, BLUE, no)
    simple_table(
        s, Inches(0.9), Inches(4.95), Inches(11.55),
        ["저장 필드", "내용"],
        [
            ["member_name / member_phone / member_email", "예약자 이름, 연락처, 이메일"],
            ["p_code / p_name", "선택한 상품코드와 상품명"],
            ["start_date / stop_date / p_cnt / room_cnt", "여행기간, 여행인원, 방갯수"],
        ],
        [0.36, 0.64],
    )
    footer(s)


def slide_history(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "09. 상담이력 확인", "같은 화면 하단에서 기존 상담 이력을 확인합니다.", 10, total)
    simple_table(
        s, Inches(0.75), Inches(1.28), Inches(11.8),
        ["상담코드", "상품명", "예약자", "연락처", "출발일", "담당", "전달직원"],
        [
            ["C260429...", "선택 상품명", "홍길동", "000-0000", "2026-07-29", "관리자", "직원명"],
            ["C260429...", "다른 상담 상품", "김고객", "000-0000", "2026-08-01", "관리자", "직원명"],
            ["C260430...", "후속 상담 상품", "이예약", "000-0000", "2026-08-15", "관리자", "직원명"],
        ],
        [0.16, 0.27, 0.12, 0.14, 0.12, 0.09, 0.10],
    )
    bullets(s, Inches(0.95), Inches(4.2), Inches(11.1), [
        "상담코드는 링크로 표시되며 클릭하면 해당 상담 상세 화면으로 다시 진입합니다.",
        "상담이력은 저장된 상담을 추적하고 중복 상담을 확인하는 용도입니다.",
        "견적등록 전 고객 정보와 상담내용이 올바르게 남아 있는지 확인합니다.",
    ], size=12.5, color=GREEN)
    footer(s)


def slide_checklist(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "10. 저장 전 체크리스트", "실수 방지를 위한 최종 확인 항목입니다.", 11, total)
    left = [
        "선택한 상품명과 상품코드가 고객 문의 상품과 같은가?",
        "여행시작날짜와 종료날짜가 맞는가?",
        "방갯수와 여행인원이 고객 요청과 같은가?",
        "예약자 이름/전화번호/이메일이 모두 입력되었는가?",
    ]
    right = [
        "전달직원이 지정되어 후속 업무 담당자가 명확한가?",
        "상담내용에 고객 요청, 안내사항, 특이사항이 남아 있는가?",
        "저장 후 상담이력 목록에 새 상담이 표시되는가?",
        "예약 또는 견적 업무가 필요하면 견적등록으로 이어갔는가?",
    ]
    card(s, Inches(0.75), Inches(1.25), Inches(5.9), Inches(4.9),
         "입력값 확인", "", BLUE)
    bullets(s, Inches(1.05), Inches(1.95), Inches(5.25), left, size=12.2, color=BLUE)
    card(s, Inches(6.9), Inches(1.25), Inches(5.65), Inches(4.9),
         "후속처리 확인", "", GREEN)
    bullets(s, Inches(7.2), Inches(1.95), Inches(5.05), right, size=12.2, color=GREEN)
    footer(s)


def slide_faq(total):
    s = prs.slides.add_slide(BLANK)
    header(s, "11. FAQ / 오류 대응", "사용 중 자주 확인할 내용입니다.", 12, total)
    faqs = [
        ("검색 결과가 안 보일 때", "투어분류, 지역분류, 출발일 범위를 넓혀 다시 검색합니다."),
        ("상품 선택 후 값이 다를 때", "상품코드와 출발일을 다시 확인하고 목록 화면에서 재선택합니다."),
        ("상담등록이 안 될 때", "여행기간, 예약자 이름, 전화번호, 이메일 필수값을 먼저 확인합니다."),
        ("상담을 예약으로 넘길 때", "상담등록으로 저장한 뒤 견적등록 버튼을 사용합니다."),
        ("삭제가 필요한 때", "상담삭제는 현재 상담 데이터를 삭제하므로 상담코드와 내용을 확인 후 실행합니다."),
    ]
    y = Inches(1.15)
    for i, (q, a) in enumerate(faqs, start=1):
        card(s, Inches(0.75), y, Inches(11.8), Inches(0.78), q, a,
             ORANGE if i < 5 else RED, i)
        y += Inches(0.92)
    footer(s)


def slide_sources(total):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    text(s, Inches(0.85), Inches(1.25), Inches(11.8), Inches(0.78),
         "사용자 매뉴얼 생성 완료", size=38, bold=True, color=WHITE)
    text(s, Inches(0.9), Inches(2.25), Inches(11.4), Inches(0.42),
         "두 첨부 화면 기준으로 상품 검색부터 상담 저장까지의 업무 흐름을 정리했습니다.",
         size=17, color=RGBColor(208, 216, 238))
    rect(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(1.6), WHITE, LINE, radius=True)
    text(s, Inches(1.15), Inches(3.82), Inches(10.9), Inches(0.35),
         "확인한 파일", size=14, bold=True, color=NAVY)
    text(s, Inches(1.15), Inches(4.28), Inches(10.9), Inches(0.42),
         "D:\\www\\prttour_myprt\\base_consult.php", size=13, color=TEXT)
    text(s, Inches(1.15), Inches(4.72), Inches(10.9), Inches(0.42),
         "D:\\www\\prttour_myprt\\base_conslut_m.php", size=13, color=TEXT)
    text(s, Inches(0.9), Inches(6.8), Inches(11.4), Inches(0.26),
         "푸른투어 인트라넷 · 예약상담등록 사용자 매뉴얼", size=10.5,
         color=RGBColor(208, 216, 238), align=PP_ALIGN.CENTER)


@dataclass
class Preview:
    title: str
    subtitle: str
    rows: Sequence[tuple[str, str]]
    dark: bool = False


def rgb(c: RGBColor):
    return (c[0], c[1], c[2])


def font(size, bold=False):
    path = r"C:/Windows/Fonts/malgunbd.ttf" if bold else r"C:/Windows/Fonts/malgun.ttf"
    if os.path.exists(path):
        return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def wrap(draw, value, fnt, max_w):
    words = str(value).split()
    if not words:
        return [""]
    lines, cur = [], words[0]
    for word in words[1:]:
        nxt = cur + " " + word
        if draw.textbbox((0, 0), nxt, font=fnt)[2] <= max_w:
            cur = nxt
        else:
            lines.append(cur)
            cur = word
    lines.append(cur)
    return lines


def draw_preview(idx, total, spec: Preview):
    W, H = 1920, 1080
    img = Image.new("RGB", (W, H), rgb(BG if not spec.dark else NAVY))
    d = ImageDraw.Draw(img)
    if spec.dark:
        d.ellipse([1400, -260, 2220, 560], fill=rgb(NAVY_DARK))
        d.ellipse([-240, 620, 420, 1280], fill=rgb(NAVY_DARK))
        d.text((118, 240), spec.title, fill=rgb(WHITE), font=font(72, True))
        d.text((125, 360), spec.subtitle, fill=(208, 216, 238), font=font(30))
    else:
        d.rectangle([0, 0, W, 126], fill=rgb(NAVY))
        d.text((64, 26), spec.title, fill=rgb(WHITE), font=font(48, True))
        d.text((68, 84), spec.subtitle, fill=(208, 216, 238), font=font(19))
        d.text((1690, 46), f"{idx} / {total}", fill=rgb(WHITE), font=font(18, True))
        y = 170
        for i, (head, body) in enumerate(spec.rows):
            x = 90 + (i % 2) * 870
            if i and i % 2 == 0:
                y += 220
            d.rounded_rectangle([x, y, x + 790, y + 170], radius=10,
                                fill=rgb(WHITE), outline=rgb(LINE), width=2)
            d.rectangle([x, y, x + 96, y + 170], fill=rgb(BLUE))
            d.text((x + 34, y + 54), str(i + 1), fill=rgb(WHITE), font=font(36, True))
            d.text((x + 128, y + 24), head, fill=rgb(TEXT), font=font(27, True))
            yy = y + 74
            for line in wrap(d, body, font(22), 610)[:3]:
                d.text((x + 128, yy), line, fill=rgb(MUTED), font=font(22))
                yy += 32
        d.text((1260, 1038), "출처: base_consult.php, base_conslut_m.php",
               fill=rgb(MUTED), font=font(16))
    out = f"{PREVIEW_PREFIX}{idx:02d}.png"
    img.save(out)
    return out


def render_previews(total):
    specs = [
        Preview("예약상담등록 사용자 매뉴얼", "상품 검색/선택 화면 + 상담예약 입력 화면", [], True),
        Preview("01. 업무 흐름", "예약상담등록은 상품 선택 후 상담 정보를 저장하는 순서입니다.", [
            ("메뉴 진입", "예약상담관리 > 예약상담등록"), ("상품 검색", "투어분류/지역/상품명/출발일"),
            ("상품 선택", "선택 버튼으로 입력 화면 이동"), ("상담 입력", "여행기간/예약자/상담내용")
        ]),
        Preview("02. 첨부 1 화면", "base_consult.php 상품 검색/선택", [
            ("검색 조건", "투어분류, 지역분류, 상품명, 출발일"), ("검색 버튼", "조건 제출 후 목록 갱신"),
            ("상품 목록", "투어분류/지역/상품코드/상품명"), ("선택", "상담 입력 화면으로 이동")
        ]),
        Preview("03. 검색 조건 입력 방법", "검색 결과는 조건에 따라 표시됩니다.", [
            ("투어분류", "p_type 조건"), ("지역분류", "area1, area2 조건"),
            ("상품명", "키워드 LIKE 검색"), ("출발일", "오늘부터 3개월 기본값")
        ]),
        Preview("04. 상품 목록과 선택 버튼", "목록에서 상품을 고르면 상담예약 입력 화면으로 이동합니다.", [
            ("상품코드", "p_code 확인"), ("상품명", "고객 문의 상품 확인"),
            ("선택 버튼", "pcode와 st 전달"), ("주의", "상품코드와 출발일 동시 확인")
        ]),
        Preview("05. 첨부 2 화면", "base_conslut_m.php 상담예약 입력", [
            ("기본정보", "투어분류/상품명/상품코드"), ("여행기간", "시작날짜와 종료날짜"),
            ("인원", "방갯수/예약인원/여행인원"), ("상담정보", "이름/연락처/이메일/내용")
        ]),
        Preview("06. 필수 입력 항목", "상담등록 버튼은 필수값을 먼저 검사합니다.", [
            ("여행기간", "시작날짜 필수"), ("이름", "예약자 이름 필수"),
            ("전화번호", "예약자 연락처 필수"), ("이메일", "예약자 이메일 필수")
        ]),
        Preview("07. 상단 버튼 역할", "상담등록, 상담삭제, 견적등록", [
            ("상담등록", "mode=save 제출"), ("상담삭제", "mode=delete 제출"),
            ("견적등록", "base_reservation_m.php 이동"), ("실무 순서", "저장 후 견적등록")
        ]),
        Preview("08. 저장 후 처리", "상담정보는 consult_info 테이블에 저장됩니다.", [
            ("상담코드", "신규 저장 시 생성"), ("저장 필드", "예약자/상품/기간/인원/내용"),
            ("완료 이동", "base_conslut_mylist.php"), ("삭제", "seq_no 기준 삭제")
        ]),
        Preview("09. 상담이력 확인", "하단 상담이력 목록", [
            ("상담코드", "상세 재진입 링크"), ("상품명", "상담 상품 확인"),
            ("예약자", "고객 식별"), ("전달직원", "후속 담당 확인")
        ]),
        Preview("10. 저장 전 체크리스트", "실수 방지를 위한 최종 확인", [
            ("상품 확인", "상품명/상품코드"), ("기간 확인", "시작/종료일"),
            ("연락처", "이름/전화/이메일"), ("후속처리", "견적등록 필요 여부")
        ]),
        Preview("11. FAQ / 오류 대응", "자주 확인할 내용", [
            ("검색 안 됨", "조건 범위 확대"), ("등록 안 됨", "필수값 확인"),
            ("값이 다름", "상품 재선택"), ("삭제", "상담코드 확인 후 실행")
        ]),
        Preview("사용자 매뉴얼 생성 완료", "확인 파일: base_consult.php, base_conslut_m.php", [], True),
    ]
    return [draw_preview(i + 1, total, spec) for i, spec in enumerate(specs)]


def build():
    total = 13
    slide_cover(total)
    slide_flow(total)
    slide_search_screen(total)
    slide_search_fields(total)
    slide_select_product(total)
    slide_form_screen(total)
    slide_required_input(total)
    slide_buttons(total)
    slide_save_logic(total)
    slide_history(total)
    slide_checklist(total)
    slide_faq(total)
    slide_sources(total)
    prs.save(OUT)
    previews = render_previews(total)
    print("SAVED:", OUT)
    print("PREVIEWS:", len(previews))
    for path in previews:
        print(path)


if __name__ == "__main__":
    build()
