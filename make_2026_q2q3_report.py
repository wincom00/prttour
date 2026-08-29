# -*- coding: utf-8 -*-
"""
2026년 2·3분기 경영실적 PPTX 생성기

기존 2026_3Q_경영실적.pptx 와 동일한 디자인 규격(색/폰트/카드/표/차트)을 그대로 쓰고,
내용만 2분기·3분기 두 분기 + 출발일(stDate)/예약일(revDate) 기준 비교로 재구성했다.

집계 기준 (기존 분기 리포트와 동일)
  - reserve_info : parent='MAIN', rev_status IN ('READY','DONE')
  - 매출 last_total / 미수 last_bal / 인원 p_cnt
  - 수금(현금흐름)은 payment_history 결제일(wdate) 기준, pay_method='init' 및 취소건 제외

실행:  C:/Windows/py.exe make_2026_q2q3_report.py
"""
import os
import re
import html
import subprocess
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
MYSQL = r"C:\laragon\bin\mysql\mysql-8.4.3-winx64\bin\mysql.exe"
OUT = os.path.join(BASE_DIR, "2026_2Q3Q_경영실적.pptx")

ASOF = datetime.date.today().strftime("%Y-%m-%d")

# ── 디자인 상수 (3Q 덱에서 추출) ────────────────────────────────
NAVY = "1B335C"
BLUE = "2E6BC6"
PALE = "B5C6DC"
TEAL = "2A9D8F"
CARD = "F2F5F9"
GRAY = "6B7280"
RED = "C0392B"
GREEN = "1E8E3E"
WHITE = "FFFFFF"
TITLE_SUB = "9AB8E0"
TITLE_SUB2 = "C8D8F0"
FONT = "맑은 고딕"


# ── DB ──────────────────────────────────────────────────────────
def db_creds():
    """dbconn.php 에서 접속정보를 읽어 온다 (비밀번호 사본을 만들지 않기 위해)."""
    src = open(os.path.join(BASE_DIR, "include", "dbconn.php"), encoding="utf-8").read()
    # 파일 상단에 주석 처리된 예전 접속정보가 남아 있어서 마지막 대입값을 쓴다
    g = lambda k: re.findall(r'\$%s\s*=\s*"([^"]+)"' % k, src)[-1]
    host = g("db_host")
    port = "3306"
    if ":" in host:
        host, port = host.split(":", 1)
    return host, port, g("db_user"), g("db_passwd"), g("db_name")


HOST, PORT, USER, PW, DBN = db_creds()


def q(sql):
    """mysql.exe 배치모드(TSV) → dict 리스트."""
    p = subprocess.run(
        [MYSQL, "-h" + HOST, "-P" + PORT, "-u" + USER, "-p" + PW, DBN,
         "--default-character-set=utf8mb4", "-B", "-e", sql],
        capture_output=True)
    out = p.stdout.decode("utf-8", "replace")
    err = p.stderr.decode("utf-8", "replace")
    if p.returncode != 0:
        raise RuntimeError(err + "\n-- SQL --\n" + sql)
    lines = [l for l in out.replace("\r\n", "\n").split("\n") if l != ""]
    if not lines:
        return []
    head = lines[0].split("\t")
    rows = []
    for l in lines[1:]:
        cells = l.split("\t")
        rows.append({head[i]: (None if cells[i] == "NULL" else cells[i])
                     for i in range(len(head))})
    return rows


def f(v):
    return float(v or 0)


def i(v):
    return int(float(v or 0))


# ── 데이터 수집 ─────────────────────────────────────────────────
RSV = "parent='MAIN' AND rev_status IN ('READY','DONE')"

print("[1/9] 분기 집계(출발일)")
_qs = q("""
SELECT YEAR(stDate) y, QUARTER(stDate) q, COUNT(*) cnt, SUM(p_cnt) pax,
       SUM(last_total) amt, SUM(last_bal) bal
  FROM reserve_info
 WHERE %s AND YEAR(stDate) IN (2025,2026) AND QUARTER(stDate) IN (2,3)
 GROUP BY 1,2""" % RSV)
ST = {(i(r["y"]), i(r["q"])): r for r in _qs}

print("[2/9] 분기 집계(예약일)")
_qr = q("""
SELECT YEAR(revDate) y, QUARTER(revDate) q, COUNT(*) cnt, SUM(p_cnt) pax,
       SUM(last_total) amt
  FROM reserve_info
 WHERE %s AND YEAR(revDate) IN (2025,2026) AND QUARTER(revDate) IN (2,3)
 GROUP BY 1,2""" % RSV)
RV = {(i(r["y"]), i(r["q"])): r for r in _qr}

print("[3/9] 월별(출발일)")
_ms = q("""
SELECT YEAR(stDate) y, MONTH(stDate) m, SUM(last_total) amt, SUM(p_cnt) pax
  FROM reserve_info
 WHERE %s AND YEAR(stDate) IN (2025,2026) AND MONTH(stDate) BETWEEN 4 AND 9
 GROUP BY 1,2""" % RSV)
MST = {(i(r["y"]), i(r["m"])): r for r in _ms}

print("[4/9] 월별(예약일)")
_mr = q("""
SELECT YEAR(revDate) y, MONTH(revDate) m, SUM(last_total) amt, SUM(p_cnt) pax
  FROM reserve_info
 WHERE %s AND YEAR(revDate) IN (2025,2026) AND MONTH(revDate) BETWEEN 4 AND 9
 GROUP BY 1,2""" % RSV)
MRV = {(i(r["y"]), i(r["m"])): r for r in _mr}

print("[5/9] 취소율")
_cx = q("""
SELECT YEAR(stDate) y, QUARTER(stDate) q,
       SUM(rev_status='CANCEL') c, COUNT(*) t
  FROM reserve_info
 WHERE parent='MAIN' AND YEAR(stDate) IN (2025,2026) AND QUARTER(stDate) IN (2,3)
 GROUP BY 1,2""")
CX = {(i(r["y"]), i(r["q"])): (f(r["c"]) / f(r["t"]) * 100 if f(r["t"]) else 0) for r in _cx}

print("[6/9] 수금(payment_history)")
PAYM = q("""
SELECT DATE_FORMAT(wdate,'%%y-%%m') ym, SUM(payment) amt
  FROM payment_history
 WHERE payment_status NOT IN ('RRQUEST','CANCEL') AND pay_method <> 'init'
   AND wdate >= '2025-09-01' AND wdate < '%s'
 GROUP BY 1 ORDER BY 1""" % (datetime.date.today() + datetime.timedelta(days=1)))

PAYQ = q("""
SELECT QUARTER(wdate) q, SUM(payment) amt
  FROM payment_history
 WHERE payment_status NOT IN ('RRQUEST','CANCEL') AND pay_method <> 'init'
   AND YEAR(wdate)=2026 AND QUARTER(wdate) IN (2,3)
 GROUP BY 1""")
PAYQ = {i(r["q"]): f(r["amt"]) for r in PAYQ}

PAYMTD = q("""
SELECT pay_method, SUM(payment) amt
  FROM payment_history
 WHERE payment_status NOT IN ('RRQUEST','CANCEL') AND pay_method <> 'init'
   AND wdate >= '2026-04-01'
 GROUP BY 1 ORDER BY 2 DESC LIMIT 5""")

print("[7/9] 상품 TOP10")
TOP = q("""
SELECT p_name, COUNT(*) cnt, SUM(p_cnt) pax, SUM(last_total) amt
  FROM reserve_info
 WHERE %s AND YEAR(stDate)=2026 AND QUARTER(stDate) IN (2,3)
 GROUP BY p_code, p_name ORDER BY amt DESC LIMIT 10""" % RSV)

print("[8/9] 예약경로 · 예약유형")
PATH = q("""
SELECT COALESCE(NULLIF(c.comment,''),'기타') nm, SUM(r.last_total) amt
  FROM reserve_info r
  LEFT JOIN code_base c
    ON CONCAT(c.lvcode1,c.lvcode2,c.lvcode3) = r.r_path
 WHERE %s AND YEAR(r.stDate)=2026 AND QUARTER(r.stDate) IN (2,3)
 GROUP BY 1 ORDER BY amt DESC LIMIT 8""" % RSV)

TTYPE = q("""
SELECT tour_type, COUNT(*) cnt, SUM(last_total) amt
  FROM reserve_info
 WHERE %s AND YEAR(stDate)=2026 AND QUARTER(stDate) IN (2,3)
 GROUP BY 1""" % RSV)
TTNAME = {"1": "직접예약", "2": "웹예약", "3": "업체예약"}

print("[9/9] 지역별")
AREA = q("""
SELECT COALESCE(NULLIF(c.comment,''),'기타') nm, SUM(r.last_total) amt
  FROM reserve_info r
  LEFT JOIN code_base c
    ON CONCAT(c.lvcode1,c.lvcode2,c.lvcode3) = r.s_area
 WHERE %s AND YEAR(r.stDate)=2026 AND QUARTER(r.stDate) IN (2,3)
 GROUP BY 1 ORDER BY amt DESC LIMIT 8""" % RSV)


# ── 파생값 ──────────────────────────────────────────────────────
def st(y, qq, k):
    r = ST.get((y, qq))
    return f(r[k]) if r else 0.0


def rv(y, qq, k):
    r = RV.get((y, qq))
    return f(r[k]) if r else 0.0


def mst(y, m, k):
    r = MST.get((y, m))
    return f(r[k]) if r else 0.0


def mrv(y, m, k):
    r = MRV.get((y, m))
    return f(r[k]) if r else 0.0


# 2·3분기 합계
def hy(y, k):
    return st(y, 2, k) + st(y, 3, k)


H26 = {k: hy(2026, k) for k in ("cnt", "pax", "amt", "bal")}
H25 = {k: hy(2025, k) for k in ("cnt", "pax", "amt", "bal")}
H26["rec"] = H26["amt"] - H26["bal"]
H25["rec"] = H25["amt"] - H25["bal"]
H26["aov"] = H26["amt"] / H26["pax"] if H26["pax"] else 0
H25["aov"] = H25["amt"] / H25["pax"] if H25["pax"] else 0


# ── 포맷 헬퍼 ───────────────────────────────────────────────────
def money(v):
    return "$" + format(int(round(v)), ",")


def num(v, unit=""):
    return format(int(round(v)), ",") + unit


def delta(cur, prev):
    """(문구, 색) — 전년 대비 증감."""
    if not prev:
        return "-", GRAY
    d = (cur - prev) / prev * 100
    mark = "▲" if d >= 0 else "▼"
    return "%s %.1f%%" % (mark, abs(d)), (GREEN if d >= 0 else RED)


def delta_pp(cur, prev):
    d = cur - prev
    return "%+.1f%%p" % d, (RED if d >= 0 else GREEN)


# ── PPTX 빌더 ───────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def txt(sl, l, t, w, h, text, sz=11, bold=False, color=GRAY,
        align=PP_ALIGN.LEFT, spacing=None):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for n, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.name = FONT
        r.font.color.rgb = RGBColor.from_string(color)
    return tb


def rect(sl, l, t, w, h, color, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = sl.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if rounded:
        sp.adjustments[0] = 0.06
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def head(sl, title, sub=None):
    rect(sl, 0.5, 0.45, 0.09, 0.42, BLUE)
    txt(sl, 0.75, 0.42, 9.0, 0.5, title, sz=24, bold=True, color=NAVY)
    if sub:
        txt(sl, 0.75, 0.95, 11.5, 0.35, sub, sz=11, color=GRAY)
    txt(sl, 11.6, 6.95, 1.4, 0.3, "푸른투어", sz=9, color=GRAY, align=PP_ALIGN.RIGHT)


def kpi(sl, l, t, label, value, sub=None, subcolor=GRAY, w=3.85, h=1.5):
    rect(sl, l, t, w, h, CARD, rounded=True)
    txt(sl, l + 0.3, t + 0.2, w - 0.55, 0.3, label, sz=12, color=GRAY)
    txt(sl, l + 0.3, t + 0.55, w - 0.55, 0.5, value, sz=26, bold=True, color=NAVY)
    if sub:
        txt(sl, l + 0.3, t + 1.1, w - 0.55, 0.3, sub, sz=11, bold=True, color=subcolor)


def table(sl, l, t, w, h, rows, widths, aligns, colors=None, hdr_sz=12, body_sz=12):
    """rows[0] = 헤더. colors[r][c] 로 개별 셀 글자색 지정(없으면 NAVY)."""
    gt = sl.shapes.add_table(len(rows), len(rows[0]),
                             Inches(l), Inches(t), Inches(w), Inches(h)).table
    for c, cw in enumerate(widths):
        gt.columns[c].width = Inches(cw)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = RGBColor.from_string(NAVY)
            else:
                cell.fill.fore_color.rgb = RGBColor.from_string(CARD if r % 2 == 0 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.alignment = aligns[c]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(hdr_sz if r == 0 else body_sz)
            run.font.bold = (r == 0 or c == 0)
            run.font.name = FONT
            if r == 0:
                run.font.color.rgb = RGBColor.from_string(WHITE)
            else:
                cc = (colors or {}).get((r, c), NAVY)
                run.font.color.rgb = RGBColor.from_string(cc)
    return gt


def bar(sl, l, t, w, h, cats, series, horizontal=False, legend=False,
        label_sz=10, cat_sz=11):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals, _color in series:
        cd.add_series(name, vals)
    ct = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = sl.shapes.add_chart(ct, Inches(l), Inches(t), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.font.name = FONT
    ch.font.size = Pt(cat_sz)
    ch.font.color.rgb = RGBColor.from_string(NAVY)
    for idx, (_n, _v, color) in enumerate(series):
        s = ch.plots[0].series[idx]
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = RGBColor.from_string(color)
    pl = ch.plots[0]
    pl.gap_width = 150
    pl.has_data_labels = True
    dl = pl.data_labels
    dl.number_format = '\\$#,##0,"K"'
    dl.number_format_is_linked = False
    dl.font.size = Pt(label_sz)
    dl.font.name = FONT
    dl.font.color.rgb = RGBColor.from_string(NAVY)
    if legend:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(cat_sz)
        ch.legend.font.name = FONT
    else:
        ch.has_legend = False
    return ch


# ── 1. 표지 ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, NAVY)
txt(s, 1.2, 2.5, 11.0, 0.6, "푸른투어", sz=16, color=TITLE_SUB)
txt(s, 1.2, 3.0, 11.0, 1.0, "2026년 2·3분기 경영실적", sz=44, bold=True, color=WHITE)
txt(s, 1.2, 4.1, 11.0, 0.5, "2026.04.01 ~ 2026.09.30 · 출발일(stDate) / 예약일(revDate) 기준",
    sz=18, color=TITLE_SUB2)
rect(s, 1.2, 4.8, 2.2, 0.04, BLUE)
txt(s, 1.2, 5.2, 11.0, 0.8,
    "작성기준일 %s  ·  자료출처 ERP(prtadmindb) reserve_info / payment_history" % ASOF,
    sz=11, color=TITLE_SUB)

# ── 2. 실적 요약 (2·3분기 합계) ─────────────────────────────────
s = prs.slides.add_slide(BLANK)
head(s, "실적 요약", "출발일 기준 확정예약(READY/DONE) · 2·3분기 합계 · 전년동기 대비")
cards = [
    ("매출액", money(H26["amt"]), delta(H26["amt"], H25["amt"])),
    ("송객인원", num(H26["pax"], "명"), delta(H26["pax"], H25["pax"])),
    ("예약건수", num(H26["cnt"], "건"), delta(H26["cnt"], H25["cnt"])),
    ("1인 객단가", money(H26["aov"]), delta(H26["aov"], H25["aov"])),
    ("수금액(누계)", money(H26["rec"]), delta(H26["rec"], H25["rec"])),
    ("미수잔액", money(H26["bal"]), delta(H26["bal"], H25["bal"])),
]
for n, (lab, val, (dtxt, dcol)) in enumerate(cards):
    kpi(s, 0.55 + (n % 3) * 4.15, 1.55 + (n // 3) * 1.75, lab, val,
        "전년동기 " + dtxt, dcol)

rect(s, 0.55, 5.25, 12.2, 1.6, CARD, rounded=True)
txt(s, 0.85, 5.45, 11.6, 1.3,
    "· 2·3분기 매출 %s / 인원 %s — 전년동기 대비 매출 %s, 인원 %s\n"
    "· 1인 객단가 %s (전년동기 %s) — 인원 감소분을 단가가 거의 그대로 상쇄\n"
    "· 미수잔액 %s (매출대비 %.1f%%) · 수금완료율 %.1f%% · 취소율 Q2 %.1f%% / Q3 %.1f%%\n"
    "· 9월 출발분은 예약 진행 중으로 분기 마감 시 증가 여지 있음 (%s 기준)" % (
        money(H26["amt"]), num(H26["pax"], "명"),
        delta(H26["amt"], H25["amt"])[0], delta(H26["pax"], H25["pax"])[0],
        money(H26["aov"]), money(H25["aov"]),
        money(H26["bal"]), H26["bal"] / H26["amt"] * 100 if H26["amt"] else 0,
        H26["rec"] / H26["amt"] * 100 if H26["amt"] else 0,
        CX.get((2026, 2), 0), CX.get((2026, 3), 0),
        ASOF),
    sz=12, color=NAVY, spacing=1.35)

# ── 3. 분기 실적 비교 ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
head(s, "분기 실적 비교", "Q2'26 / Q3'26 · 각 전년동기 대비")


def qmet(y, qq):
    amt, pax, cnt, bal = (st(y, qq, "amt"), st(y, qq, "pax"),
                          st(y, qq, "cnt"), st(y, qq, "bal"))
    return {
        "매출액": (amt, money),
        "송객인원": (pax, lambda v: num(v, "명")),
        "예약건수": (cnt, lambda v: num(v, "건")),
        "1인 객단가": (amt / pax if pax else 0, money),
        "수금액": (amt - bal, money),
        "미수잔액": (bal, money),
        "신규수주(예약일)": (rv(y, qq, "amt"), money),
    }


rows = [["구분", "Q2'26", "Q2'25", "전년비", "Q3'26", "Q3'25", "전년비"]]
cellcolor = {}
labels = ["매출액", "송객인원", "예약건수", "1인 객단가", "수금액", "미수잔액", "신규수주(예약일)"]
m226, m225, m326, m325 = qmet(2026, 2), qmet(2025, 2), qmet(2026, 3), qmet(2025, 3)
for n, lab in enumerate(labels):
    v26, fmt = m226[lab]
    v25, _ = m225[lab]
    w26, _ = m326[lab]
    w25, _ = m325[lab]
    d2, c2 = delta(v26, v25)
    d3, c3 = delta(w26, w25)
    rows.append([lab, fmt(v26), fmt(v25), d2, fmt(w26), fmt(w25), d3])
    cellcolor[(n + 1, 3)] = c2
    cellcolor[(n + 1, 6)] = c3
# 취소율
d2, c2 = delta_pp(CX.get((2026, 2), 0), CX.get((2025, 2), 0))
d3, c3 = delta_pp(CX.get((2026, 3), 0), CX.get((2025, 3), 0))
rows.append(["취소율", "%.1f%%" % CX.get((2026, 2), 0), "%.1f%%" % CX.get((2025, 2), 0), d2,
             "%.1f%%" % CX.get((2026, 3), 0), "%.1f%%" % CX.get((2025, 3), 0), d3])
cellcolor[(len(rows) - 1, 3)] = c2
cellcolor[(len(rows) - 1, 6)] = c3

A = PP_ALIGN
table(s, 0.55, 1.45, 12.2, 2.9, rows,
      [2.3, 1.75, 1.75, 1.3, 1.75, 1.75, 1.3],
      [A.LEFT, A.RIGHT, A.RIGHT, A.RIGHT, A.RIGHT, A.RIGHT, A.RIGHT],
      cellcolor, body_sz=11)
txt(s, 0.55, 6.50, 12.2, 0.6,
    "* 매출·인원·미수는 출발일(stDate) 기준 / 신규수주는 예약일(revDate) 기준 · "
    "취소율은 전체 예약 대비 CANCEL 건수 비중", sz=10, color=GRAY)

# ── 4. 기준별 매출 비교 (stDate vs revDate) ─────────────────────
s = prs.slides.add_slide(BLANK)
head(s, "집계기준별 매출 비교", "같은 예약을 출발일(stDate) / 예약일(revDate) 어느 날짜에 붙이는가에 따른 차이")
months = [4, 5, 6, 7, 8, 9]
bar(s, 0.55, 1.50, 7.4, 4.6,
    ["%d월" % m for m in months],
    [("출발일 기준", [mst(2026, m, "amt") for m in months], BLUE),
     ("예약일 기준", [mrv(2026, m, "amt") for m in months], PALE)],
    legend=True)

rect(s, 8.25, 1.50, 4.5, 4.6, CARD, rounded=True)
txt(s, 8.55, 1.75, 3.9, 0.3, "분기 합계", sz=12, bold=True, color=NAVY)
yy = 2.25
for qq in (2, 3):
    txt(s, 8.55, yy, 3.9, 0.3, "%d분기" % qq, sz=13, bold=True, color=NAVY)
    txt(s, 8.55, yy + 0.35, 2.3, 0.3, "출발일 기준", sz=11, color=GRAY)
    txt(s, 10.55, yy + 0.35, 1.95, 0.3, money(st(2026, qq, "amt")), sz=11, bold=True,
        color=NAVY, align=A.RIGHT)
    txt(s, 8.55, yy + 0.67, 2.3, 0.3, "예약일 기준", sz=11, color=GRAY)
    txt(s, 10.55, yy + 0.67, 1.95, 0.3, money(rv(2026, qq, "amt")), sz=11, bold=True,
        color=NAVY, align=A.RIGHT)
    gap = rv(2026, qq, "amt") - st(2026, qq, "amt")
    txt(s, 8.55, yy + 0.99, 2.3, 0.3, "차이", sz=11, color=GRAY)
    txt(s, 10.55, yy + 0.99, 1.95, 0.3, ("%s%s" % ("+" if gap >= 0 else "-", money(abs(gap)))),
        sz=11, bold=True, color=(GREEN if gap >= 0 else RED), align=A.RIGHT)
    yy += 1.65
txt(s, 8.55, 5.6, 3.9, 0.4,
    "예약일 기준은 그 기간에 '판' 금액,\n출발일 기준은 그 기간에 '나간' 금액이다.",
    sz=10, color=GRAY, spacing=1.2)
txt(s, 0.55, 6.35, 12.2, 0.5,
    "* 예약일 기준 9월은 %s 현재 아직 판매 구간이 시작되지 않아 사실상 0 — 3분기 예약일 기준 수치는 7~8월 두 달치로 봐야 한다."
    % ASOF, sz=10, color=GRAY)

# ── 5. 월별 매출 추이 ───────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
head(s, "월별 매출 추이", "2·3분기 월별 매출(출발일 기준) · 전년 동월 대비")
bar(s, 0.55, 1.50, 7.4, 4.6,
    ["%d월" % m for m in months],
    [("2025년", [mst(2025, m, "amt") for m in months], PALE),
     ("2026년", [mst(2026, m, "amt") for m in months], BLUE)],
    legend=True)

rect(s, 8.25, 1.50, 4.5, 4.6, CARD, rounded=True)
txt(s, 8.55, 1.70, 3.9, 0.3, "월별 상세 (2026년)", sz=12, bold=True, color=NAVY)
yy = 2.10
for m in months:
    cur, prv = mst(2026, m, "amt"), mst(2025, m, "amt")
    d, c = delta(cur, prv)
    txt(s, 8.55, yy, 1.4, 0.3, "%d월" % m, sz=12, bold=True, color=NAVY)
    txt(s, 9.60, yy, 1.7, 0.3, money(cur), sz=12, bold=True, color=NAVY, align=A.RIGHT)
    txt(s, 11.40, yy, 1.1, 0.3, d, sz=11, bold=True, color=c, align=A.RIGHT)
    txt(s, 8.55, yy + 0.28, 3.9, 0.3,
        "   인원 %s (전년 %s)" % (num(mst(2026, m, "pax"), "명"), num(mst(2025, m, "pax"), "명")),
        sz=10, color=GRAY)
    yy += 0.62
txt(s, 8.55, 5.85, 3.9, 0.6,
    "2·3분기 합계\n%s (전년 %s)" % (money(H26["amt"]), money(H25["amt"])),
    sz=13, bold=True, color=NAVY, spacing=1.3)
txt(s, 0.55, 6.40, 12.2, 0.5,
    "* 9월분은 출발 전 예약 포함(진행 중), 분기 마감 시 변동 가능", sz=10, color=GRAY)

# ── 6. 수금 및 미수 현황 ────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
head(s, "수금 및 미수 현황", "수금은 결제일(wdate) 기준 실입금, 미수는 출발일 기준 잔액")
bar(s, 0.55, 1.50, 8.0, 4.4,
    [r["ym"] for r in PAYM],
    [("월 수금액", [f(r["amt"]) for r in PAYM], BLUE)], label_sz=9, cat_sz=10)

cash23 = PAYQ.get(2, 0) + PAYQ.get(3, 0)
boxes = [
    ("2·3분기 수금액 (4/1~%d/%d)" % (datetime.date.today().month, datetime.date.today().day),
     money(cash23)),
    ("2분기 / 3분기 수금액", "%s / %s" % (money(PAYQ.get(2, 0)), money(PAYQ.get(3, 0)))),
    ("2·3분기 미수잔액", money(H26["bal"])),
    ("매출 대비 미수율", "%.1f%%" % (H26["bal"] / H26["amt"] * 100 if H26["amt"] else 0)),
    ("수금 완료율", "%.1f%%" % (H26["rec"] / H26["amt"] * 100 if H26["amt"] else 0)),
]
yy = 1.50
for lab, val in boxes:
    rect(s, 8.85, yy, 3.9, 0.78, CARD, rounded=True)
    txt(s, 9.10, yy + 0.12, 3.4, 0.3, lab, sz=10, color=GRAY)
    txt(s, 9.10, yy + 0.38, 3.4, 0.35, val, sz=15, bold=True, color=NAVY)
    yy += 0.90

txt(s, 0.55, 6.05, 8.0, 0.3, "2·3분기 결제수단별 수금", sz=11, bold=True, color=NAVY)
txt(s, 0.55, 6.35, 8.0, 0.4,
    "   ·   ".join("%s %s" % (r["pay_method"] or "기타", money(f(r["amt"]))) for r in PAYMTD),
    sz=10, color=GRAY)
txt(s, 8.85, 6.35, 3.9, 0.4, "* 최근월 수금은 %s까지 집계" % ASOF, sz=10, color=GRAY,
    align=A.RIGHT)

# ── 7. 상품별 매출 TOP 10 ───────────────────────────────────────
s = prs.slides.add_slide(BLANK)
head(s, "상품별 매출 TOP 10", "2026년 2·3분기 출발 기준")
rows = [["순위", "상품명", "예약건수", "인원", "매출액"]]
for n, r in enumerate(TOP):
    # p_name 에는 목록 화면용 색상 마크업(<span style=…><strong>…)이 섞여 있다
    nm = html.unescape(re.sub(r"<[^>]+>", "", r["p_name"] or "")).strip()
    nm = re.sub(r"\s+", " ", nm)
    if len(nm) > 42:
        nm = nm[:42] + "…"
    rows.append([n + 1, nm, num(f(r["cnt"])), num(f(r["pax"])), money(f(r["amt"]))])
table(s, 0.55, 1.50, 12.2, 3.12, rows,
      [0.9, 6.9, 1.4, 1.2, 1.8],
      [A.CENTER, A.LEFT, A.RIGHT, A.RIGHT, A.RIGHT], body_sz=11)
top_sum = sum(f(r["amt"]) for r in TOP)
txt(s, 0.55, 6.75, 12.2, 0.4,
    "* TOP10 합계 %s — 2·3분기 매출의 %.1f%%" % (
        money(top_sum), top_sum / H26["amt"] * 100 if H26["amt"] else 0),
    sz=10, color=GRAY)

# ── 8. 예약경로 · 예약유형별 매출 ───────────────────────────────
s = prs.slides.add_slide(BLANK)
head(s, "예약경로 · 예약유형별 매출", "2026년 2·3분기 출발 기준")
pl = list(reversed(PATH))
bar(s, 0.55, 1.60, 7.3, 4.6, [r["nm"] for r in pl],
    [("매출", [f(r["amt"]) for r in pl], BLUE)], horizontal=True, cat_sz=11)

txt(s, 8.25, 1.60, 4.5, 0.3, "예약유형별", sz=12, bold=True, color=NAVY)
tt = sorted(TTYPE, key=lambda r: -f(r["amt"]))
yy = 2.0
for r in tt:
    rect(s, 8.25, yy, 4.5, 0.95, CARD, rounded=True)
    txt(s, 8.50, yy + 0.14, 2.5, 0.3, TTNAME.get(r["tour_type"], "기타"), sz=12,
        bold=True, color=NAVY)
    txt(s, 8.50, yy + 0.50, 2.6, 0.3, num(f(r["cnt"]), "건"), sz=11, color=GRAY)
    share = f(r["amt"]) / H26["amt"] * 100 if H26["amt"] else 0
    txt(s, 10.40, yy + 0.14, 2.1, 0.3, "%.1f%%" % share, sz=11, color=GRAY, align=A.RIGHT)
    txt(s, 10.40, yy + 0.45, 2.1, 0.4, money(f(r["amt"])), sz=14, bold=True, color=NAVY,
        align=A.RIGHT)
    yy += 1.10
if PATH:
    txt(s, 0.55, 6.50, 12.2, 0.4,
        "* 최대 경로: %s %s (%.1f%%)" % (
            PATH[0]["nm"], money(f(PATH[0]["amt"])),
            f(PATH[0]["amt"]) / H26["amt"] * 100 if H26["amt"] else 0),
        sz=10, color=GRAY)

# ── 9. 고객 지역별 매출 ─────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
head(s, "고객 지역별 매출", "예약 고객 소재지(통계지역분류) 기준 · 2026년 2·3분기")
al = list(reversed(AREA))
bar(s, 0.55, 1.60, 7.3, 4.7, [r["nm"] for r in al],
    [("매출", [f(r["amt"]) for r in al], TEAL)], horizontal=True, cat_sz=11)

rect(s, 8.25, 1.60, 4.5, 4.7, CARD, rounded=True)
txt(s, 8.55, 1.85, 3.9, 0.3, "지역별 구성비", sz=12, bold=True, color=NAVY)
yy = 2.35
for r in AREA:
    amt = f(r["amt"])
    txt(s, 8.55, yy, 2.3, 0.3, r["nm"], sz=11, color=NAVY)
    txt(s, 10.30, yy, 1.3, 0.3, "$%sK" % num(amt / 1000), sz=11, bold=True, color=NAVY,
        align=A.RIGHT)
    txt(s, 11.60, yy, 0.9, 0.3, "%.1f%%" % (amt / H26["amt"] * 100 if H26["amt"] else 0),
        sz=11, color=GRAY, align=A.RIGHT)
    yy += 0.48

# ── 10. 종합 분석 및 유의사항 ───────────────────────────────────
s = prs.slides.add_slide(BLANK)
head(s, "종합 분석 및 유의사항")
path_top2 = sum(f(r["amt"]) for r in PATH[:2]) / H26["amt"] * 100 if PATH and H26["amt"] else 0
blocks = [
    ("실적 요약",
     "· 2·3분기 매출 %s, 전년동기 %s 대비 %s\n"
     "· 송객인원 %s (%s), 1인 객단가 %s (%s)\n"
     "· 분기별로는 Q2 %s / Q3 %s — 3분기는 9월 출발분이 아직 채워지는 중" % (
         money(H26["amt"]), money(H25["amt"]),
         "%+.1f%%" % ((H26["amt"] - H25["amt"]) / H25["amt"] * 100 if H25["amt"] else 0),
         num(H26["pax"], "명"),
         "%+.1f%%" % ((H26["pax"] - H25["pax"]) / H25["pax"] * 100 if H25["pax"] else 0),
         money(H26["aov"]),
         "%+.1f%%" % ((H26["aov"] - H25["aov"]) / H25["aov"] * 100 if H25["aov"] else 0),
         money(st(2026, 2, "amt")), money(st(2026, 3, "amt")))),
    ("수익 구조",
     "· 매출 상위 10개 상품이 2·3분기 매출의 %.1f%% 차지 — 상품 집중도 높음\n"
     "· 예약경로 상위 2개(%s)가 %.1f%% 비중\n"
     "· 수금완료율 %.1f%%, 미수잔액 %s — 출발 전 잔금 회수 관리 필요" % (
         top_sum / H26["amt"] * 100 if H26["amt"] else 0,
         ", ".join(r["nm"] for r in PATH[:2]), path_top2,
         H26["rec"] / H26["amt"] * 100 if H26["amt"] else 0, money(H26["bal"]))),
    ("데이터 기준 및 유의사항",
     "· 집계 기준: reserve_info 중 parent='MAIN', 예약상태 READY/DONE\n"
     "· 출발일(stDate) = 실적 귀속 기준 / 예약일(revDate) = 수주 시점 기준 — 두 기준은 서로 대체가 아니라 보완\n"
     "· 수금액은 payment_history 결제일(wdate) 기준, pay_method='init' 및 취소건 제외\n"
     "· 작성기준일 %s — 9월 출발분은 예약·수금이 진행 중으로 마감 시 증가 예상\n"
     "· 원가(가이드·차량·호텔 정산) 미반영 매출 기준 자료로 손익 판단에는 정산 데이터 추가 필요" % ASOF),
]
yy = 1.40
for title, body in blocks:
    h = 1.55 if title != "데이터 기준 및 유의사항" else 1.95
    rect(s, 0.55, yy, 12.2, h, CARD, rounded=True)
    txt(s, 0.85, yy + 0.15, 11.6, 0.3, title, sz=13, bold=True, color=NAVY)
    txt(s, 0.85, yy + 0.55, 11.6, h - 0.6, body, sz=11, color=NAVY, spacing=1.3)
    yy += h + 0.2

prs.save(OUT)
print("생성 완료:", OUT)
