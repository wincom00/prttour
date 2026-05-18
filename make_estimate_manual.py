# -*- coding: utf-8 -*-
"""
estimate_list.php / estimate_form.php 사용자 매뉴얼 PPTX 생성.

기존 input_batch_user_manual.pptx와 같은 네이비 헤더, 목차, 단계별 카드 형식을
따르되, 내용은 견적 목록/견적 입력 화면 기준으로 구성한다.
"""
from __future__ import annotations

import os
import textwrap
import time
from dataclasses import dataclass
from typing import Iterable, List, Sequence, Tuple

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt


# ===== input_batch_user_manual.pptx 계열 색상 =====
NAVY = RGBColor(0x13, 0x11, 0x76)
NAVY_DARK = RGBColor(0x0A, 0x09, 0x55)
BLUE = RGBColor(0x0B, 0x5B, 0xD3)
GREEN = RGBColor(0x28, 0xA7, 0x45)
TEAL = RGBColor(0x20, 0xC9, 0x97)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)
RED = RGBColor(0xDC, 0x35, 0x45)
GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)
GRAY_TXT = RGBColor(0x6C, 0x75, 0x7D)
DARK_TXT = RGBColor(0x21, 0x25, 0x29)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE7, 0xF3, 0xFF)
LIGHT_GREEN = RGBColor(0xE8, 0xF7, 0xEE)
LIGHT_YELLOW = RGBColor(0xFF, 0xF8, 0xE1)
LIGHT_RED = RGBColor(0xFF, 0xEE, 0xEE)

KFONT = "Malgun Gothic"
OUT = r"d:/www/prttour_myprt/estimate_user_manual_with_view_intranet.pptx"
PREVIEW_PREFIX = r"d:/www/prttour_myprt/estimate_user_manual_with_view_slide_"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def set_text(
    shape,
    text,
    size=14,
    bold=False,
    color=DARK_TXT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.MIDDLE,
    font=KFONT,
):
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = str(line)
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def add_text(slide, x, y, w, h, text, **kw):
    sh = slide.shapes.add_textbox(x, y, w, h)
    set_text(sh, text, **kw)
    return sh


def add_header_bar(slide, title, subtitle=None, page=None, total=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), NAVY)
    add_text(slide, Inches(0.4), Inches(0.12), Inches(10), Inches(0.45), title,
             size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.4), Inches(0.52), Inches(10), Inches(0.35), subtitle,
                 size=11, color=RGBColor(0xCB, 0xD3, 0xE8))
    if page is not None:
        add_text(slide, Inches(11.5), Inches(0.28), Inches(1.6), Inches(0.4),
                 f"{page} / {total}", size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.RIGHT)


def add_footer(slide, source="estimate_list.php / estimate_form.php / estimate_view.php"):
    add_rect(slide, 0, Inches(7.15), SW, Inches(0.35), GRAY_BG)
    add_text(slide, Inches(0.4), Inches(7.18), Inches(8), Inches(0.3),
             "푸른투어 인트라넷 · 견적 사용자 매뉴얼",
             size=9, color=GRAY_TXT)
    add_text(slide, Inches(8.4), Inches(7.18), Inches(4.5), Inches(0.3),
             source, size=9, color=GRAY_TXT, align=PP_ALIGN.RIGHT)


def bullet(slide, x, y, w, items, size=13, bullet_color=NAVY, gap=0.36):
    cur_y = y
    for item in items:
        ic = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, cur_y + Emu(45000),
                                    Inches(0.22), Inches(0.22))
        ic.fill.solid()
        ic.fill.fore_color.rgb = bullet_color
        ic.line.fill.background()
        set_text(ic, "✓", size=9, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.32), cur_y, w - Inches(0.32), Inches(gap),
                 item, size=size, color=DARK_TXT)
        cur_y += Inches(gap)


def step_card(slide, x, y, w, h, no, title, body, color=NAVY):
    add_rect(slide, x, y, w, h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
    add_rect(slide, x, y, Inches(0.82), h, color)
    add_text(slide, x, y, Inches(0.82), h, no,
             size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(1.0), y + Inches(0.18), w - Inches(1.2), Inches(0.38),
             title, size=13, bold=True, color=DARK_TXT,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(1.0), y + Inches(0.66), w - Inches(1.2), h - Inches(0.82),
             body, size=10.5, color=GRAY_TXT,
             anchor=MSO_ANCHOR.TOP)


def table_like(slide, x, y, w, h, headers, rows, col_widths=None):
    add_rect(slide, x, y, w, h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
    if col_widths is None:
        col_widths = [1 / len(headers)] * len(headers)
    header_h = Inches(0.48)
    add_rect(slide, x, y, w, header_h, LIGHT_BLUE, line=RGBColor(0xDD, 0xDD, 0xDD))
    cur_x = x
    for hdr, ratio in zip(headers, col_widths):
        cw = int(w * ratio)
        add_text(slide, cur_x + Inches(0.06), y + Inches(0.05), cw - Inches(0.12), header_h - Inches(0.10),
                 hdr, size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        cur_x += cw
    row_h = int((h - header_h) / max(1, len(rows)))
    for r_idx, row in enumerate(rows):
        ry = y + header_h + row_h * r_idx
        cur_x = x
        if r_idx % 2 == 1:
            add_rect(slide, x, ry, w, row_h, RGBColor(0xFB, 0xFC, 0xFE))
        for txt, ratio in zip(row, col_widths):
            cw = int(w * ratio)
            add_text(slide, cur_x + Inches(0.06), ry + Inches(0.03), cw - Inches(0.12), row_h - Inches(0.06),
                     txt, size=9.0, color=DARK_TXT, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)
            cur_x += cw


# ===== 슬라이드 작성 =====
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, GRAY_BG)
    add_rect(s, 0, 0, SW, Inches(4.5), NAVY)
    for x, y, w in [(Inches(10.2), Inches(-1.5), Inches(5)), (Inches(-1.5), Inches(2.8), Inches(3.5))]:
        deco = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, w)
        deco.fill.solid()
        deco.fill.fore_color.rgb = NAVY_DARK
        deco.line.fill.background()

    add_text(s, Inches(0.8), Inches(0.7), Inches(6), Inches(0.4),
             "푸른투어 인트라넷 · 사용자 매뉴얼", size=14, bold=True,
             color=RGBColor(0xCB, 0xD3, 0xE8))
    add_text(s, Inches(0.8), Inches(1.35), Inches(11.5), Inches(1.15),
             "견적 목록 · 견적 입력", size=52, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(2.48), Inches(11), Inches(0.55),
             "Estimate List & Breakdown Quotation Guide",
             size=21, color=RGBColor(0xCB, 0xD3, 0xE8))
    add_text(s, Inches(0.8), Inches(3.22), Inches(11.2), Inches(0.55),
             "검색, 신규 작성, 항목별 금액 입력, 합계 확인, 저장까지 한 번에 보는 업무 흐름",
             size=16, color=WHITE)

    cards = [
        ("대상 페이지", "estimate_list.php", "견적 조회 · 수정 · 내보내기", NAVY),
        ("작성 화면", "estimate_form.php", "BREAKDOWN QUOTATION 입력", GREEN),
        ("핵심 결과", "TOTAL / 1인 요금", "항목별 소계 자동 합산", ORANGE),
    ]
    for i, (eyebrow, title, desc, color) in enumerate(cards):
        x = Inches(0.8 + i * 4.05)
        add_rect(s, x, Inches(5.0), Inches(3.8), Inches(1.6), WHITE, line=color)
        add_text(s, x + Inches(0.18), Inches(5.1), Inches(3.45), Inches(0.35),
                 eyebrow, size=12, bold=True, color=color)
        add_text(s, x + Inches(0.18), Inches(5.45), Inches(3.45), Inches(0.45),
                 title, size=17, bold=True, color=DARK_TXT)
        add_text(s, x + Inches(0.18), Inches(5.95), Inches(3.45), Inches(0.4),
                 desc, size=11, color=GRAY_TXT)
    add_text(s, Inches(0.8), Inches(6.85), Inches(11), Inches(0.35),
             "푸른투어 인트라넷", size=11, bold=True, color=NAVY)


def slide_toc(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "목차 · CONTENTS", "견적 목록과 견적 입력 화면 기준 사용자 흐름", 2, total)
    items = [
        ("01", "페이지 목적", "견적 메뉴에서 처리하는 업무"),
        ("02", "전체 업무 흐름", "목록 → 신규/수정 → 저장 → 내보내기"),
        ("03", "견적 목록 화면", "검색 조건과 결과 테이블"),
        ("04", "신규/수정 진입", "새 견적과 기존 견적 수정"),
        ("05", "기본 정보 입력", "TO, PAX, FOC, 일정, GROUP NAME"),
        ("06", "HOTEL 입력", "호텔 행 추가/복제/삭제와 소계"),
        ("07", "일자별 항목", "MEAL, TRANSPORTATION, OVERTIME"),
        ("08", "기타 비용", "TICKET, GUIDE, ETC, TIP, PROFIT"),
        ("09", "합계 확인", "TOTAL TOUR FEE와 1인 요금"),
        ("10", "저장/내보내기", "저장 후 목록 수정/출력 흐름"),
        ("11", "PDF 프리뷰", "엑셀/PDF/이메일/수정/목록 버튼"),
        ("12", "체크리스트", "입력 전후 확인 사항"),
        ("13", "FAQ", "사용 중 자주 생기는 질문"),
    ]
    box_w, box_h = Inches(6.0), Inches(0.64)
    sx, sy = Inches(0.5), Inches(1.1)
    for i, (no, title, sub) in enumerate(items):
        col, row = i % 2, i // 2
        x = sx + (box_w + Inches(0.3)) * col
        y = sy + (box_h + Inches(0.10)) * row
        add_rect(s, x, y, box_w, box_h, WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, x, y, Inches(0.85), box_h, NAVY)
        add_text(s, x, y + Inches(0.05), Inches(0.85), Inches(0.45), no,
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(1.05), y + Inches(0.10), box_w - Inches(1.2), Inches(0.30),
                 title, size=13, bold=True, color=DARK_TXT)
        add_text(s, x + Inches(1.05), y + Inches(0.44), box_w - Inches(1.2), Inches(0.25),
                 sub, size=9.5, color=GRAY_TXT)
    add_footer(s)


def slide_overview(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "01. 페이지 목적", "견적 조회와 BREAKDOWN QUOTATION 입력", 3, total)
    add_rect(s, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.9), LIGHT_BLUE)
    add_text(s, Inches(0.7), Inches(1.23), Inches(11.9), Inches(0.58),
             "견적 목록에서 기존 견적을 찾고, 견적 입력 화면에서 항목별 비용을 작성한 뒤 총액과 1인 요금을 저장합니다.",
             size=15, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    cards = [
        ("목록", "검색 · 조회", ["견적번호, GROUP, TO로 검색", "작성일 기간 조건 검색", "최근 300건을 id DESC 순서로 표시"], NAVY),
        ("입력", "항목별 작성", ["기본 정보 입력", "HOTEL / MEAL / TRANSPORTATION", "OVERTIME / TICKET / GUIDE / ETC"], GREEN),
        ("결과", "합계 · 저장", ["항목별 소계 자동 계산", "TOTAL TOUR FEE 확인", "1인 요금 저장 및 출력"], ORANGE),
    ]
    for i, (label, title, items, color) in enumerate(cards):
        x = Inches(0.55 + i * 4.18)
        add_rect(s, x, Inches(2.35), Inches(3.9), Inches(4.35), WHITE, line=color)
        add_text(s, x + Inches(0.2), Inches(2.55), Inches(0.65), Inches(0.5),
                 label, size=15, bold=True, color=color)
        add_text(s, x + Inches(0.2), Inches(3.05), Inches(3.45), Inches(0.45),
                 title, size=20, bold=True, color=DARK_TXT)
        bullet(s, x + Inches(0.25), Inches(3.75), Inches(3.4), items,
               size=12.2, bullet_color=color, gap=0.46)
    add_footer(s)


def slide_flow(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "02. 전체 업무 흐름", "목록에서 시작해 저장과 내보내기까지", 4, total)
    steps = [
        ("1", "목록 검색", "견적번호 / GROUP / TO 또는 작성일 기간으로 견적을 찾습니다.", BLUE),
        ("2", "새 견적", "+ 새 견적으로 빈 BREAKDOWN QUOTATION 화면을 엽니다.", GREEN),
        ("3", "기본 정보", "TO, PAX, FOC, 총인원, 그룹명, 작성일, 여행기간을 입력합니다.", TEAL),
        ("4", "비용 입력", "HOTEL, MEAL, 차량, 오버타임, 입장권, 가이드/기사, 기타비용을 입력합니다.", ORANGE),
        ("5", "합계 확인", "각 항목 소계와 TOTAL TOUR FEE, 1인 요금을 확인합니다.", RED),
        ("6", "저장/출력", "저장 후 목록에서 수정하거나 새 창으로 내보내기를 실행합니다.", NAVY),
    ]
    x0, y0 = Inches(0.7), Inches(1.35)
    for i, (no, title, body, color) in enumerate(steps):
        row, col = divmod(i, 3)
        x = x0 + col * Inches(4.15)
        y = y0 + row * Inches(2.35)
        add_rect(s, x, y, Inches(3.65), Inches(1.75), WHITE, line=color)
        add_rect(s, x, y, Inches(0.65), Inches(1.75), color)
        add_text(s, x, y + Inches(0.48), Inches(0.65), Inches(0.5), no,
                 size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.85), y + Inches(0.24), Inches(2.6), Inches(0.34),
                 title, size=15, bold=True, color=DARK_TXT)
        add_text(s, x + Inches(0.85), y + Inches(0.70), Inches(2.6), Inches(0.82),
                 body, size=11.2, color=GRAY_TXT)
    add_footer(s)


def slide_list(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "03. 견적 목록 화면", "estimate_list.php: 검색 조건과 결과 테이블", 5, total)
    add_rect(s, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.55), LIGHT_BLUE)
    add_text(s, Inches(0.7), Inches(1.14), Inches(11.8), Inches(0.35),
             "검색 조건은 q, from, to이며 결과는 estimate_master에서 최신순 최대 300건을 표시합니다.",
             size=13, bold=True, color=NAVY)
    table_like(
        s, Inches(0.65), Inches(1.9), Inches(12.0), Inches(2.35),
        ["견적번호", "GROUP", "TO", "PAX", "여행기간", "총액", "1인요금", "작성/수정일", "액션"],
        [
            ["EST-...", "GROUP NAME", "거래처", "총인원", "시작~종료", "grand_total", "per_pax", "wdate/updated_at", "수정"],
            ["검색 결과", "그룹명", "담당 TO", "total_pax", "start~end", "합계", "인당 금액", "날짜", "내보내기"],
        ],
        [0.11, 0.13, 0.12, 0.07, 0.15, 0.10, 0.10, 0.12, 0.10],
    )
    step_card(s, Inches(0.65), Inches(4.65), Inches(3.85), Inches(1.55),
              "A", "검색어 입력", "견적번호, 그룹명, TO명 중 하나를 입력합니다.", BLUE)
    step_card(s, Inches(4.75), Inches(4.65), Inches(3.85), Inches(1.55),
              "B", "작성일 기간", "from/to 날짜로 작성일 범위를 제한합니다.", GREEN)
    step_card(s, Inches(8.85), Inches(4.65), Inches(3.85), Inches(1.55),
              "C", "액션 선택", "수정은 입력화면, 내보내기는 새 창 출력 화면으로 이동합니다.", ORANGE)
    add_footer(s, "estimate_list.php")


def slide_entry(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "04. 신규/수정 진입", "+ 새 견적과 수정 버튼의 사용 차이", 6, total)
    add_rect(s, Inches(0.65), Inches(1.15), Inches(5.85), Inches(4.95), WHITE, line=GREEN)
    add_text(s, Inches(0.9), Inches(1.38), Inches(5.3), Inches(0.42),
             "+ 새 견적", size=23, bold=True, color=GREEN)
    bullet(s, Inches(0.95), Inches(2.05), Inches(5.1), [
        "목록 상단의 + 새 견적 버튼을 클릭합니다.",
        "estimate_form.php가 id 없이 열립니다.",
        "빈 견적 화면에서 기본 정보와 비용 항목을 새로 입력합니다.",
        "저장 시 estimate_no는 EST-시간값 형식으로 생성됩니다.",
    ], size=13.2, bullet_color=GREEN, gap=0.55)
    add_rect(s, Inches(6.85), Inches(1.15), Inches(5.85), Inches(4.95), WHITE, line=NAVY)
    add_text(s, Inches(7.1), Inches(1.38), Inches(5.3), Inches(0.42),
             "기존 견적 수정", size=23, bold=True, color=NAVY)
    bullet(s, Inches(7.15), Inches(2.05), Inches(5.1), [
        "목록 각 행의 수정 버튼을 클릭합니다.",
        "estimate_form.php?id=견적ID 형식으로 열립니다.",
        "estimate_master와 estimate_items 기존 데이터가 화면에 로드됩니다.",
        "수정 후 저장하면 같은 견적 ID 기준으로 갱신됩니다.",
    ], size=13.2, bullet_color=NAVY, gap=0.55)
    add_footer(s, "estimate_list.php / estimate_form.php")


def slide_basic(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "05. 기본 정보 입력", "견적 계산의 기준값을 먼저 입력", 7, total)
    fields = [
        ("TO", "거래처/담당", "견적을 받을 대상"),
        ("PAX", "인원", "유료 인원"),
        ("FOC", "무상 인원", "총인원에 포함"),
        ("총인원", "자동 계산", "PAX + FOC"),
        ("GROUP NAME", "그룹명", "행사/단체명"),
        ("작성일", "wdate", "기본값 오늘"),
        ("여행 시작일", "start_date", "일자별 항목 기준"),
        ("여행 종료일", "end_date", "일자별 항목 기준"),
    ]
    for i, (name, label, desc) in enumerate(fields):
        row, col = divmod(i, 4)
        x = Inches(0.65 + col * 3.13)
        y = Inches(1.25 + row * 1.35)
        add_rect(s, x, y, Inches(2.8), Inches(0.95), WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_text(s, x + Inches(0.15), y + Inches(0.12), Inches(2.5), Inches(0.25),
                 name, size=12, bold=True, color=NAVY)
        add_text(s, x + Inches(0.15), y + Inches(0.40), Inches(2.5), Inches(0.22),
                 label, size=10.5, color=DARK_TXT)
        add_text(s, x + Inches(0.15), y + Inches(0.64), Inches(2.5), Inches(0.2),
                 desc, size=9, color=GRAY_TXT)
    add_rect(s, Inches(0.65), Inches(4.55), Inches(12.0), Inches(1.4), LIGHT_YELLOW)
    add_text(s, Inches(0.9), Inches(4.75), Inches(11.4), Inches(0.35),
             "입력 순서 팁", size=16, bold=True, color=ORANGE)
    add_text(s, Inches(0.9), Inches(5.16), Inches(11.4), Inches(0.55),
             "먼저 PAX/FOC와 여행 시작·종료일을 넣어야 MEAL, TRANSPORTATION, OVERTIME 날짜열 생성 버튼을 정확하게 사용할 수 있습니다.",
             size=13.5, color=DARK_TXT)
    add_footer(s, "estimate_form.php")


def slide_hotel(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "06. HOTEL 입력", "행 추가, 복제, 삭제 후 숙박 소계를 확인", 8, total)
    table_like(
        s, Inches(0.65), Inches(1.25), Inches(12.0), Inches(2.25),
        ["지역", "날짜", "요일", "호텔명", "방수", "요금(USD)", "박수", "합계"],
        [
            ["NYC/NIA", "숙박일", "자동", "호텔 선택", "room", "rate", "night", "자동"],
            ["행 추가", "복제 가능", "삭제 가능", "호텔 옵션", "입력", "입력", "입력", "방수×요금×박수"],
        ],
        [0.09, 0.11, 0.07, 0.26, 0.09, 0.12, 0.08, 0.18],
    )
    cards = [
        ("+ 행추가", "새 호텔 비용 행을 추가합니다.", GREEN),
        ("행 복제", "선택한 행을 복제해 비슷한 숙박 조건을 빠르게 작성합니다.", BLUE),
        ("행삭제", "불필요한 호텔 행을 삭제합니다.", ORANGE),
    ]
    for i, (title, body, color) in enumerate(cards):
        step_card(s, Inches(0.65 + i * 4.1), Inches(4.0), Inches(3.85), Inches(1.35),
                  chr(65 + i), title, body, color)
    add_rect(s, Inches(0.65), Inches(5.75), Inches(12), Inches(0.7), LIGHT_GREEN)
    add_text(s, Inches(0.9), Inches(5.88), Inches(11.4), Inches(0.38),
             "계산 기준: 방수 × 요금(USD) × 박수 = HOTEL 행 합계, 모든 행 합계가 HOTEL 소계로 표시됩니다.",
             size=13, bold=True, color=GREEN)
    add_footer(s, "estimate_form.php")


def slide_date_matrix(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "07. 일자별 항목", "MEAL, TRANSPORTATION, OVERTIME 날짜열 생성", 9, total)
    blocks = [
        ("MEAL", "조식 / 중식 / 석식", "여행 시작~종료일 기준 날짜열 생성\n일자별 식사 단가와 인원으로 소계 계산", GREEN),
        ("TRANSPORTATION", "일자별 차량료", "차량명 입력 후 날짜별 금액 입력\n차량수와 날짜별 금액으로 소계 계산", BLUE),
        ("OVERTIME", "일자별 오버타임", "각 날짜 칸에 금액과 사유 입력\n건수와 날짜별 금액으로 소계 계산", ORANGE),
    ]
    for i, (title, sub, body, color) in enumerate(blocks):
        x = Inches(0.65 + i * 4.1)
        add_rect(s, x, Inches(1.25), Inches(3.85), Inches(4.8), WHITE, line=color)
        add_rect(s, x, Inches(1.25), Inches(3.85), Inches(0.55), color)
        add_text(s, x + Inches(0.15), Inches(1.35), Inches(3.55), Inches(0.3),
                 title, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.2), Inches(2.05), Inches(3.45), Inches(0.35),
                 sub, size=15, bold=True, color=DARK_TXT)
        bullet(s, x + Inches(0.25), Inches(2.65), Inches(3.35),
               body.split("\n") + ["버튼: 여행일로 날짜열 생성"], size=12,
               bullet_color=color, gap=0.50)
    add_footer(s, "estimate_form.php")


def slide_other_costs(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "08. 기타 비용 입력", "입장권, 가이드/기사, 기타경비, 팁, 회사 수익금", 10, total)
    items = [
        ("TICKET", "입장지 / 단가 / 인원 / 합계", "입장권 코드 옵션을 선택하고 수량 기준 비용을 입력합니다.", BLUE),
        ("GUIDE", "항목 / 기간 / 인원·대수 / 단가", "가이드비, 드라이버 비용 등 인력성 비용을 입력합니다.", GREEN),
        ("ETC", "항목 / 인원·수량 / 단가", "식수, 라벨, 기타 실비성 경비를 입력합니다.", ORANGE),
        ("TIP", "가이드 팁 / 기사 팁 / 일수", "일수는 여행 기간 기준으로 자동 반영됩니다.", TEAL),
        ("PROFIT", "마진 금액 / 비고", "회사 수익금 또는 별도 마진을 입력합니다.", RED),
    ]
    for i, (title, field, desc, color) in enumerate(items):
        y = Inches(1.15 + i * 1.05)
        add_rect(s, Inches(0.65), y, Inches(12.0), Inches(0.82), WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, Inches(0.65), y, Inches(1.55), Inches(0.82), color)
        add_text(s, Inches(0.65), y + Inches(0.12), Inches(1.55), Inches(0.36),
                 title, size=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.4), y + Inches(0.10), Inches(3.4), Inches(0.28),
                 field, size=12.3, bold=True, color=DARK_TXT)
        add_text(s, Inches(2.4), y + Inches(0.43), Inches(9.8), Inches(0.25),
                 desc, size=10.5, color=GRAY_TXT)
    add_footer(s, "estimate_form.php")


def slide_totals(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "09. 합계 확인", "항목별 소계에서 TOTAL TOUR FEE와 1인 요금까지", 11, total)
    add_rect(s, Inches(0.65), Inches(1.15), Inches(12), Inches(1.2), LIGHT_BLUE)
    add_text(s, Inches(0.9), Inches(1.34), Inches(11.4), Inches(0.45),
             "하단 합계 영역에서 모든 항목 소계가 한 줄로 표시되고, 총 투어비와 1인 요금이 계산됩니다.",
             size=15, bold=True, color=NAVY)
    labels = ["HOTEL", "MEAL", "TRANS", "OVERTIME", "TICKET", "GUIDE", "ETC", "TIP", "PROFIT"]
    for i, label in enumerate(labels):
        row, col = divmod(i, 3)
        x = Inches(0.85 + col * 4.0)
        y = Inches(2.8 + row * 0.75)
        add_rect(s, x, y, Inches(3.55), Inches(0.48), WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_text(s, x + Inches(0.15), y + Inches(0.08), Inches(1.5), Inches(0.25),
                 label, size=11.5, bold=True, color=DARK_TXT)
        add_text(s, x + Inches(1.9), y + Inches(0.08), Inches(1.4), Inches(0.25),
                 "0", size=11.5, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    add_rect(s, Inches(0.85), Inches(5.35), Inches(5.8), Inches(0.72), NAVY)
    add_text(s, Inches(1.1), Inches(5.50), Inches(3.3), Inches(0.3),
             "10) TOTAL TOUR FEE", size=15, bold=True, color=WHITE)
    add_text(s, Inches(4.5), Inches(5.50), Inches(1.8), Inches(0.3),
             "자동 합계", size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    add_rect(s, Inches(6.85), Inches(5.35), Inches(5.8), Inches(0.72), GREEN)
    add_text(s, Inches(7.1), Inches(5.50), Inches(3.3), Inches(0.3),
             "11) 1인당 요금", size=15, bold=True, color=WHITE)
    add_text(s, Inches(10.5), Inches(5.50), Inches(1.8), Inches(0.3),
             "총액 ÷ 인원", size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    add_footer(s, "estimate_form.php")


def slide_save_export(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "10. 저장/내보내기", "저장 후 목록에서 수정하거나 출력 화면을 엽니다", 12, total)
    add_rect(s, Inches(0.65), Inches(1.15), Inches(5.85), Inches(4.9), WHITE, line=NAVY)
    add_text(s, Inches(0.9), Inches(1.38), Inches(5.25), Inches(0.42),
             "저장", size=23, bold=True, color=NAVY)
    bullet(s, Inches(0.95), Inches(2.05), Inches(5.1), [
        "화면 하단 저장 버튼을 클릭합니다.",
        "estimate_save.php로 master와 items 데이터가 전송됩니다.",
        "저장 성공 후 화면의 estimate_id가 갱신됩니다.",
        "saveStatus 영역에서 저장 결과를 확인합니다.",
    ], size=13.1, bullet_color=NAVY, gap=0.55)
    add_rect(s, Inches(6.85), Inches(1.15), Inches(5.85), Inches(4.9), WHITE, line=ORANGE)
    add_text(s, Inches(7.1), Inches(1.38), Inches(5.25), Inches(0.42),
             "목록 액션", size=23, bold=True, color=ORANGE)
    bullet(s, Inches(7.15), Inches(2.05), Inches(5.1), [
        "수정: estimate_form.php?id=견적ID로 다시 엽니다.",
        "내보내기: estimate_view.php?id=견적ID를 새 창으로 엽니다.",
        "목록은 작성일과 수정일을 함께 보여줍니다.",
        "검색 결과가 없으면 결과 없음 메시지가 표시됩니다.",
    ], size=13.1, bullet_color=ORANGE, gap=0.55)
    add_footer(s, "estimate_list.php / estimate_form.php")


def slide_view_preview(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "11. PDF 프리뷰 / 내보내기 화면", "estimate_view.php: 저장된 견적을 출력용 화면으로 확인", 13, total)

    add_rect(s, Inches(0.65), Inches(1.10), Inches(12.0), Inches(0.70), LIGHT_BLUE)
    add_text(s, Inches(0.9), Inches(1.24), Inches(11.4), Inches(0.35),
             "목록의 내보내기 버튼 또는 PDF 버튼으로 저장된 견적의 출력용 BREAKDOWN QUOTATION 화면을 확인합니다.",
             size=13.5, bold=True, color=NAVY)

    table_like(
        s, Inches(0.65), Inches(2.08), Inches(12.0), Inches(1.75),
        ["상단 버튼", "역할", "이동/처리"],
        [
            ["엑셀", "견적 데이터를 엑셀로 다운로드", "estimate_export_breakdown.php?action=download_excel"],
            ["PDF", "PDF 출력/생성 화면 호출", "estimate_export_breakdown.php?action=pdf"],
            ["이메일 발송", "현재 견적 기준 이메일 발송", "send_breakdown_email.php"],
            ["수정 / 목록", "입력 화면 또는 목록으로 이동", "estimate_form.php / estimate_list.php"],
        ],
        [0.18, 0.32, 0.50],
    )

    blocks = [
        ("기본 정보", "PAX, FOC, 총인원, TO, 여행 시작/종료일, 작성일, GROUP을 상단 카드로 표시합니다.", BLUE),
        ("섹션별 출력", "HOTEL, MEAL, TRANSPORTATION, OVERTIME, TICKET, GUIDE, ETC, TIP, PROFIT을 표 형식으로 보여줍니다.", GREEN),
        ("합계 영역", "각 섹션 소계를 pill 형태로 모아 보여주고 TOTAL TOUR FEE와 1인당 요금을 크게 표시합니다.", ORANGE),
    ]
    for i, (title, body, color) in enumerate(blocks):
        step_card(s, Inches(0.65 + i * 4.1), Inches(4.12), Inches(3.85), Inches(1.42),
                  chr(65 + i), title, body, color)

    add_rect(s, Inches(0.65), Inches(5.88), Inches(12.0), Inches(0.65), LIGHT_YELLOW)
    add_text(s, Inches(0.9), Inches(6.01), Inches(11.4), Inches(0.32),
             "출력 시 action-buttons는 숨김 처리되고, 표와 합계 영역은 A4 출력에 맞게 page-break 방지 스타일을 사용합니다.",
             size=12.6, bold=True, color=ORANGE)
    add_footer(s, "estimate_view.php")


def slide_checklist(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "12. 체크리스트", "견적 저장 전후 확인 사항", 14, total)
    groups = [
        ("입력 전", ["TO와 GROUP NAME 확인", "PAX/FOC 입력 후 총인원 확인", "여행 시작일/종료일 입력"], BLUE),
        ("비용 입력 중", ["일자별 항목 날짜열 생성", "호텔 박수/요금/방수 확인", "오버타임 금액과 사유 확인"], GREEN),
        ("저장 전", ["TOTAL TOUR FEE 확인", "1인당 요금 확인", "마진/팁 누락 여부 확인"], ORANGE),
        ("저장 후", ["목록에서 견적번호 확인", "수정 버튼으로 재조회 확인", "내보내기 화면에서 출력 상태 확인"], NAVY),
    ]
    for i, (title, items, color) in enumerate(groups):
        row, col = divmod(i, 2)
        x = Inches(0.65 + col * 6.15)
        y = Inches(1.25 + row * 2.45)
        add_rect(s, x, y, Inches(5.75), Inches(2.0), WHITE, line=color)
        add_rect(s, x, y, Inches(5.75), Inches(0.45), color)
        add_text(s, x + Inches(0.18), y + Inches(0.08), Inches(5.4), Inches(0.25),
                 title, size=13.5, bold=True, color=WHITE)
        bullet(s, x + Inches(0.25), y + Inches(0.72), Inches(5.25), items,
               size=12.2, bullet_color=color, gap=0.40)
    add_footer(s, "estimate_list.php / estimate_form.php")


def slide_faq(total):
    s = prs.slides.add_slide(BLANK)
    add_header_bar(s, "13. FAQ", "사용 중 자주 확인할 내용", 15, total)
    faqs = [
        ("Q1", "검색 결과가 없을 때", "검색어(q), 작성일 from/to 조건을 줄이거나 빈 값으로 다시 조회합니다."),
        ("Q2", "날짜열이 안 보일 때", "여행 시작일과 종료일을 먼저 입력한 뒤 날짜열 생성 버튼을 누릅니다."),
        ("Q3", "총인원이 이상할 때", "PAX와 FOC 값을 확인합니다. 총인원은 PAX + FOC 자동 계산값입니다."),
        ("Q4", "합계가 안 맞을 때", "각 행의 수량, 단가, 박수, 건수 입력값과 항목별 소계를 먼저 확인합니다."),
        ("Q5", "기존 견적을 다시 열 때", "목록의 수정 버튼으로 estimate_form.php?id=견적ID 화면을 엽니다."),
    ]
    y = Inches(1.15)
    for q, title, ans in faqs:
        add_rect(s, Inches(0.65), y, Inches(12.0), Inches(0.86), WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
        add_rect(s, Inches(0.65), y, Inches(0.85), Inches(0.86), NAVY)
        add_text(s, Inches(0.65), y + Inches(0.14), Inches(0.85), Inches(0.32),
                 q, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.7), y + Inches(0.08), Inches(10.6), Inches(0.30),
                 title, size=12.5, bold=True, color=DARK_TXT)
        add_text(s, Inches(1.7), y + Inches(0.44), Inches(10.6), Inches(0.25),
                 ans, size=10.4, color=GRAY_TXT)
        y += Inches(0.96)
    add_footer(s, "estimate_list.php / estimate_form.php")


def slide_thanks(total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, NAVY)
    for x, y, w in [(Inches(-2), Inches(-2), Inches(6)), (Inches(9.5), Inches(4), Inches(6))]:
        deco = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, w)
        deco.fill.solid()
        deco.fill.fore_color.rgb = NAVY_DARK
        deco.line.fill.background()
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0),
             "Thank You", size=66, bold=True, color=WHITE)
    add_text(s, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.6),
             "견적 목록과 견적 입력 화면은 목록 조회 → 비용 입력 → 저장 → 출력 흐름으로 사용합니다.",
             size=19, color=RGBColor(0xCB, 0xD3, 0xE8))
    add_rect(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.45), WHITE)
    add_text(s, Inches(1.0), Inches(4.65), Inches(11), Inches(0.35),
             "확인 파일", size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.0), Inches(5.05), Inches(11), Inches(0.55),
             "estimate_list.php / estimate_form.php", size=16, bold=True, color=DARK_TXT)
    add_text(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.4),
             "푸른투어 인트라넷 · estimate manual",
             size=11, color=RGBColor(0xCB, 0xD3, 0xE8), align=PP_ALIGN.CENTER)


# ===== PNG 미리보기 생성 =====
@dataclass
class PreviewSlide:
    title: str
    subtitle: str
    rows: Sequence[Tuple[str, str]]
    dark: bool = False


def rgb_tuple(color: RGBColor):
    return (color[0], color[1], color[2])


def load_font(size: int, bold: bool = False):
    base = r"C:/Windows/Fonts/malgunbd.ttf" if bold else r"C:/Windows/Fonts/malgun.ttf"
    if os.path.exists(base):
        return ImageFont.truetype(base, size)
    return ImageFont.load_default()


def wrap_text(draw, text, font, max_width):
    words = str(text).split()
    if not words:
        return [""]
    lines, cur = [], words[0]
    for word in words[1:]:
        trial = cur + " " + word
        if draw.textbbox((0, 0), trial, font=font)[2] <= max_width:
            cur = trial
        else:
            lines.append(cur)
            cur = word
    lines.append(cur)
    return lines


def draw_preview(idx: int, spec: PreviewSlide):
    W, H = 1920, 1080
    img = Image.new("RGB", (W, H), rgb_tuple(GRAY_BG))
    d = ImageDraw.Draw(img)
    f_title = load_font(54, True)
    f_sub = load_font(25)
    f_card_title = load_font(27, True)
    f_body = load_font(23)
    f_small = load_font(18)

    if spec.dark:
        d.rectangle([0, 0, W, H], fill=rgb_tuple(NAVY))
        d.ellipse([1450, -220, 2200, 530], fill=rgb_tuple(NAVY_DARK))
        d.ellipse([-240, 420, 430, 1090], fill=rgb_tuple(NAVY_DARK))
        d.text((115, 280), spec.title, fill=rgb_tuple(WHITE), font=load_font(76, True))
        d.text((120, 415), spec.subtitle, fill=(203, 211, 232), font=f_sub)
        out = f"{PREVIEW_PREFIX}{idx:02d}.png"
        img.save(out)
        return out

    d.rectangle([0, 0, W, 130], fill=rgb_tuple(NAVY))
    d.text((58, 28), spec.title, fill=rgb_tuple(WHITE), font=f_title)
    d.text((62, 86), spec.subtitle, fill=(203, 211, 232), font=f_small)
    d.text((1700, 48), f"{idx} / 16", fill=rgb_tuple(WHITE), font=f_small)

    y = 170
    for i, (head, body) in enumerate(spec.rows):
        x = 90 + (i % 2) * 870
        if i and i % 2 == 0:
            y += 240
        d.rounded_rectangle([x, y, x + 790, y + 185], radius=8,
                            fill=rgb_tuple(WHITE), outline=(221, 221, 221), width=2)
        d.rectangle([x, y, x + 112, y + 185], fill=rgb_tuple(NAVY))
        d.text((x + 34, y + 58), str(i + 1), fill=rgb_tuple(WHITE), font=load_font(40, True))
        d.text((x + 145, y + 25), head, fill=rgb_tuple(DARK_TXT), font=f_card_title)
        line_y = y + 76
        for line in wrap_text(d, body, f_body, 590)[:3]:
            d.text((x + 145, line_y), line, fill=rgb_tuple(GRAY_TXT), font=f_body)
            line_y += 34

    d.rectangle([0, 1030, W, H], fill=rgb_tuple(GRAY_BG))
    d.text((58, 1040), "푸른투어 인트라넷 · 견적 사용자 매뉴얼",
           fill=rgb_tuple(GRAY_TXT), font=f_small)
    out = f"{PREVIEW_PREFIX}{idx:02d}.png"
    img.save(out)
    return out


def render_previews():
    specs = [
        PreviewSlide("견적 목록 · 견적 입력", "Estimate List & Breakdown Quotation Guide",
                     [("대상 페이지", "estimate_list.php"), ("작성 화면", "estimate_form.php")], True),
        PreviewSlide("목차 · CONTENTS", "견적 목록과 견적 입력 화면 기준 사용자 흐름",
                     [("페이지 목적", "목록 조회와 견적 작성"), ("업무 흐름", "신규/수정/저장/내보내기"), ("목록 화면", "검색과 액션"), ("입력 화면", "항목별 비용 입력")]),
        PreviewSlide("01. 페이지 목적", "견적 조회와 BREAKDOWN QUOTATION 입력",
                     [("목록", "검색 · 조회 · 수정 · 내보내기"), ("입력", "기본정보와 비용 항목 작성"), ("결과", "합계와 1인 요금 저장")]),
        PreviewSlide("02. 전체 업무 흐름", "목록에서 시작해 저장과 내보내기까지",
                     [("목록 검색", "견적번호, 그룹명, TO, 작성일 기간으로 검색"), ("새 견적", "빈 입력 화면으로 이동"), ("비용 입력", "항목별 비용과 날짜별 금액 입력"), ("저장/출력", "저장 후 수정 또는 내보내기")]),
        PreviewSlide("03. 견적 목록 화면", "estimate_list.php 검색 조건과 결과 테이블",
                     [("검색 조건", "q, from, to 값으로 조회"), ("결과 열", "견적번호, GROUP, TO, PAX, 기간, 총액"), ("액션", "수정과 내보내기 버튼 사용")]),
        PreviewSlide("04. 신규/수정 진입", "+ 새 견적과 수정 버튼의 사용 차이",
                     [("새 견적", "estimate_form.php가 id 없이 열림"), ("기존 수정", "estimate_form.php?id=견적ID로 로드"), ("저장 기준", "저장 후 estimate_id 갱신")]),
        PreviewSlide("05. 기본 정보 입력", "견적 계산의 기준값을 먼저 입력",
                     [("TO/PAX/FOC", "총인원은 PAX + FOC 자동 계산"), ("GROUP/작성일", "그룹명과 작성일 입력"), ("여행기간", "일자별 항목 날짜열 기준")]),
        PreviewSlide("06. HOTEL 입력", "행 추가, 복제, 삭제 후 숙박 소계 확인",
                     [("호텔 행", "지역, 날짜, 호텔명, 방수, 요금, 박수 입력"), ("계산", "방수 × 요금 × 박수"), ("소계", "모든 호텔 행 합산")]),
        PreviewSlide("07. 일자별 항목", "MEAL, TRANSPORTATION, OVERTIME 날짜열 생성",
                     [("MEAL", "조식/중식/석식 단가와 인원"), ("TRANSPORTATION", "일자별 차량료와 차량수"), ("OVERTIME", "일자별 금액과 사유 입력")]),
        PreviewSlide("08. 기타 비용 입력", "TICKET, GUIDE, ETC, TIP, PROFIT",
                     [("TICKET/GUIDE", "입장권과 가이드/기사 비용"), ("ETC/TIP", "기타경비와 팁"), ("PROFIT", "회사 수익금/마진")]),
        PreviewSlide("09. 합계 확인", "TOTAL TOUR FEE와 1인 요금",
                     [("소계", "항목별 금액 확인"), ("총액", "TOTAL TOUR FEE 자동 계산"), ("1인 요금", "총액 기준 인당 금액 확인")]),
        PreviewSlide("10. 저장/내보내기", "저장 후 목록에서 수정하거나 출력 화면 열기",
                     [("저장", "estimate_save.php로 데이터 전송"), ("수정", "목록 수정 버튼으로 재조회"), ("내보내기", "목록 내보내기 버튼으로 새 창 출력")]),
        PreviewSlide("11. PDF 프리뷰 / 내보내기 화면", "estimate_view.php 출력용 화면 확인",
                     [("상단 버튼", "엑셀, PDF, 이메일 발송, 수정, 목록"), ("섹션별 출력", "HOTEL부터 PROFIT까지 표 형식 표시"), ("최종 합계", "섹션별 소계와 TOTAL TOUR FEE 표시")]),
        PreviewSlide("12. 체크리스트", "견적 저장 전후 확인 사항",
                     [("입력 전", "인원과 여행기간 확인"), ("저장 전", "총액과 1인 요금 확인"), ("저장 후", "목록에서 수정/내보내기 확인")]),
        PreviewSlide("13. FAQ", "사용 중 자주 확인할 내용",
                     [("검색 결과 없음", "검색어 또는 작성일 조건을 줄여 다시 조회"), ("날짜열 없음", "여행기간 입력 후 날짜열 생성"), ("합계 확인", "항목별 소계와 총액을 순서대로 확인")]),
        PreviewSlide("Thank You", "견적 목록과 견적 입력 화면 매뉴얼", [], True),
    ]
    return [draw_preview(i + 1, spec) for i, spec in enumerate(specs)]


def build():
    total = 16
    slide_cover()
    slide_toc(total)
    slide_overview(total)
    slide_flow(total)
    slide_list(total)
    slide_entry(total)
    slide_basic(total)
    slide_hotel(total)
    slide_date_matrix(total)
    slide_other_costs(total)
    slide_totals(total)
    slide_save_export(total)
    slide_view_preview(total)
    slide_checklist(total)
    slide_faq(total)
    slide_thanks(total)

    try:
        prs.save(OUT)
        saved = OUT
    except PermissionError:
        saved = OUT.replace(".pptx", f"_v{int(time.time())}.pptx")
        prs.save(saved)
    previews = render_previews()
    print("SAVED:", saved)
    print("PREVIEWS:", len(previews))
    for path in previews:
        print(path)


if __name__ == "__main__":
    build()
