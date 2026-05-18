# -*- coding: utf-8 -*-
from __future__ import annotations

import json
import math
import re
import subprocess
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "prod_month_sales_trend_2025_2026_ytd.pptx"
PREVIEW_PREFIX = ROOT / "prod_month_sales_trend_slide_"

KFONT = "Malgun Gothic"
NAVY = RGBColor(0x13, 0x11, 0x76)
INK = RGBColor(0x1D, 0x25, 0x34)
MUTED = RGBColor(0x64, 0x70, 0x80)
BLUE = RGBColor(0x1F, 0x6F, 0xD5)
CYAN = RGBColor(0x21, 0xB6, 0xC7)
GREEN = RGBColor(0x2F, 0x9D, 0x58)
RED = RGBColor(0xD8, 0x44, 0x44)
ORANGE = RGBColor(0xED, 0x8A, 0x19)
GRAY_BG = RGBColor(0xF5, 0xF7, 0xFA)
LINE = RGBColor(0xD7, 0xDF, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def run_sales_query() -> dict:
    php = r"""
<?php
include 'include/dbconn.php';
$start = '2025-01-01';
$end = '2026-05-01';
$baseWhere = "b.rev_status = 'DONE'
  AND a.p_code NOT LIKE '%PICKUP%'
  AND a.p_code NOT LIKE '%SENDING%'
  AND b.parent = 'MAIN'
  AND b.revDate >= '$start' AND b.revDate < '$end'
  AND MONTH(b.revDate) BETWEEN 1 AND 4";

$out = array('monthly'=>array(), 'products'=>array(), 'pday'=>array());

$sql = "SELECT YEAR(b.revDate) AS yy, MONTH(b.revDate) AS mm,
        SUM(b.last_total) AS sales, SUM(b.p_cnt) AS pax, COUNT(*) AS bookings
        FROM product_master a INNER JOIN reserve_info b ON a.p_code = b.p_code
        WHERE $baseWhere
        GROUP BY YEAR(b.revDate), MONTH(b.revDate)
        ORDER BY yy, mm";
$res = mysql_query($sql, $dbConn);
while ($row = mysql_fetch_assoc($res)) { $out['monthly'][] = $row; }

$sql = "SELECT YEAR(b.revDate) AS yy, b.p_code, b.p_name, a.p_day,
        SUM(b.last_total) AS sales, SUM(b.p_cnt) AS pax, COUNT(*) AS bookings
        FROM product_master a INNER JOIN reserve_info b ON a.p_code = b.p_code
        WHERE $baseWhere
        GROUP BY YEAR(b.revDate), b.p_code, b.p_name, a.p_day
        ORDER BY yy, sales DESC";
$res = mysql_query($sql, $dbConn);
while ($row = mysql_fetch_assoc($res)) { $row['p_name'] = strip_tags($row['p_name']); $out['products'][] = $row; }

$sql = "SELECT YEAR(b.revDate) AS yy, a.p_day,
        SUM(b.last_total) AS sales, SUM(b.p_cnt) AS pax, COUNT(*) AS bookings
        FROM product_master a INNER JOIN reserve_info b ON a.p_code = b.p_code
        WHERE $baseWhere
        GROUP BY YEAR(b.revDate), a.p_day
        ORDER BY yy, a.p_day";
$res = mysql_query($sql, $dbConn);
while ($row = mysql_fetch_assoc($res)) { $out['pday'][] = $row; }

echo json_encode($out, JSON_UNESCAPED_UNICODE);
?>
"""
    proc = subprocess.run(
        ["php"],
        input=php,
        cwd=str(ROOT),
        text=True,
        encoding="utf-8",
        capture_output=True,
        check=True,
    )
    stdout = proc.stdout.strip()
    start = stdout.find("{")
    if start > 0:
        stdout = stdout[start:]
    return json.loads(stdout)


def money(v: float, digits: int = 1) -> str:
    return f"${v/1_000_000:.{digits}f}M"


def num(v: float) -> str:
    return f"{v:,.0f}"


def pct(v: float) -> str:
    sign = "+" if v > 0 else ""
    return f"{sign}{v:.1f}%"


def clean_name(name: str, max_len: int = 34) -> str:
    text = re.sub(r"\s+", " ", name).strip()
    return text if len(text) <= max_len else text[: max_len - 1] + "…"


def summarize(data: dict) -> dict:
    months = [1, 2, 3, 4]
    month_name = {1: "1월", 2: "2월", 3: "3월", 4: "4월"}
    monthly = {(int(r["yy"]), int(r["mm"])): r for r in data["monthly"]}
    sales = {y: [float(monthly.get((y, m), {}).get("sales", 0)) for m in months] for y in (2025, 2026)}
    pax = {y: [float(monthly.get((y, m), {}).get("pax", 0)) for m in months] for y in (2025, 2026)}
    bookings = {y: [float(monthly.get((y, m), {}).get("bookings", 0)) for m in months] for y in (2025, 2026)}
    total25, total26 = sum(sales[2025]), sum(sales[2026])
    pax25, pax26 = sum(pax[2025]), sum(pax[2026])
    bk25, bk26 = sum(bookings[2025]), sum(bookings[2026])

    product_map = {}
    for r in data["products"]:
        code = r["p_code"]
        item = product_map.setdefault(
            code,
            {
                "code": code,
                "name": clean_name(r["p_name"], 42),
                "p_day": r["p_day"],
                "s2025": 0.0,
                "s2026": 0.0,
                "p2025": 0,
                "p2026": 0,
            },
        )
        yy = int(r["yy"])
        item[f"s{yy}"] += float(r["sales"])
        item[f"p{yy}"] += int(float(r["pax"]))
    products = list(product_map.values())
    for item in products:
        item["delta"] = item["s2026"] - item["s2025"]
        item["pct"] = (item["delta"] / item["s2025"] * 100) if item["s2025"] else None
    products.sort(key=lambda x: x["delta"])

    pday_map = {}
    for r in data["pday"]:
        key = int(float(r["p_day"]))
        item = pday_map.setdefault(key, {"p_day": key, "s2025": 0.0, "s2026": 0.0, "p2025": 0, "p2026": 0})
        yy = int(r["yy"])
        item[f"s{yy}"] += float(r["sales"])
        item[f"p{yy}"] += int(float(r["pax"]))
    pday = list(pday_map.values())
    for item in pday:
        item["delta"] = item["s2026"] - item["s2025"]
        item["pct"] = (item["delta"] / item["s2025"] * 100) if item["s2025"] else None
    pday.sort(key=lambda x: x["delta"])

    return {
        "months": months,
        "month_names": [month_name[m] for m in months],
        "sales": sales,
        "pax": pax,
        "bookings": bookings,
        "total25": total25,
        "total26": total26,
        "delta": total26 - total25,
        "delta_pct": (total26 - total25) / total25 * 100,
        "pax25": pax25,
        "pax26": pax26,
        "pax_pct": (pax26 - pax25) / pax25 * 100,
        "bk25": bk25,
        "bk26": bk26,
        "bk_pct": (bk26 - bk25) / bk25 * 100,
        "avg_pax25": total25 / pax25 if pax25 else 0,
        "avg_pax26": total26 / pax26 if pax26 else 0,
        "avg_booking25": total25 / bk25 if bk25 else 0,
        "avg_booking26": total26 / bk26 if bk26 else 0,
        "products": products,
        "pday": pday,
    }


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, radius=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.8)
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, body, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    sh = slide.shapes.add_textbox(x, y, w, h)
    tf = sh.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(body)
    run.font.name = KFONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return sh


def header(slide, title, kicker="", page=None):
    text(slide, Inches(0.55), Inches(0.35), Inches(9.8), Inches(0.45), kicker, size=10, bold=True, color=BLUE)
    text(slide, Inches(0.55), Inches(0.72), Inches(10.8), Inches(0.55), title, size=26, bold=True, color=INK)
    if page:
        text(slide, Inches(12.0), Inches(0.48), Inches(0.8), Inches(0.35), str(page), size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    rect(slide, Inches(0.55), Inches(1.33), Inches(12.2), Inches(0.02), LINE)


def footer(slide, src="출처: prod_month.php 기준 DB 집계"):
    rect(slide, Inches(0.55), Inches(7.05), Inches(12.2), Inches(0.01), LINE)
    text(slide, Inches(0.55), Inches(7.12), Inches(11.9), Inches(0.25), src, size=8, color=MUTED)


def kpi(slide, x, y, label, value, sub, color=INK):
    text(slide, x, y, Inches(2.5), Inches(0.25), label, size=10, bold=True, color=MUTED)
    text(slide, x, y + Inches(0.28), Inches(2.5), Inches(0.45), value, size=26, bold=True, color=color)
    text(slide, x, y + Inches(0.78), Inches(2.5), Inches(0.32), sub, size=10, color=MUTED)


def add_line_chart(slide, x, y, w, h, cats, series):
    chart_data = CategoryChartData()
    chart_data.categories = cats
    for name, vals in series:
        chart_data.add_series(name, vals)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, w, h, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.number_format = '$0.0'
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.plots[0].series[0].format.line.color.rgb = BLUE
    chart.plots[0].series[1].format.line.color.rgb = RED
    return chart


def add_column_chart(slide, x, y, w, h, cats, series):
    chart_data = CategoryChartData()
    chart_data.categories = cats
    for name, vals in series:
        chart_data.add_series(name, vals)
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.number_format = '$0.0'
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = RED
    return chart


def add_bullet_list(slide, x, y, items, color=INK):
    cur = y
    for item in items:
        rect(slide, x, cur + Inches(0.08), Inches(0.08), Inches(0.08), color, shape=MSO_SHAPE.OVAL)
        text(slide, x + Inches(0.18), cur, Inches(5.6), Inches(0.38), item, size=12, color=INK)
        cur += Inches(0.48)


def slide_cover(s):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, SW, SH, GRAY_BG)
    rect(slide, Inches(0.55), Inches(0.45), Inches(0.15), Inches(6.3), BLUE)
    text(slide, Inches(0.9), Inches(0.75), Inches(6.0), Inches(0.4), "PROD_MONTH.PHP SALES DIAGNOSTIC", size=12, bold=True, color=BLUE)
    text(slide, Inches(0.9), Inches(1.25), Inches(9.4), Inches(1.25), "상반기 매출추이\n변동현황 분석", size=42, bold=True, color=INK)
    text(slide, Inches(0.95), Inches(2.75), Inches(7.8), Inches(0.45), "2025년 1-4월 vs 2026년 1-4월 | 2026-04-30 현재월 기준", size=15, color=MUTED)
    text(slide, Inches(9.15), Inches(1.2), Inches(3.2), Inches(0.35), "YTD 매출 증감", size=12, bold=True, color=MUTED)
    text(slide, Inches(8.85), Inches(1.62), Inches(3.8), Inches(0.8), pct(s["delta_pct"]), size=48, bold=True, color=RED, align=PP_ALIGN.RIGHT)
    text(slide, Inches(9.05), Inches(2.45), Inches(3.55), Inches(0.4), f"{money(s['total25'])} → {money(s['total26'])}", size=15, color=INK, align=PP_ALIGN.RIGHT)
    rect(slide, Inches(8.9), Inches(3.25), Inches(3.7), Inches(0.02), LINE)
    kpi(slide, Inches(8.95), Inches(3.55), "감소액", money(abs(s["delta"])), "전년 동기간 대비", RED)
    kpi(slide, Inches(8.95), Inches(4.85), "인원수", f"{num(s['pax25'])} → {num(s['pax26'])}", pct(s["pax_pct"]), ORANGE)
    footer(slide, "출처: prod_month.php SQL 기준, reserve_info.last_total 월별 합계")


def slide_basis(s):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "집계 기준과 한 줄 결론", "분석 범위", 2)
    text(slide, Inches(0.65), Inches(1.65), Inches(5.4), Inches(0.5), "prod_month.php 기준", size=20, bold=True)
    add_bullet_list(slide, Inches(0.75), Inches(2.25), [
        "매출액 기준: reserve_info.last_total 합계",
        "기간: 2025년 1-4월, 2026년 1-4월",
        "조건: rev_status='DONE', parent='MAIN'",
        "제외: PICKUP / SENDING 상품코드",
        "월 기준: reserve_info.revDate의 월",
    ], BLUE)
    text(slide, Inches(6.6), Inches(1.65), Inches(5.8), Inches(0.45), "핵심 수치", size=20, bold=True)
    kpi(slide, Inches(6.7), Inches(2.25), "매출", f"{money(s['total25'])} → {money(s['total26'])}", pct(s["delta_pct"]), RED)
    kpi(slide, Inches(9.7), Inches(2.25), "예약건수", f"{num(s['bk25'])} → {num(s['bk26'])}", pct(s["bk_pct"]), ORANGE)
    kpi(slide, Inches(6.7), Inches(3.75), "인원수", f"{num(s['pax25'])} → {num(s['pax26'])}", pct(s["pax_pct"]), ORANGE)
    avg_pct = (s["avg_booking26"] - s["avg_booking25"]) / s["avg_booking25"] * 100
    kpi(slide, Inches(9.7), Inches(3.75), "건당 매출", f"${s['avg_booking25']:,.0f} → ${s['avg_booking26']:,.0f}", pct(avg_pct), RED)
    text(slide, Inches(0.75), Inches(5.75), Inches(11.7), Inches(0.45), "결론: 시장 전체 전망은 회복 방향이나, 내부 데이터는 핵심 장기 인바운드·특정 계정 상품의 감소가 전체 하락을 주도합니다.", size=17, bold=True, color=INK)
    footer(slide)


def slide_monthly(s):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "월별 매출은 모든 월에서 전년 대비 하락", "매출 추이", 3)
    vals25 = [v / 1_000_000 for v in s["sales"][2025]]
    vals26 = [v / 1_000_000 for v in s["sales"][2026]]
    add_line_chart(slide, Inches(0.75), Inches(1.65), Inches(7.4), Inches(4.6), s["month_names"], [("2025", vals25), ("2026", vals26)])
    text(slide, Inches(8.6), Inches(1.75), Inches(3.9), Inches(0.45), "월별 YoY", size=20, bold=True)
    y = Inches(2.35)
    for i, m in enumerate(s["month_names"]):
        d = s["sales"][2026][i] - s["sales"][2025][i]
        p = d / s["sales"][2025][i] * 100
        text(slide, Inches(8.75), y, Inches(0.8), Inches(0.3), m, size=12, bold=True, color=MUTED)
        text(slide, Inches(9.55), y, Inches(1.4), Inches(0.3), money(d, 2), size=14, bold=True, color=RED)
        text(slide, Inches(11.05), y, Inches(1.0), Inches(0.3), pct(p), size=14, bold=True, color=RED, align=PP_ALIGN.RIGHT)
        y += Inches(0.55)
    text(slide, Inches(8.75), Inches(5.2), Inches(3.5), Inches(0.7), "4월 하락폭이 가장 큼: 전년 대비 약 -$1.33M, -50.9%", size=16, bold=True, color=RED)
    footer(slide)


def slide_delta(s):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "감소액의 37%가 4월에서 발생", "변동현황", 4)
    deltas = [(s["sales"][2026][i] - s["sales"][2025][i]) / 1_000_000 for i in range(4)]
    add_column_chart(slide, Inches(0.75), Inches(1.65), Inches(7.5), Inches(4.7), s["month_names"], [("감소액", deltas)])
    total_drop = abs(s["delta"])
    y = Inches(1.85)
    text(slide, Inches(8.65), y, Inches(3.6), Inches(0.5), "감소 기여도", size=20, bold=True)
    y += Inches(0.7)
    for i, m in enumerate(s["month_names"]):
        contrib = abs(s["sales"][2026][i] - s["sales"][2025][i]) / total_drop * 100
        bar_w = Inches(2.6 * contrib / 40)
        text(slide, Inches(8.75), y, Inches(0.7), Inches(0.25), m, size=11, bold=True, color=MUTED)
        rect(slide, Inches(9.55), y + Inches(0.03), bar_w, Inches(0.15), RED)
        text(slide, Inches(12.0), y - Inches(0.04), Inches(0.65), Inches(0.25), f"{contrib:.0f}%", size=11, bold=True, color=RED, align=PP_ALIGN.RIGHT)
        y += Inches(0.5)
    text(slide, Inches(8.75), Inches(5.1), Inches(3.7), Inches(0.9), "3-4월로 갈수록 하락률이 확대됩니다. 이는 시즌 진입 전환기에 상품 공급/확정/판매 타이밍 문제가 겹쳤을 가능성을 보여줍니다.", size=14, color=INK)
    footer(slide)


def slide_mix(s):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "하락은 9일·10일 장기 인바운드와 특정 계정 상품에 집중", "상품/일정 믹스", 5)
    top = s["products"][:6]
    x0, y0 = Inches(0.7), Inches(1.65)
    cols = [0.0, 1.35, 5.85, 7.25, 8.65, 10.05]
    heads = ["코드", "상품명", "2025", "2026", "증감", "YoY"]
    for c, h in zip(cols, heads):
        text(slide, x0 + Inches(c), y0, Inches(1.3 if h != "상품명" else 4.2), Inches(0.3), h, size=10, bold=True, color=MUTED)
    y = y0 + Inches(0.42)
    for item in top:
        rect(slide, x0, y + Inches(0.33), Inches(11.7), Inches(0.01), LINE)
        text(slide, x0 + Inches(cols[0]), y, Inches(1.25), Inches(0.35), item["code"], size=10, bold=True)
        text(slide, x0 + Inches(cols[1]), y, Inches(4.2), Inches(0.35), item["name"], size=9)
        text(slide, x0 + Inches(cols[2]), y, Inches(1.2), Inches(0.35), money(item["s2025"], 2), size=10)
        text(slide, x0 + Inches(cols[3]), y, Inches(1.2), Inches(0.35), money(item["s2026"], 2), size=10)
        text(slide, x0 + Inches(cols[4]), y, Inches(1.2), Inches(0.35), money(item["delta"], 2), size=10, bold=True, color=RED)
        text(slide, x0 + Inches(cols[5]), y, Inches(0.8), Inches(0.35), pct(item["pct"] or 0), size=10, bold=True, color=RED, align=PP_ALIGN.RIGHT)
        y += Inches(0.52)
    p9 = next((x for x in s["pday"] if x["p_day"] == 9), None)
    p10 = next((x for x in s["pday"] if x["p_day"] == 10), None)
    text(slide, Inches(0.75), Inches(5.55), Inches(5.4), Inches(0.7), f"9일 상품: {money(p9['s2025'])} → {money(p9['s2026'])} ({pct(p9['pct'])})", size=16, bold=True, color=RED)
    text(slide, Inches(6.7), Inches(5.55), Inches(5.4), Inches(0.7), f"10일 상품: {money(p10['s2025'])} → {money(p10['s2026'])} ({pct(p10['pct'])})", size=16, bold=True, color=RED)
    footer(slide)


def slide_factors(s):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "외부요인은 압박, 내부요인은 하락폭을 키움", "요인 분해", 6)
    text(slide, Inches(0.85), Inches(1.65), Inches(5.3), Inches(0.45), "외부요인", size=22, bold=True, color=BLUE)
    add_bullet_list(slide, Inches(0.95), Inches(2.25), [
        "환율 부담: 원/달러가 2025·2026년 4월 모두 1,400원대 중후반에서 움직여 한국발 수요의 달러 체감가를 압박",
        "항공/교통비 부담: 2026년 3월 미국 교통 CPI에서 운송비·항공요금이 상승 요인으로 작용",
        "시장 자체는 회복 전망: NTTO는 미국 국제방문객이 2026년에 8,500만명까지 증가할 것으로 전망",
    ], BLUE)
    text(slide, Inches(6.8), Inches(1.65), Inches(5.3), Inches(0.45), "내부요인", size=22, bold=True, color=RED)
    p9 = next((x for x in s["pday"] if x["p_day"] == 9), None)
    add_bullet_list(slide, Inches(6.9), Inches(2.25), [
        f"핵심 장기 인바운드 약화: 9일 상품 감소액 {money(p9['delta'])}, 전체 감소의 약 {abs(p9['delta'])/abs(s['delta'])*100:.0f}%",
        f"대형 계정/대표 상품 집중 리스크: SIN2052만 {money(s['products'][0]['delta'])} 감소",
        f"볼륨과 단가가 동시 하락: 예약건수 {pct(s['bk_pct'])}, 건당 매출 {pct((s['avg_booking26']-s['avg_booking25'])/s['avg_booking25']*100)}",
    ], RED)
    footer(slide, "출처: Federal Reserve H.10, BTS/BLS CPI, NTTO Forecast, prod_month.php DB 집계")


def slide_actions(s):
    slide = prs.slides.add_slide(BLANK)
    header(slide, "우선순위는 4월 회복, 9일 상품 재점검, 데이터 인식 정리", "대응 제안", 7)
    text(slide, Inches(0.85), Inches(1.65), Inches(11.8), Inches(0.55), "다음 액션", size=24, bold=True)
    actions = [
        ("1", "4월/5월 파이프라인 점검", "확정 지연·미수금·0매출 예약을 분리해 실제 매출 인식 누락 여부 확인"),
        ("2", "SIN2052 및 9일 상품 회복", "롯데관광 단독·미동부/캐나다 장기 상품의 가격, 출발일, 객실/좌석 조건 재검토"),
        ("3", "환율형 가격표 운영", "원화 체감가가 높은 구간에는 조기예약/분납/포함사항 강조로 가격 저항 완화"),
        ("4", "항공권 믹스 관리", "항공권 매출은 방어됐지만 투어 상품 하락을 상쇄하지 못함. 패키지 부가판매 연결 필요"),
    ]
    y = Inches(2.35)
    for no, title, body in actions:
        rect(slide, Inches(0.9), y, Inches(0.45), Inches(0.45), BLUE, shape=MSO_SHAPE.OVAL)
        text(slide, Inches(0.9), y + Inches(0.05), Inches(0.45), Inches(0.25), no, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(slide, Inches(1.55), y - Inches(0.02), Inches(3.3), Inches(0.28), title, size=15, bold=True)
        text(slide, Inches(4.9), y - Inches(0.02), Inches(7.0), Inches(0.42), body, size=12, color=INK)
        y += Inches(0.82)
    text(slide, Inches(0.9), Inches(6.25), Inches(11.7), Inches(0.55), "공식 외부 출처: Federal Reserve H.10(환율), BTS/BLS CPI(항공·교통비), NTTO Forecast(미국 국제방문 전망)", size=10, color=MUTED)
    footer(slide)


def build_deck(summary):
    for fn in [slide_cover, slide_basis, slide_monthly, slide_delta, slide_mix, slide_factors, slide_actions]:
        fn(summary)
    prs.save(OUT)


def get_font(size, bold=False):
    candidates = [
        r"C:\Windows\Fonts\malgunbd.ttf" if bold else r"C:\Windows\Fonts\malgun.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for p in candidates:
        if Path(p).exists():
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def pil_color(rgb):
    return (rgb[0], rgb[1], rgb[2])


def draw_wrapped(d, xy, text_value, font, fill, width, line_spacing=6):
    x, y = xy
    words = str(text_value).split()
    line = ""
    for word in words:
        test = (line + " " + word).strip()
        if d.textbbox((0, 0), test, font=font)[2] <= width:
            line = test
        else:
            d.text((x, y), line, font=font, fill=fill)
            y += font.size + line_spacing
            line = word
    if line:
        d.text((x, y), line, font=font, fill=fill)
    return y + font.size


def draw_preview(summary, idx, title, bullets=None):
    img = Image.new("RGB", (1600, 900), (245, 247, 250))
    d = ImageDraw.Draw(img)
    title_font = get_font(44, True)
    body_font = get_font(25)
    small_font = get_font(18)
    d.rectangle([70, 60, 85, 805], fill=pil_color(BLUE))
    d.text((115, 85), title, font=title_font, fill=pil_color(INK))
    d.line([115, 160, 1510, 160], fill=pil_color(LINE), width=2)
    if idx == 1:
        d.text((115, 245), "상반기 매출추이 변동현황 분석", font=get_font(60, True), fill=pil_color(INK))
        d.text((1130, 245), pct(summary["delta_pct"]), font=get_font(70, True), fill=pil_color(RED), anchor="ra")
        d.text((930, 335), f"{money(summary['total25'])} → {money(summary['total26'])}", font=body_font, fill=pil_color(INK))
    elif idx in (3, 4):
        vals = summary["sales"]
        maxv = max(max(vals[2025]), max(vals[2026]))
        base_x, base_y, W, H = 170, 710, 980, 430
        d.line([base_x, base_y, base_x + W, base_y], fill=(180, 190, 200), width=2)
        for yv, color, label in [(vals[2025], pil_color(BLUE), "2025"), (vals[2026], pil_color(RED), "2026")]:
            pts = []
            for i, v in enumerate(yv):
                x = base_x + i * W / 3
                y = base_y - (v / maxv) * H
                pts.append((x, y))
            d.line(pts, fill=color, width=5)
            for p in pts:
                d.ellipse([p[0]-7, p[1]-7, p[0]+7, p[1]+7], fill=color)
            d.text((base_x + W + 35, 260 if label == "2025" else 300), label, font=small_font, fill=color)
        for i, m in enumerate(summary["month_names"]):
            d.text((base_x + i * W / 3 - 15, base_y + 18), m, font=small_font, fill=pil_color(MUTED))
    else:
        y = 220
        for b in bullets or []:
            d.ellipse([125, y + 8, 145, y + 28], fill=pil_color(BLUE))
            y = draw_wrapped(d, (165, y), b, body_font, pil_color(INK), 1220) + 25
    out = f"{PREVIEW_PREFIX}{idx:02d}.png"
    img.save(out)
    return out


def render_previews(summary):
    specs = [
        ("매출추이 분석 표지", []),
        ("집계 기준과 핵심 수치", [f"매출 {money(summary['total25'])} → {money(summary['total26'])} ({pct(summary['delta_pct'])})", f"예약건수 {num(summary['bk25'])} → {num(summary['bk26'])}", "prod_month.php 기준: DONE / MAIN / PICKUP·SENDING 제외"]),
        ("월별 매출 추이", []),
        ("월별 감소 기여", []),
        ("상품/일정 믹스", [f"최대 감소 상품: {summary['products'][0]['code']} {money(summary['products'][0]['delta'])}", "9일·10일 장기 상품 하락이 전체 감소를 주도", "항공권 매출은 방어됐지만 패키지 하락을 상쇄하지 못함"]),
        ("외부요인과 내부요인", ["환율·항공요금은 수요 저항 요인", "NTTO 전망은 전체 미국 방문 수요 회복을 가리킴", "내부적으로는 핵심 장기 인바운드와 대형 계정 상품 약화가 결정적"]),
        ("대응 제안", ["4월/5월 파이프라인 및 0매출 예약 점검", "SIN2052와 9일 상품 가격·좌석·출발일 재정비", "환율형 프로모션과 패키지 부가판매 강화"]),
    ]
    return [draw_preview(summary, i + 1, title, bullets) for i, (title, bullets) in enumerate(specs)]


def main():
    data = run_sales_query()
    summary = summarize(data)
    build_deck(summary)
    previews = render_previews(summary)
    print("PPTX:", OUT)
    print("PREVIEWS:", len(previews))
    for p in previews:
        print(p)
    print("SUMMARY:", json.dumps({
        "sales_2025": summary["total25"],
        "sales_2026": summary["total26"],
        "delta_pct": summary["delta_pct"],
        "pax_pct": summary["pax_pct"],
        "bookings_pct": summary["bk_pct"],
    }, ensure_ascii=False))


if __name__ == "__main__":
    main()
